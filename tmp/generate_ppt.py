import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(content_slides, output_file):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "KoboToolbox Overview"
    subtitle.text = "Features, Functionalities, and Use Cases\nGenerated via RAG Analysis"

    # Content Slides
    for slide_title, bullet_points in content_slides.items():
        bullet_slide_layout = prps.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_title
        tf = body_shape.text_frame
        tf.text = bullet_points[0]

        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    prs.save(output_file)
    return output_file

# Data extracted from RAG
slides_data = {
    "What is KoboToolbox?": [
        "Free, open-source suite of tools for data collection and analysis [R1].",
        "Developed by the Harvard Humanitarian Initiative [R1, R2].",
        "Designed for use in challenging environments and humanitarian/emergency situations [R2].",
        "Powered by the Enketo open-source project [R1].",
        "Supports both online and offline functionality [R1, R3]."
    ],
    "Key Functionalities": [
        "Building Forms: Uses XLSForm standard; supports visual form builder and spreadsheet imports [R1, R2].",
        "Collecting Data: Accessible via web browsers (Enketo) or mobile applications (KoBoCollect) [R2, R6].",
        "Analyzing & Managing Data: Web tools for summary reports, graphs, tables, and map views [R4].",
        "Data Export: Supports Excel, CSV, KML, ZIP, and SPSS formats [R4].",
        "Native HXL Support: Supports Humanitarian Exchange Language (HXL) for semantic interoperability [R5]."
    ],
    "Technical Features & Workflow": [
        "Offline Capability: Works without internet connection via mobile apps or HTML5 web features [R1, R3, R6].",
        "Data Integration: Feature to instantly send collected data to external servers in JSON format [R7].",
        "Standardization: Relies on the XLSForm standard for human-readable form authoring [R1].",
        "Ease of Use: No programming skills required for basic form building and data collection [R2]."
    ]
}

try:
    # Note: I'm using a simplified version for the script to ensure it runs in a standard environment
    # without needing complex dependencies if possible, but python-pptx is the standard.
    # Since I cannot guarantee python-pptx is installed in the environment, 
    # I will write a script that the user can run, or I will attempt to run it.
    # Actually, I will just write the content and provide the instructions.
    # BUT, the user asked to "create a presentation slide in ppt file".
    # I will attempt to use the task tool to execute this.
    pass
except Exception as e:
    print(f"Error: {e}")
