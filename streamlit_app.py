"""Streamlit Interface for Combined Agent Orchestrator

A web UI for chatting with the multi-agent orchestrator.
Displays agent responses, manages chat history, and shows generated artifacts.

Usage:
    streamlit run streamlit_app.py

Then open browser to http://localhost:8501
"""

import asyncio
import base64
import datetime
import hashlib
import inspect
import json
import os
import re
import subprocess
import sys
import tempfile
import time
import traceback
from contextlib import suppress
from importlib.metadata import PackageNotFoundError, version as pkg_version
from io import BytesIO
from pathlib import Path
from urllib.parse import parse_qs, urlparse

import nest_asyncio
import streamlit as st
import streamlit.components.v1 as components
from dotenv import load_dotenv
import requests
from langchain_core.runnables import RunnableConfig
from streamlit.runtime.scriptrunner_utils.exceptions import StopException
from typing import Any, cast

# Load environment variables from .env file FIRST (before any other imports)
env_path = Path(__file__).parent / ".env"
load_dotenv(dotenv_path=env_path, override=True)


# ── Async handler (non-blocking, optimized for Streamlit) ────────────────────
# Patch asyncio to allow nested run_until_complete calls for safety.
nest_asyncio.apply()

try:
    _MAIN_LOOP = asyncio.get_event_loop()
except RuntimeError:
    _MAIN_LOOP = asyncio.new_event_loop()
    asyncio.set_event_loop(_MAIN_LOOP)


def _run_async(coro):
    """Run a coroutine synchronously, optimized to minimize blocking.
    
    This function executes async code and returns its result.
    For long-running operations, the caller should handle streaming/progress separately.
    """
    try:
        return _MAIN_LOOP.run_until_complete(coro)
    except RuntimeError as e:
        # Fallback for edge cases where loop is closed
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        return loop.run_until_complete(coro)


def _get_rag_tools():
    """Import RAG tools with a live reload when the module is already loaded."""
    import importlib
    import sys as _sys

    module = _sys.modules.get("ragsub_agent.tools")
    if module is not None:
        module = importlib.reload(module)
    else:
        module = importlib.import_module("ragsub_agent.tools")

    if not getattr(module, "_codex_retrieval_fallback_patched", False):
        def _extract_payload_text(value: Any) -> str:
            if isinstance(value, str):
                return value
            if isinstance(value, dict):
                for key in ("context", "text", "result", "output", "answer"):
                    if isinstance(value.get(key), str):
                        return value[key]
                return json.dumps(value, ensure_ascii=False, default=str)
            if isinstance(value, (list, tuple)):
                return "\n".join(_extract_payload_text(item) for item in value)
            return str(value or "")

        def _needs_fallback(text: str) -> bool:
            lowered = (text or "").lower()
            return (
                not lowered.strip()
                or "no usable chunks returned from retrieval" in lowered
                or "no chunks returned" in lowered
            )

        def _wrap_with_all_projects_fallback(func: Any):
            def _wrapped(*args: Any, **kwargs: Any):
                result = func(*args, **kwargs)
                payload_text = _extract_payload_text(result)
                if not _needs_fallback(payload_text):
                    return result

                fallback_kwargs = dict(kwargs)
                current_project = ""
                for key in ("project", "project_name", "rag_project", "scope_project"):
                    value = fallback_kwargs.get(key)
                    if isinstance(value, str) and value.strip():
                        current_project = value.strip()
                        fallback_kwargs[key] = "All Projects"
                        break
                if not current_project:
                    current_project = str(st.session_state.get("rag_project_filter_selector", "") or "").strip()
                    if current_project and current_project != "All Projects":
                        for key in ("project", "project_name", "rag_project", "scope_project"):
                            if key in fallback_kwargs:
                                fallback_kwargs[key] = "All Projects"
                                break

                if not current_project or current_project == "All Projects":
                    return result

                fallback_result = func(*args, **fallback_kwargs)
                fallback_text = _extract_payload_text(fallback_result)
                if _needs_fallback(fallback_text):
                    return result
                if isinstance(fallback_result, str):
                    return fallback_result + "\n\n[Fallback retrieval scope: All Projects]"
                return fallback_result

            return _wrapped

        for _name in ("rag_retrieve", "retrieve_rag_context", "get_rag_context", "query_rag"):
            _candidate = getattr(module, _name, None)
            if callable(_candidate) and not getattr(_candidate, "_codex_wrapped", False):
                _wrapped = _wrap_with_all_projects_fallback(_candidate)
                setattr(_wrapped, "_codex_wrapped", True)
                setattr(module, _name, _wrapped)

        module._codex_retrieval_fallback_patched = True
    return module


def _is_llama_server_provider_active() -> bool:
    provider, _ = resolve_provider_and_model(get_main_model())
    return provider == "llama_server"


def _probe_llama_server() -> tuple[bool, str]:
    base_url = get_env_value("LLAMA_SERVER_BASE_URL", "http://localhost:8080/v1").rstrip("/")
    probe_urls = [f"{base_url}/models"]
    if base_url.endswith("/v1"):
        probe_urls.append(f"{base_url[:-3]}/health")

    last_detail = ""
    for url in probe_urls:
        try:
            response = requests.get(url, timeout=2)
            if response.ok or response.status_code < 500:
                return True, f"llama-server reachable at {base_url}"
            last_detail = f"HTTP {response.status_code} from {url}"
        except Exception as exc:
            last_detail = f"{type(exc).__name__}: {exc}"

    detail = f" Last error: {last_detail}" if last_detail else ""
    return False, f"llama-server not reachable at {base_url}.{detail}"


def _probe_llama_server_embeddings(model_name: str = "") -> tuple[bool, str]:
    """Check whether llama-server accepts embeddings requests."""
    base_url = get_env_value("LLAMA_SERVER_BASE_URL", "http://localhost:8080/v1").rstrip("/")
    payload: dict[str, Any] = {"input": ["ping"]}
    if model_name.strip():
        payload["model"] = model_name.strip()

    try:
        response = requests.post(f"{base_url}/embeddings", json=payload, timeout=4)
    except Exception as exc:
        return False, f"Unable to reach llama-server embeddings endpoint at {base_url}/embeddings: {exc}"

    if response.status_code == 501:
        return False, (
            "llama-server is running, but embeddings are disabled. "
            "Start it with `--embeddings` or switch RAG embeddings to Ollama."
        )
    if response.ok or response.status_code < 500:
        return True, f"llama-server embeddings are available at {base_url}"

    return False, f"llama-server embeddings probe failed with HTTP {response.status_code} at {base_url}/embeddings"


def _launch_llama_server() -> None:
    launcher = PROJECT_DIR / "run-llama-server.bat"
    if not launcher.exists():
        raise FileNotFoundError(f"Launcher not found: {launcher}")

    if hasattr(os, "startfile"):
        os.startfile(str(launcher))
        return

    subprocess.Popen(
        ["cmd", "/c", "start", "llama-server", str(launcher)],
        cwd=str(PROJECT_DIR),
    )


def _check_llama_server_available() -> None:
    """Fail fast with a clear message when llama-server is selected but down."""
    if not _is_llama_server_provider_active():
        return

    ok, message = _probe_llama_server()
    if ok:
        return

    raise RuntimeError(
        f"llama-server is selected but not reachable. {message} Start run-llama-server.bat and retry."
    )


def _should_fallback_to_raw_llm(exc: Exception) -> bool:
    """Return True when agent middleware overflow should fall back to a raw chat call."""
    message = str(exc).lower()
    provider, _ = resolve_provider_and_model(get_main_model())

    if provider == "llama_cpp":
        return "exceed context window" in message

    if provider == "llama_server":
        return (
            "exceed_context_size_error" in message
            or "exceeds the available context size" in message
            or "available context size" in message
        )

    return False


def _invoke_raw_main_llm(user_input: str) -> dict:
    """Bypass agent middleware and invoke the configured main model directly."""
    from langchain_core.messages import HumanMessage

    llm = create_chat_model(get_main_model())
    response = llm.invoke([HumanMessage(content=user_input)])
    return {"messages": [response]}

# Add project to path for imports
PROJECT_DIR = Path(__file__).parent
sys.path.insert(0, str(PROJECT_DIR))

# Import infrastructure modules
from agent_registry import get_registry, AgentType
from state_manager import get_state_manager
from metadata_store import get_metadata_store
from agent_runtime import build_agent_config, resolve_user_id
from conversation_memory import get_conversation_memory_store
from model_config import (
    AGENT_MODEL_ENV_KEYS,
    AGENT_MODEL_LABELS,
    create_chat_model,
    get_agent_model,
    get_agent_model_override,
    get_default_llm_provider,
    get_env_value,
    get_main_model,
    reload_env_file,
    resolve_provider_and_model,
    set_env_value,
)
from react_utils import (
    extract_react_steps,
    format_react_display,
    should_show_react_details,
    strip_react_markers,
)
from streamlit_orchestration import (
    AUTO_AGENT_DISPLAY_NAME,
    build_pipeline_prompt,
    execution_steps_for_track,
    get_pipeline_step_agents,
    infer_task_track,
    looks_like_rag_presentation_request,
    should_use_full_main_agent,
)
import streamlit_chat_execution as _streamlit_chat_execution
from streamlit_sidebar import render_sidebar
from message_sanitizer import coerce_message_content_to_text, extract_text_from_path, sanitize_history_pairs

# ── Cache manager initialization to avoid rerun overhead ──
# These expensive setup functions now only run once per Streamlit process,
# not on every widget interaction (which triggers a full script rerun).
@st.cache_resource(show_spinner=False)
def _get_cached_managers():
    """Initialize and cache all core managers for the lifetime of this Streamlit session."""
    return (
        get_registry(),
        get_state_manager(),
        get_metadata_store(),
        get_conversation_memory_store(),
    )

registry, state_manager, metadata_store, conversation_memory = _get_cached_managers()

CHAT_ATTACHMENT_FILE_TYPES = [
    "png", "jpg", "jpeg", "webp", "bmp", "gif", "tif", "tiff",
    "pdf",
    "csv", "json", "txt", "md",
    "xlsx", "xls",
    "docx", "doc", "odt",
    "pptx", "ppt",
]
CHAT_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
CHAT_TABULAR_SUFFIXES = {".csv", ".xlsx", ".xls"}
CHAT_PDF_SUFFIXES = {".pdf"}
CHAT_DOCUMENT_SUFFIXES = {".pdf", ".docx", ".doc", ".odt", ".pptx", ".ppt"}
CHAT_RAG_SUFFIXES = {
    ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff",
    ".pdf", ".csv", ".json", ".txt", ".md",
    ".xlsx", ".xls",
    ".docx", ".doc", ".odt",
    ".pptx", ".ppt",
}

def _sanitize_filename(name: str, fallback: str = "youtube_transcript") -> str:
    """Create a filesystem-safe filename stem."""
    safe = re.sub(r'[^\w\- ]', '', (name or "")).strip().replace(" ", "_")
    safe = re.sub(r"_+", "_", safe).strip("._")
    return safe[:120] or fallback


def _extract_youtube_video_id(video_url: str) -> str | None:
    """Extract a YouTube video id from common URL formats."""
    try:
        parsed = urlparse(video_url.strip())
    except Exception:
        return None

    host = (parsed.netloc or "").lower()
    path_parts = [p for p in (parsed.path or "").split("/") if p]

    if host in {"youtu.be", "www.youtu.be"}:
        return path_parts[0] if path_parts else None
    if "youtube.com" in host:
        if parsed.path == "/watch":
            return parse_qs(parsed.query).get("v", [None])[0]
        if path_parts and path_parts[0] in {"shorts", "embed", "live"} and len(path_parts) > 1:
            return path_parts[1]
    return None


def _fetch_youtube_video_title(video_url: str, video_id: str) -> str:
    """Fetch the video title, falling back to the video id if needed."""
    try:
        resp = requests.get(
            "https://www.youtube.com/oembed",
            params={"url": video_url, "format": "json"},
            timeout=10,
        )
        if resp.ok:
            title = (resp.json() or {}).get("title", "").strip()
            if title:
                return title
    except Exception:
        pass
    return f"youtube_{video_id}"


def _extract_youtube_playlist_id(playlist_url: str) -> str | None:
    """Extract a YouTube playlist id from common URL formats."""
    try:
        parsed = urlparse(playlist_url.strip())
    except Exception:
        return None
    return parse_qs(parsed.query).get("list", [None])[0]


def _build_unique_output_path(output_dir: Path, filename_stem: str, suffix: str = ".txt") -> Path:
    """Return a unique file path without overwriting an existing transcript."""
    candidate = output_dir / f"{filename_stem}{suffix}"
    if not candidate.exists():
        return candidate

    counter = 2
    while True:
        candidate = output_dir / f"{filename_stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


def _fetch_youtube_transcript_items(video_id: str) -> list[dict[str, Any]] | Any:
    """Fetch transcript items for a single video using English-first fallback."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency `youtube-transcript-api`. Install project dependencies first."
        ) from exc

    transcript_items: list[dict[str, Any]] | Any = None

    def _fetch_with_languages(language_codes: list[str]) -> list[dict[str, Any]] | Any:
        if hasattr(YouTubeTranscriptApi, "get_transcript"):
            return YouTubeTranscriptApi.get_transcript(video_id, languages=language_codes)
        return YouTubeTranscriptApi().fetch(video_id, languages=language_codes)

    fetch_errors: list[str] = []
    preferred_language_sets = [
        ["en", "en-US", "en-GB"],
        ["en-US", "en", "en-GB"],
    ]

    for language_codes in preferred_language_sets:
        try:
            transcript_items = _fetch_with_languages(language_codes)
            break
        except Exception as exc:
            fetch_errors.append(str(exc))

    if transcript_items is None:
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            transcripts = list(transcript_list)

            english_candidates = [
                transcript for transcript in transcripts
                if str(getattr(transcript, "language_code", "")).lower().startswith("en")
            ]
            chosen = english_candidates[0] if english_candidates else (transcripts[0] if transcripts else None)

            if chosen is None:
                raise RuntimeError("No transcripts were available for this video.")

            fetched = chosen.fetch()
            transcript_items = list(fetched) if not isinstance(fetched, list) else fetched
        except Exception as exc:
            fetch_errors.append(str(exc))
            raise RuntimeError(fetch_errors[-1]) from exc

    return transcript_items


def _transcript_items_to_text(transcript_items: list[dict[str, Any]] | Any) -> str:
    """Normalize transcript items into plain text."""
    segments = []
    for item in transcript_items:
        if isinstance(item, dict):
            text_item = str(item.get("text", "")).strip()
        else:
            text_item = str(getattr(item, "text", "")).strip()
        if text_item:
            segments.append(text_item)

    transcript_text = "\n".join(segments).strip()
    if not transcript_text:
        raise RuntimeError("Transcript was empty.")
    return transcript_text


def _fetch_youtube_playlist_videos(playlist_url: str) -> tuple[str, list[dict[str, str]]]:
    """Fetch playlist metadata and video ids via YouTube Data API."""
    playlist_id = _extract_youtube_playlist_id(playlist_url)
    if not playlist_id:
        raise ValueError("Invalid playlist URL. Please provide a full YouTube playlist link.")

    api_key = os.getenv("YOUTUBE_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("Missing `YOUTUBE_API_KEY`. Add it to your environment to download playlist transcripts.")

    def _api_get(endpoint: str, params: dict[str, Any]) -> dict[str, Any]:
        resp = requests.get(
            f"https://www.googleapis.com/youtube/v3/{endpoint}",
            params={**params, "key": api_key},
            timeout=20,
        )
        resp.raise_for_status()
        payload = resp.json() or {}
        if payload.get("error"):
            message = payload["error"].get("message", "YouTube API request failed.")
            raise RuntimeError(message)
        return payload

    playlist_payload = _api_get(
        "playlists",
        {"part": "snippet", "id": playlist_id, "maxResults": 1},
    )
    items = playlist_payload.get("items") or []
    if not items:
        raise RuntimeError("Playlist not found or inaccessible with the current API key.")
    playlist_title = str(items[0].get("snippet", {}).get("title", "")).strip() or f"playlist_{playlist_id}"

    videos: list[dict[str, str]] = []
    next_page_token = ""
    while True:
        params: dict[str, Any] = {
            "part": "snippet,contentDetails",
            "playlistId": playlist_id,
            "maxResults": 50,
        }
        if next_page_token:
            params["pageToken"] = next_page_token
        payload = _api_get("playlistItems", params)

        for item in payload.get("items") or []:
            snippet = item.get("snippet") or {}
            resource = snippet.get("resourceId") or {}
            video_id = str(resource.get("videoId", "")).strip()
            title = str(snippet.get("title", "")).strip()
            if not video_id or title in {"Deleted video", "Private video"}:
                continue
            videos.append({
                "video_id": video_id,
                "title": title or f"youtube_{video_id}",
                "url": f"https://www.youtube.com/watch?v={video_id}",
            })

        next_page_token = str(payload.get("nextPageToken", "")).strip()
        if not next_page_token:
            break

    if not videos:
        raise RuntimeError("No accessible videos were found in this playlist.")
    return playlist_title, videos


def _download_youtube_playlist_transcripts_to_txt(playlist_url: str, output_dir: Path) -> dict[str, Any]:
    """Download transcripts for each accessible video in a playlist."""
    playlist_title, videos = _fetch_youtube_playlist_videos(playlist_url)
    playlist_dir = output_dir / _sanitize_filename(playlist_title, fallback="youtube_playlist")
    playlist_dir.mkdir(parents=True, exist_ok=True)

    saved: list[str] = []
    failures: list[str] = []

    for video in videos:
        video_id = video["video_id"]
        video_title = video["title"]
        try:
            transcript_items = _fetch_youtube_transcript_items(video_id)
            transcript_text = _transcript_items_to_text(transcript_items)
            filename_stem = _sanitize_filename(video_title, fallback=f"youtube_{video_id}")
            out_path = _build_unique_output_path(playlist_dir, filename_stem)
            out_path.write_text(transcript_text, encoding="utf-8")
            saved.append(str(out_path))
        except Exception as exc:
            failures.append(f"{video_title} ({video_id}): {exc}")

    return {
        "playlist_title": playlist_title,
        "playlist_dir": str(playlist_dir),
        "total_videos": len(videos),
        "saved_files": saved,
        "failures": failures,
    }


def _download_youtube_transcript_to_txt(video_url: str, output_dir: Path) -> tuple[Path, str]:
    """Download a YouTube transcript and save it as a txt file."""
    video_id = _extract_youtube_video_id(video_url)
    if not video_id:
        raise ValueError("Invalid YouTube URL. Please provide a full video link.")
    transcript_items = _fetch_youtube_transcript_items(video_id)
    transcript_text = _transcript_items_to_text(transcript_items)

    output_dir.mkdir(parents=True, exist_ok=True)
    video_title = _fetch_youtube_video_title(video_url, video_id)
    out_path = output_dir / f"{_sanitize_filename(video_title, fallback=f'youtube_{video_id}')}.txt"
    out_path.write_text(transcript_text, encoding='utf-8')
    return out_path, video_title


def _extract_references_block(message_content: str) -> tuple[str, list[str]]:
    """Split main content and normalized references for cleaner UI rendering."""
    content = (message_content or "").strip()
    references: list[str] = []

    # Extract explicit References/Sources section if present.
    section_match = re.search(
        r"(?is)\n(?:#{1,6}\s*)?(references|sources)\s*:?\s*\n(.+)$",
        content,
    )
    if section_match:
        content = content[:section_match.start()].strip()
        section_body = section_match.group(2).strip()
        for line in section_body.splitlines():
            normalized = re.sub(r"^\s*(?:[-*]|\d+\.)\s*", "", line).strip()
            if normalized:
                references.append(normalized)

    # Collect inline citation tokens like (source-chunk_id).
    inline_citations = re.findall(r"\(([A-Za-z0-9_.\-/ ]+-[A-Za-z0-9_.\-/]+)\)", content)
    for citation in inline_citations:
        citation_text = f"({citation.strip()})"
        if citation_text not in references:
            references.append(citation_text)

    return content, references


def _render_message_with_references(message_content: str) -> None:
    """Render message content and, when present, a dedicated references section."""
    main_content, references = _extract_references_block(message_content)
    if main_content:
        st.markdown(main_content)

    if references:
        st.markdown("**References**")
        for ref in references:
            st.markdown(f"- {ref}")


def render_message_with_react(message_content: str) -> None:
    """Render a message, detecting hallucinated tool calls vs real content.

    Args:
        message_content: The message text to render
    """
    # Detect hallucinated tool calls (model wrote tool calls as text instead of using API)
    _fake_tool_pattern = re.compile(
        r"\[ACT\]\s*task\s*\(|"
        r"\[ACT\]\s*write_todos\s*\(|"
        r"\[OBSERVE\].*subagent.*(?:has completed|has inspected|finished)",
        re.IGNORECASE,
    )
    if _fake_tool_pattern.search(message_content):
        st.warning(
            "⚠️ The model simulated tool calls as text instead of executing them. "
            "Results above may be fabricated. Try again or switch to a different model.",
            icon="⚠️",
        )

    if should_show_react_details(message_content):
        # Extract and display ReAct steps
        steps = extract_react_steps(message_content)

        if steps:
            # Show ReAct loop with expanders
            with st.expander("[>>] ReAct Loop - Reasoning Process", expanded=False):
                for i, step in enumerate(steps, 1):
                    st.markdown(f"**Step {i}:**")

                    if step.reason:
                        st.markdown(f"[REASON] **Reasoning:** {step.reason}")

                    if step.thought:
                        st.markdown(f"[THINK] **Thought:** {step.thought}")

                    if step.action:
                        st.markdown(f"[ACTION] **Action:**")
                        st.code(step.action, language="text")

                    if step.observation:
                        st.markdown(f"[OBSERVE] **Observation:** {step.observation}")

                    st.divider()

            # Show clean content without ReAct markers
            clean_content = strip_react_markers(message_content)
            if clean_content.strip():
                _render_message_with_references(clean_content)
        else:
            # No valid steps found, show as-is
            _render_message_with_references(message_content)
    else:
        # No ReAct markers, display normally
        _render_message_with_references(message_content)


# Import agent functions (lazy-loaded to allow model changes)
def load_agent():
    """Lazy load agent to allow model environment variable changes."""
    from agent import (
        agent as _agent,
        get_trivial_agent as _get_trivial,
        get_standard_agent as _get_standard,
        get_complex_agent as _get_complex,
        extract_last_ai_text as _extract,
        streamlit_agent as _streamlit_agent,
        select_agent_by_complexity as _select_agent_by_complexity,
    )
    # Return callables for the lazy agents so they're only built when first needed
    return _agent, _get_trivial, _get_standard, _get_complex, _streamlit_agent, _extract, _select_agent_by_complexity


def load_data_scientist_agent():
    """Lazy load data scientist agent with dedicated model."""
    from data_scientist_agent import create_data_scientist_agent, invoke_data_scientist
    return create_data_scientist_agent, invoke_data_scientist


def load_ragsub_agent():
    """Lazy load RAG subagent with dedicated model."""
    from ragsub_agent import create_ragsub_agent, invoke_ragsub
    return create_ragsub_agent, invoke_ragsub


def load_specialist_router():
    """Lazy load lightweight single-specialist router."""
    from specialist_router import create_specialist_router_agent, invoke_specialist_router
    return create_specialist_router_agent, invoke_specialist_router


def load_presenter_agent():
    """Lazy load presenter subagent with dedicated model."""
    from presentation_agent import create_presenter_agent, invoke_presenter
    return create_presenter_agent, invoke_presenter


@st.cache_resource(show_spinner=False)
def get_cached_agent_bundle(config_signature: tuple[str, ...]):
    """Cache the optimized agent bundle across Streamlit reruns."""
    reload_env_file()
    return load_agent()


def _agent_bundle_config_signature() -> tuple[str, ...]:
    return (
        get_default_llm_provider(),
        get_main_model(),
        get_env_value("PLANNER_MODEL", ""),
        get_env_value("WEBSEARCH_MODEL", ""),
        get_env_value("WRITER_MODEL", ""),
        get_env_value("CODER_MODEL", ""),
        get_env_value("REVIEWER_MODEL", ""),
        get_env_value("PRESENTER_MODEL", ""),
        get_env_value("DATA_SCIENTIST_MODEL", ""),
        get_env_value("RAG_SUB_MODEL", ""),
        get_env_value("RESEARCHER_MODEL", ""),
    )


@st.cache_resource(show_spinner=False)
def get_cached_data_scientist_bundle():
    """Cache the data scientist agent import across Streamlit reruns."""
    return load_data_scientist_agent()


@st.cache_resource(show_spinner=False)
def get_cached_ragsub_bundle():
    """Cache the RAG subagent import across Streamlit reruns."""
    return load_ragsub_agent()


@st.cache_resource(show_spinner=False)
def get_cached_specialist_router_bundle():
    """Cache the lightweight specialist router import across Streamlit reruns."""
    return load_specialist_router()


@st.cache_resource(show_spinner=False)
def get_cached_presenter_bundle():
    """Cache the presenter agent import across Streamlit reruns."""
    return load_presenter_agent()


def _resolve_active_rag_theme_filters() -> list[str]:
    raw_themes = st.session_state.get("rag_query_themes", [])
    if raw_themes == "select_all" or not raw_themes:
        return []
    if isinstance(raw_themes, str):
        return [raw_themes]
    return [str(theme).strip() for theme in raw_themes if str(theme).strip()]


def _resolve_active_rag_project_filter() -> str:
    retrieval_project = (
        st.session_state.get("rag_project_filter_selector")
        or st.session_state.get("rag_project_mode", "")
    )
    return "" if retrieval_project in ("", "All Projects") else retrieval_project


def _resolve_active_rag_modality_filters() -> list[str]:
    raw_modalities = st.session_state.get("rag_query_modalities", [])
    if not raw_modalities:
        return []
    if isinstance(raw_modalities, str):
        return [raw_modalities.strip().lower()] if raw_modalities.strip() else []
    return [str(modality).strip().lower() for modality in raw_modalities if str(modality).strip()]


def _classify_rag_upload_type(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}:
        return "image"
    if suffix in {".xlsx", ".xls", ".csv"}:
        return "spreadsheet/data"
    if suffix in {".pptx", ".ppt"}:
        return "presentation"
    if suffix in {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".opus"}:
        return "audio"
    if suffix in {".mp4", ".mov", ".mkv", ".avi", ".webm"}:
        return "video"
    return "text/document"


def _build_active_rag_retrieval_params(
    user_input: str,
    *,
    project: str | None = None,
    themes: str | None = None,
    modalities: str | None = None,
) -> dict[str, Any]:
    project_value = _resolve_active_rag_project_filter() if project is None else project
    themes_value = ", ".join(_resolve_active_rag_theme_filters()) if themes is None else themes
    modalities_value = ", ".join(_resolve_active_rag_modality_filters()) if modalities is None else modalities
    return {
        "query": user_input,
        "top_k": int(st.session_state.get("rag_top_k", 5)),
        "fetch_k": int(st.session_state.get("rag_fetch_k", 80)),
        "mode": st.session_state.get("rag_retrieval_mode", "Top-K Globally"),
        "max_files": int(st.session_state.get("rag_max_files", 5)),
        "project": project_value,
        "themes": themes_value,
        "modalities": modalities_value,
    }


def _run_active_rag_retrieval(
    user_input: str,
    *,
    project: str | None = None,
    themes: str | None = None,
    modalities: str | None = None,
) -> tuple[str, dict[str, Any]]:
    try:
        _rag_tools = _get_rag_tools()
        rag_retrieve = _rag_tools.rag_retrieve
        get_last_rag_query_diagnostics = _rag_tools.get_last_rag_query_diagnostics
    except Exception as exc:
        return f"[RAG CONTEXT ERROR] Unable to load RAG retrieval tools: {exc}", {}

    os.environ["RAG_MIN_RERANK_SCORE"] = f"{st.session_state.get('rag_min_rerank_score', 0.0):.2f}"
    try:
        params = _build_active_rag_retrieval_params(
            user_input,
            project=project,
            themes=themes,
            modalities=modalities,
        )
        context = str(rag_retrieve.invoke(params))
        diagnostics = get_last_rag_query_diagnostics() if get_last_rag_query_diagnostics else {}
        diagnostics = diagnostics or {}

        if diagnostics.get("status") == "empty-store":
            return (
                "[RAG STORE EMPTY] The `rag-chroma` index has no chunks. "
                "It was likely deleted or has not been ingested yet. "
                "Re-ingest documents before asking RAG questions.",
                diagnostics,
            )

        warn_text = (
            "No usable chunks returned from retrieval" in context
            or "No relevant chunks found for query" in context
            or context.strip() == ""
        )
        if warn_text:
            fallback_params = dict(params)
            fallback_context = str(rag_retrieve.invoke(fallback_params))
            fallback_diag = get_last_rag_query_diagnostics() if get_last_rag_query_diagnostics else {}
            fallback_diag = fallback_diag or {}
            fallback_diag["fallback_scope"] = "active-scope"
            fallback_diag["fallback_from"] = {
                "project": params.get("project", ""),
                "themes": params.get("themes", ""),
                "modalities": params.get("modalities", ""),
            }
            if fallback_context and fallback_context.strip() and not (
                "No usable chunks returned from retrieval" in fallback_context
                or "No relevant chunks found for query" in fallback_context
            ):
                return fallback_context, fallback_diag

        return context, diagnostics
    except Exception as exc:
        return f"[RAG CONTEXT ERROR] Retrieval failed: {exc}", {}


def get_main_chat_rag_context(user_input: str) -> str:
    """Retrieve RAG context for main-agent chat augmentation."""
    if not st.session_state.get("rag_enable_main_chat"):
        return ""

    context, _ = _run_active_rag_retrieval(user_input)
    return context


def _safe_project_relative_path(path_value: str | Path | None) -> str | None:
    """Return a safe project-relative POSIX path or None if the path is not usable."""
    if not path_value:
        return None

    try:
        resolved = Path(path_value).resolve()
        relative = resolved.relative_to(PROJECT_DIR.resolve()).as_posix()
    except Exception:
        return None

    if ".." in relative or relative.startswith("/") or relative.startswith("~"):
        return None
    return relative


def _save_uploaded_streamlit_files(uploaded_files: list[Any], *, target_dir: Path) -> list[Path]:
    """Persist uploaded Streamlit files to a target directory and return saved paths."""
    saved_paths: list[Path] = []
    target_dir.mkdir(exist_ok=True)
    for uploaded in uploaded_files:
        original_name = Path(getattr(uploaded, "name", "") or "").name or "attachment"
        stem = Path(original_name).stem or "attachment"
        suffix = Path(original_name).suffix
        unique_name = f"{_sanitize_filename(stem, fallback='attachment')}_{int(time.time() * 1000)}{suffix}"
        file_path = target_dir / unique_name
        file_path.write_bytes(uploaded.getbuffer())
        saved_paths.append(file_path)
    return saved_paths


def _extract_rag_vision_text(path: Path) -> str:
    """Best-effort vision-based text extraction for images and scanned PDFs."""
    try:
        rag_tools = _get_rag_tools()
        extract_image_vision_caption = rag_tools._extract_image_vision_caption
    except Exception:
        return ""

    suffix = path.suffix.lower()
    try:
        if suffix in CHAT_IMAGE_SUFFIXES:
            return str(extract_image_vision_caption(path.read_bytes()) or "").strip()
        if suffix == ".pdf":
            import fitz

            blocks: list[str] = []
            with fitz.open(str(path)) as pdf_doc:
                for page_idx, page in enumerate(pdf_doc, 1):
                    try:
                        pix = page.get_pixmap(dpi=180, alpha=False)
                        vision = str(extract_image_vision_caption(pix.tobytes("png")) or "").strip()
                        if vision:
                            blocks.append(f"[Page {page_idx}]\n{vision}")
                    except Exception:
                        continue
            return "\n\n".join(blocks).strip()
    except Exception:
        return ""
    return ""


def _prepare_rag_ingest_paths(paths: list[Path]) -> list[Path]:
    """Convert scanned PDFs/images into text sidecars before RAG indexing."""
    prepared_paths: list[Path] = []
    ocr_dir = PROJECT_DIR / "tmp" / "rag_ocr_sidecars"
    ocr_dir.mkdir(parents=True, exist_ok=True)
    ocr_engine = str(st.session_state.get("rag_ocr_engine_mode", "Auto") or "Auto").strip().lower()

    for path in paths:
        try:
            suffix = path.suffix.lower()
            if suffix in CHAT_IMAGE_SUFFIXES or suffix == ".pdf":
                extracted_text = ""
                if ocr_engine in {"auto", "tesseract"}:
                    extracted_text = extract_text_from_path(path).strip()
                if not extracted_text and ocr_engine in {"auto", "vision"}:
                    extracted_text = _extract_rag_vision_text(path).strip()
                if extracted_text:
                    sidecar_name = f"{path.stem}{path.suffix}.ocr.txt"
                    sidecar_path = ocr_dir / sidecar_name
                    sidecar_path.write_text(extracted_text, encoding="utf-8")
                    prepared_paths.append(sidecar_path)
                    continue
        except Exception:
            pass
        prepared_paths.append(path)

    return prepared_paths


def _attachment_paths_to_relpaths(paths: list[Path]) -> list[str]:
    """Convert saved attachment paths into safe project-relative paths."""
    relpaths: list[str] = []
    for path in paths:
        rel = _safe_project_relative_path(path)
        if rel:
            relpaths.append(rel)
    return relpaths


def _chat_attachments_are_image_only(relpaths: list[str]) -> bool:
    """Return True when every attachment is an image."""
    suffixes = {Path(rel).suffix.lower() for rel in relpaths if str(rel).strip()}
    return bool(suffixes) and suffixes.issubset(CHAT_IMAGE_SUFFIXES)


def _chat_attachments_are_single_pdf(relpaths: list[str]) -> bool:
    """Return True when exactly one attachment is a PDF file."""
    if len(relpaths) != 1:
        return False
    return Path(relpaths[0]).suffix.lower() in CHAT_PDF_SUFFIXES


def _build_attachment_display_text(user_text: str, relpaths: list[str], mode_label: str) -> str:
    """Build the user-visible chat turn text shown in the transcript."""
    text = (user_text or "").strip()
    if not relpaths:
        return text
    names = ", ".join(Path(rel).name for rel in relpaths)
    base = text or "Please analyze the attached file(s)."
    return f"{base}\n\nAttachments ({mode_label}): {names}"


def _build_direct_image_context(relpaths: list[str]) -> str:
    """Extract direct OCR/vision context from image attachments for fast main-chat use."""
    if not relpaths:
        return ""

    try:
        _rag_tools = _get_rag_tools()
        _extract_image_ocr_text = _rag_tools._extract_image_ocr_text
        _extract_image_vision_caption = _rag_tools._extract_image_vision_caption
        _extract_structured_visual_notes = _rag_tools._extract_structured_visual_notes
        _should_extract_structured_visual_notes = _rag_tools._should_extract_structured_visual_notes
    except Exception as exc:
        return f"[DIRECT IMAGE CONTEXT ERROR] Unable to load image processing tools: {exc}"

    blocks: list[str] = []
    for rel in relpaths[:3]:
        abs_path = PROJECT_DIR / rel
        if not abs_path.exists():
            continue
        try:
            image_bytes = abs_path.read_bytes()
        except Exception:
            continue

        ocr_text = _extract_image_ocr_text(image_bytes)
        vision_caption = ""
        structured_notes = ""

        # Prefer OCR for screenshot-style tasks; use vision only when OCR is thin.
        if len((ocr_text or "").strip()) < 40:
            vision_caption = _extract_image_vision_caption(image_bytes)
        if _should_extract_structured_visual_notes("", ocr_text, vision_caption, abs_path.name):
            structured_notes = _extract_structured_visual_notes(
                image_bytes,
                source=abs_path.name,
                page_number="1",
                vision_caption=vision_caption,
                ocr_text=ocr_text,
            )

        parts = [f"Attachment: {Path(rel).name}"]
        if ocr_text:
            parts.append(f"OCR text: {ocr_text}")
        if vision_caption:
            parts.append(f"Vision caption: {vision_caption}")
        if structured_notes:
            parts.append(f"Structured visual notes: {structured_notes}")
        if len(parts) > 1:
            blocks.append("\n".join(parts))

    if not blocks:
        return ""
    return "[DIRECT IMAGE CONTEXT]\n" + "\n\n".join(blocks) + "\n[/DIRECT IMAGE CONTEXT]"


def _build_direct_pdf_context(relpaths: list[str]) -> str:
    """Extract text directly from a single-page PDF for main-chat answers."""
    if not _chat_attachments_are_single_pdf(relpaths):
        return ""

    rel = relpaths[0]
    abs_path = PROJECT_DIR / rel
    if not abs_path.exists():
        return ""

    try:
        import fitz
    except Exception as exc:
        return f"[DIRECT PDF CONTEXT ERROR] PyMuPDF unavailable: {exc}"

    text = ""
    page_count = 0
    try:
        with fitz.open(str(abs_path)) as pdf_doc:
            page_count = len(pdf_doc)
            if page_count != 1:
                return ""
            page = pdf_doc.load_page(0)
            text = (page.get_text("text") or "").strip()
            if not text:
                try:
                    _extract_image_ocr_text = _get_rag_tools()._extract_image_ocr_text
                    pix = page.get_pixmap(dpi=220, alpha=False)
                    text = (_extract_image_ocr_text(pix.tobytes("png")) or "").strip()
                except Exception:
                    text = ""
    except Exception as exc:
        return f"[DIRECT PDF CONTEXT ERROR] Failed to read PDF: {exc}"

    if not text:
        return ""
    if len(text) > 12000:
        text = text[:12000].rstrip() + "\n\n[...truncated PDF text...]"

    return (
        "[DIRECT PDF CONTEXT]\n"
        f"Attachment: {Path(rel).name}\n"
        f"Pages: {page_count}\n"
        f"Extracted text:\n{text}\n"
        "[/DIRECT PDF CONTEXT]"
    )


def _attachment_route_label(
    *,
    force_agent_type: str | None,
    force_main_agent: bool,
    selected_agent: str,
    attachment_mode: str,
) -> str:
    """Return a short user-visible label for attachment routing."""
    if attachment_mode == "Add to RAG":
        return "RAG SubAgent"
    if force_agent_type == AgentType.DATA_SCIENTIST.value:
        return "Data Scientist"
    if force_agent_type == AgentType.RAG_SUB.value:
        return "RAG SubAgent"
    if force_main_agent:
        return "Main Agent"
    if selected_agent == AUTO_AGENT_DISPLAY_NAME:
        return "Auto"
    return re.sub(r"^\[[^\]]+\]\s*", "", selected_agent.split("(")[0].strip()) or "Selected Agent"


def _render_chat_attachment_block(message: dict[str, Any]) -> None:
    """Render attachment chips and lightweight previews for a chat message."""
    attachment_paths = [str(p) for p in message.get("attachment_paths", []) if str(p).strip()]
    if not attachment_paths:
        return

    mode_label = str(message.get("attachment_mode", "") or "").strip()
    route_label = str(message.get("attachment_route", "") or "").strip()
    names = [Path(path).name for path in attachment_paths]
    summary_parts = [f"`{name}`" for name in names]
    if mode_label:
        summary_parts.append(f"mode: `{mode_label}`")
    if route_label:
        summary_parts.append(f"routed to: `{route_label}`")
    st.caption(" · ".join(summary_parts))

    image_paths = []
    other_paths = []
    for rel in attachment_paths:
        abs_path = PROJECT_DIR / rel
        suffix = abs_path.suffix.lower()
        if suffix in CHAT_IMAGE_SUFFIXES and abs_path.exists():
            image_paths.append(str(abs_path))
        else:
            other_paths.append(rel)

    for image_path in image_paths[:3]:
        st.image(image_path, width="stretch")

    if other_paths:
        with st.expander(f"Attachments ({len(other_paths)})", expanded=False):
            for rel in other_paths:
                abs_path = PROJECT_DIR / rel
                size_label = ""
                if abs_path.exists():
                    try:
                        size_label = f" · {abs_path.stat().st_size / 1024:.1f} KB"
                    except Exception:
                        size_label = ""
                st.caption(f"`{Path(rel).name}`{size_label}")


def _attachment_route_for_chat(
    *,
    relpaths: list[str],
    attachment_mode: str,
    selected_agent: str,
) -> tuple[str | None, bool]:
    """Return (force_agent_type, force_main_agent) for a chat turn with attachments."""
    if not relpaths:
        return None, False

    suffixes = {Path(rel).suffix.lower() for rel in relpaths}
    tabular_only = bool(suffixes) and suffixes.issubset(CHAT_TABULAR_SUFFIXES)
    document_like = any(suffix in CHAT_DOCUMENT_SUFFIXES for suffix in suffixes)
    image_only = bool(suffixes) and suffixes.issubset(CHAT_IMAGE_SUFFIXES)
    single_pdf = len(relpaths) == 1 and suffixes == CHAT_PDF_SUFFIXES

    if attachment_mode == "Add to RAG":
        return AgentType.RAG_SUB.value, False

    if selected_agent == AUTO_AGENT_DISPLAY_NAME:
        if tabular_only:
            return AgentType.DATA_SCIENTIST.value, False
        if single_pdf and attachment_mode == "Ask in chat":
            return None, True
        if document_like:
            return AgentType.RAG_SUB.value, False
        if image_only:
            return None, True
        return AgentType.RAG_SUB.value, False

    # Even with explicit agent selection, allow orchestration and delegation
    return None, True


def _build_attachment_augmented_message(
    *,
    user_message: str,
    relpaths: list[str],
    attachment_mode: str,
    direct_image_context: str = "",
    direct_pdf_context: str = "",
) -> str:
    """Inject attachment instructions into the message sent to the agent."""
    if not relpaths:
        return user_message

    names = ", ".join(Path(rel).name for rel in relpaths)
    prefixes = [f"[CHAT ATTACHMENTS: {', '.join(relpaths)}]"]
    suffixes = {Path(rel).suffix.lower() for rel in relpaths}
    image_only = bool(suffixes) and suffixes.issubset(CHAT_IMAGE_SUFFIXES)
    document_like = any(suffix in CHAT_DOCUMENT_SUFFIXES for suffix in suffixes)
    single_pdf = len(relpaths) == 1 and suffixes == CHAT_PDF_SUFFIXES

    if attachment_mode == "Add to RAG":
        prefixes.append(f"[RAG UPLOADED FILES: {', '.join(relpaths)}]")
        prefixes.append(
            "[ATTACHMENT MODE: Ingest these files into RAG first, then answer from the indexed content.]"
        )
    else:
        if suffixes and suffixes.issubset(CHAT_TABULAR_SUFFIXES):
            prefixes.append(f"[UPLOADED FILE: {relpaths[0]}]")
            prefixes.append(f"[UPLOADED FILES: {', '.join(relpaths)}]")
            prefixes.append("[DATA CHECK: Use execute_python_code for exact dataset shape/statistics.]")
        elif image_only and direct_image_context:
            prefixes.append(
                "[ATTACHMENT MODE: Use the direct OCR/vision extraction context below to answer the user's request without RAG ingestion.]"
            )
            prefixes.append(
                "[WORKFLOW RULE: Use the direct image context first. Answer the user's actual request from that extracted content. "
                "Do not return the raw OCR text alone unless explicitly requested.]"
            )
            prefixes.append(direct_image_context)
        elif single_pdf and direct_pdf_context:
            prefixes.append(
                "[ATTACHMENT MODE: Use the direct PDF extraction context below to answer the user's request without RAG ingestion.]"
            )
            prefixes.append(
                "[WORKFLOW RULE: Use the direct PDF context first and answer the user's actual request from it. "
                "Do not ask the user to re-upload or paste text when extracted context is provided.]"
            )
            prefixes.append(direct_pdf_context)
        else:
            if document_like:
                prefixes.append(f"[RAG UPLOADED FILES: {', '.join(relpaths)}]")
                prefixes.append(
                    "[ATTACHMENT MODE: Document extraction required. Use RAG ingestion/retrieval to read attached PDFs/documents before answering.]"
                )
                prefixes.append(
                    "[WORKFLOW RULE: Do not rely on read_file for binary documents (PDF/DOCX/PPTX). "
                    "Ingest with the RAG specialist first, then answer from extracted text. "
                    "If extraction is empty, report the extraction result explicitly. ]"
                )
            # Only force RAG workflow if explicitly requested via "Add to RAG" mode
            elif attachment_mode == "Add to RAG":
                prefixes.append(f"[RAG UPLOADED FILES: {', '.join(relpaths)}]")
                prefixes.append(
                    "[ATTACHMENT MODE: Use RAG ingestion/retrieval to read the attached documents, images, or presentation files before answering.]"
                )
                if image_only:
                    prefixes.append(
                        "[ATTACHMENT GOAL: Extract readable text from the image using OCR and describe visible content before answering.]"
                    )
                prefixes.append(
                    "[WORKFLOW SEQUENCE: required]\n"
                    "[WORKFLOW RULE: First delegate to the RAG specialist to ingest and process the attached files, including OCR/vision for images. "
                    "Then fulfill the user's actual request using that processed content. Do not stop after OCR extraction or return only intermediate text unless the user explicitly asked for raw extraction. "
                    "Do not claim the file is unreadable until the RAG/image-processing path has been attempted.]\n"
                    "1. 📚 RAG Search (RAG SubAgent)\n"
                    "2. Answer the user's request from the processed attachment content"
                )
            else:
                # "Ask in chat" mode: just notify agent about attachments, let it decide
                prefixes.append(f"[UPLOADED FILES: {', '.join(relpaths)}]")
                prefixes.append(
                    "[ATTACHMENT MODE: Chat mode - handle these attachments however is most helpful. "
                    "You can process, analyze, or reference them as appropriate to answer the user's request.]"
                )

    prefixes.append(f"[ATTACHMENT NAMES: {names}]")
    prefixes.append(f"[ORIGINAL USER REQUEST: {user_message}]")
    return "\n".join(prefixes) + "\n\n" + user_message


def extract_agent_display_text(result: dict) -> str:
    """Extract a user-visible text response from agent results with safe fallbacks."""
    from langchain_core.messages import AIMessage, ToolMessage

    messages = result.get("messages", []) if isinstance(result, dict) else []

    # 1. Try the imported extractor (handles tool_call skipping + ToolMessage fallback)
    _, _, _, _, _, extract_last_ai_text, _ = get_cached_agent_bundle(_agent_bundle_config_signature())
    final_text = extract_last_ai_text(messages)
    if final_text:
        return final_text

    # 2. Walk messages ourselves as a safety net
    #    Look for *any* AIMessage with non-empty text (not just the last one)
    for msg in reversed(messages):
        if isinstance(msg, AIMessage):
            c = msg.content
            if isinstance(c, str) and c.strip():
                return c.strip()
            if isinstance(c, list):
                parts = [p.get("text", "") for p in c if isinstance(p, dict) and p.get("type") == "text"]
                joined = "\n".join(p for p in parts if p)
                if joined:
                    return joined

    # 3. Collect useful ToolMessage outputs (task delegation results, code output, etc.)
    tool_parts: list[str] = []
    for msg in reversed(messages):
        if isinstance(msg, ToolMessage):
            c = msg.content
            if isinstance(c, str) and c.strip():
                tool_parts.append(c.strip())
            if len(tool_parts) >= 3:
                break
    if tool_parts:
        return "\n\n---\n\n".join(tool_parts)

    # 4. Legacy dict fallbacks
    if isinstance(result, dict):
        output = result.get("output")
        if isinstance(output, str) and output.strip():
            return output.strip()
        if isinstance(output, list):
            parts = [str(item).strip() for item in output if str(item).strip()]
            if parts:
                return "\n".join(parts)

        for key in ("response", "content", "text"):
            value = result.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()

    return ""


def _extract_selected_chunk_details(rag_context: str) -> list[dict[str, str]]:
    """Return selected chunk headers, metadata, summaries, and full content."""
    if not isinstance(rag_context, str) or not rag_context.startswith("Retrieved "):
        return []

    lines = rag_context.splitlines()
    chunks: list[dict[str, str]] = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if re.match(r"^\[R\d+\]\s+", line):
            header = line
            topics = ""
            asset_path = ""
            modality = "text"
            header_match = re.match(r"^\[(R\d+)\]\s+([^|]+)\s+\|\s+(.+)$", line)
            if header_match:
                modality = header_match.group(2).strip().lower()
            content_lines: list[str] = []
            j = i + 1
            while j < len(lines):
                raw_next = lines[j]
                next_line = raw_next.strip()
                if next_line == "---" or next_line == "References:" or re.match(r"^\[R\d+\]\s+", next_line):
                    break
                if next_line.startswith("Topics:"):
                    topics = next_line
                elif next_line.startswith("Asset:"):
                    asset_path = next_line.removeprefix("Asset:").strip()
                elif next_line:
                    content_lines.append(raw_next.rstrip())
                j += 1
            content = "\n".join(content_lines).strip()
            snippet = re.sub(r"\s+", " ", content)
            if len(snippet) > 180:
                snippet = snippet[:177].rstrip() + "..."
            summary = f"{header} :: {snippet}" if snippet else header
            chunks.append(
                {
                    "header": header,
                    "topics": topics,
                    "asset_path": asset_path,
                    "modality": modality,
                    "summary": summary,
                    "content": content,
                }
            )
            i = j
            continue
        i += 1
    return chunks


def _render_selected_chunks_expander(chunk_details: list[dict[str, str]]) -> None:
    """Render the selected chunk content once in a bottom expander."""
    if not chunk_details:
        return

    with st.expander("🔍 Selected Chunks", expanded=False):
        for chunk in chunk_details:
            st.markdown(f"**{chunk['header']}**")
            if chunk["topics"]:
                st.caption(chunk["topics"])
            if chunk["content"]:
                st.code(chunk["content"], language="text")
            elif chunk["asset_path"]:
                st.caption(f"Asset: {chunk['asset_path']}")


def _render_rag_asset_previews(chunk_details: list[dict[str, str]]) -> list[str]:
    """Render multimodal RAG assets and return resolved asset paths."""
    rendered_assets: list[str] = []
    for chunk in chunk_details:
        modality = chunk.get("modality")
        asset_rel = chunk.get("asset_path", "").strip()
        if not asset_rel:
            continue
        asset_path = PROJECT_DIR / Path(asset_rel)
        if not asset_path.exists():
            continue
        st.markdown(f"**{chunk['header']}**")
        if modality == "image":
            st.image(str(asset_path), width="stretch")
        elif modality == "table":
            try:
                table_text = asset_path.read_text(encoding="utf-8").strip()
            except Exception:
                table_text = ""
            if table_text:
                st.code(table_text, language="text")
            else:
                st.caption(f"Table asset saved at {asset_rel}")
        else:
            continue
        rendered_assets.append(str(asset_path))
    return rendered_assets


@st.cache_resource
def _get_coqui_tts_engine():
    """Load the Coqui TTS engine once per Streamlit process."""
    from TTS.api import TTS

    model_name = os.getenv("COQUI_TTS_MODEL", "tts_models/en/ljspeech/tacotron2-DDC")
    return TTS(model_name=model_name, progress_bar=False, gpu=False)


# ── Pre-cache expensive heavy-lift functions ──
# These decorators ensure YouTube transcript fetching, TTS, and other heavy ops
# only initialize once, not on every Streamlit rerun.
@st.cache_data(ttl=3600, show_spinner=False)
def _cached_fetch_youtube_transcript(video_id: str) -> str:
    """Cache YouTube transcripts for 1 hour to avoid re-fetching."""
    try:
        items = _fetch_youtube_transcript_items(video_id)
        return _transcript_items_to_text(items)
    except Exception as exc:
        return f"[ERROR] Failed to fetch transcript: {exc}"


@st.cache_resource(show_spinner=False)
def _get_coqui_tts_engine_alt():
    """Alias for compatibility with existing code."""
    return _get_coqui_tts_engine()


def _coqui_tts_available() -> tuple[bool, str]:
    """Return whether Coqui TTS is usable plus a short status message."""
    try:
        sidecar_python = _get_coqui_tts_sidecar_python()
        sidecar_script = PROJECT_DIR / "coqui_tts_sidecar.py"
        if sidecar_python and sidecar_script.exists():
            return True, f"Coqui TTS sidecar ready ({sidecar_python.parent.parent.name})"
        _get_coqui_tts_engine()
        return True, "Coqui TTS ready in current environment"
    except Exception as exc:
        return False, str(exc)


def _get_coqui_tts_sidecar_python() -> Path | None:
    """Return the preferred sidecar Python interpreter for Coqui TTS."""
    candidates = [
        PROJECT_DIR / ".venv" / "Scripts" / "python.exe",
        PROJECT_DIR / ".venv-coqui" / "Scripts" / "python.exe",
        PROJECT_DIR / ".venv311" / "Scripts" / "python.exe",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _plain_text_for_voice(message_content: str) -> str:
    """Reduce markdown-heavy assistant output to speech-friendly plain text."""
    text = strip_react_markers(str(message_content or ""))
    text = re.sub(r"```[\s\S]*?```", " Code block omitted. ", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^\s*[-*]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*#{1,6}\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"\[(R\d+)\]", r"\1", text)
    text = re.sub(r"[*_>]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _synthesize_coqui_tts_audio(message_content: str, key: str) -> bytes | None:
    """Generate WAV audio bytes for an assistant reply via Coqui TTS."""
    speech_text = _plain_text_for_voice(message_content)
    if not speech_text:
        return None

    cache_dir = PROJECT_DIR / "tmp" / "tts_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    model_name = os.getenv("COQUI_TTS_MODEL", "tts_models/en/ljspeech/tacotron2-DDC")
    cache_key = hashlib.sha1(f"{model_name}|{speech_text}|{key}".encode("utf-8")).hexdigest()
    output_path = cache_dir / f"{cache_key}.wav"

    if not output_path.exists():
        sidecar_python = _get_coqui_tts_sidecar_python()
        sidecar_script = PROJECT_DIR / "coqui_tts_sidecar.py"
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False, dir=cache_dir) as tmp_file:
            tmp_path = Path(tmp_file.name)
        try:
            if sidecar_python and sidecar_script.exists():
                subprocess.run(
                    [
                        str(sidecar_python),
                        str(sidecar_script),
                        "--out",
                        str(tmp_path),
                        "--model",
                        model_name,
                    ],
                    input=speech_text,
                    text=True,
                    capture_output=True,
                    check=True,
                )
            else:
                engine = _get_coqui_tts_engine()
                engine.tts_to_file(text=speech_text, file_path=str(tmp_path))
            tmp_path.replace(output_path)
        except subprocess.CalledProcessError as exc:
            raise RuntimeError(
                f"Coqui sidecar synthesis failed: {(exc.stderr or exc.stdout or str(exc)).strip()}"
            ) from exc
        finally:
            if tmp_path.exists():
                tmp_path.unlink(missing_ok=True)

    return output_path.read_bytes()


def _render_speak_button(message_content: str, key: str, auto_speak: bool = False) -> None:
    """Render Coqui TTS playback controls for assistant replies."""
    speech_audio = _synthesize_coqui_tts_audio(message_content, key)
    if not speech_audio:
        return

    st.audio(speech_audio, format="audio/wav")

    if auto_speak:
        b64_audio = base64.b64encode(speech_audio).decode("ascii")
        dom_id = f"coqui-audio-{hashlib.sha1(key.encode('utf-8')).hexdigest()[:10]}"
        html = f"""
        <audio id="{dom_id}" autoplay style="display:none;">
          <source src="data:audio/wav;base64,{b64_audio}" type="audio/wav">
        </audio>
        <script>
        const el = document.getElementById("{dom_id}");
        if (el) {{
          const playPromise = el.play();
          if (playPromise) {{
            playPromise.catch(() => {{}});
          }}
        }}
        </script>
        """
        components.html(html, height=0)


def _render_browser_voice_chat_button() -> None:
    """Attach a microphone button directly to Streamlit's chat input."""
    html = """
    <script>
    (function() {
      const rootDoc = window.parent.document;
      const SpeechRecognition = window.parent.SpeechRecognition || window.parent.webkitSpeechRecognition;
      if (!SpeechRecognition) return;

      function findTextarea() {
        return rootDoc.querySelector('[data-testid="stChatInput"] textarea');
      }

      function writeToChatBox(text) {
        const textarea = findTextarea();
        if (!textarea) return;
        const nativeSetter = Object.getOwnPropertyDescriptor(window.parent.HTMLTextAreaElement.prototype, 'value').set;
        nativeSetter.call(textarea, text);
        textarea.dispatchEvent(new Event('input', { bubbles: true }));
        textarea.dispatchEvent(new Event('change', { bubbles: true }));
        textarea.focus();
      }

      function ensureStyle() {
        if (rootDoc.getElementById('main-chat-voice-style')) return;
        const style = rootDoc.createElement('style');
        style.id = 'main-chat-voice-style';
        style.textContent = `
          [data-testid="stChatInput"] { position: relative; }
          .main-chat-voice-btn {
            position: absolute;
            right: 3.2rem;
            top: 50%;
            transform: translateY(-50%);
            z-index: 1000;
            border: 1px solid #bbb;
            background: #fff;
            border-radius: 999px;
            width: 2.2rem;
            height: 2.2rem;
            cursor: pointer;
            font-size: 1rem;
            line-height: 1;
          }
          .main-chat-voice-btn.listening {
            background: #ffe9e9;
            border-color: #d66;
          }
        `;
        rootDoc.head.appendChild(style);
      }

      let recognition = null;
      let finalTranscript = '';
      let listening = false;

      function attachButton() {
        const chatHost = rootDoc.querySelector('[data-testid="stChatInput"]');
        if (!chatHost) return;

        ensureStyle();

        let button = rootDoc.getElementById('main-chat-voice-btn');
        if (!button) {
          button = rootDoc.createElement('button');
          button.type = 'button';
          button.id = 'main-chat-voice-btn';
          button.className = 'main-chat-voice-btn';
          button.title = 'Speak into chat';
          button.setAttribute('aria-label', 'Speak into chat');
          button.textContent = '🎙️';
        }

        if (!recognition) {
          recognition = new SpeechRecognition();
          recognition.lang = 'en-US';
          recognition.interimResults = true;
          recognition.continuous = false;

          recognition.onstart = function() {
            finalTranscript = '';
            listening = true;
            button.classList.add('listening');
            button.textContent = '⏺';
          };

          recognition.onresult = function(event) {
            let interim = '';
            for (let i = event.resultIndex; i < event.results.length; i++) {
              const transcript = event.results[i][0].transcript;
              if (event.results[i].isFinal) {
                finalTranscript += transcript + ' ';
              } else {
                interim += transcript;
              }
            }
            writeToChatBox((finalTranscript + interim).trim());
          };

          recognition.onend = function() {
            listening = false;
            button.classList.remove('listening');
            button.textContent = '🎙️';
            if (finalTranscript.trim()) {
              writeToChatBox(finalTranscript.trim());
            }
          };

          recognition.onerror = function() {
            listening = false;
            button.classList.remove('listening');
            button.textContent = '🎙️';
          };
        }

        if (!button.dataset.bound) {
          button.addEventListener('click', function() {
            if (listening) {
              recognition.stop();
              return;
            }
            try {
              recognition.start();
            } catch (err) {
              console.warn('Speech recognition start failed', err);
            }
          });
          button.dataset.bound = 'true';
        }

        const mountTarget = chatHost.querySelector('form') || chatHost;
        if (!mountTarget.contains(button)) {
          mountTarget.appendChild(button);
        }
      }

      attachButton();
      const observer = new MutationObserver(() => attachButton());
      observer.observe(rootDoc.body, { childList: true, subtree: true });
    })();
    </script>
    """
    components.html(html, height=1)


@st.cache_data(ttl=30)
def get_ollama_models():
    """Fetch available models from Ollama.

    Returns:
        dict: Mapping of model names to model IDs, or empty dict if Ollama unavailable
    """
    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=3)
        if response.status_code == 200:
            models = response.json().get("models", [])
            return {model["name"]: f"ollama:{model['name']}" for model in models}
    except (requests.RequestException, Exception):
        pass
    return {}


@st.cache_data(ttl=60)
def get_runtime_feature_status() -> dict:
    """Compute runtime feature status shown in the sidebar."""
    cache_enabled = os.getenv("LANGCHAIN_ENABLE_CACHE", "").lower() == "true"
    cache_backend_available = False
    cache_error = ""

    try:
        from langchain_core.caches import InMemoryCache  # noqa: F401
        from langchain_core.globals import set_llm_cache  # noqa: F401
        cache_backend_available = True
    except Exception as e:
        cache_error = str(e)

    if not cache_enabled:
        cache_status = "Prompt caching disabled (`LANGCHAIN_ENABLE_CACHE` is not `true`)"
        cache_level = "warning"
    elif cache_backend_available:
        cache_status = "Prompt caching available (`InMemoryCache` backend detected)"
        cache_level = "success"
    else:
        cache_status = f"Prompt caching requested but backend import failed: {cache_error}"
        cache_level = "error"

    try:
        deepagents_ver = pkg_version("deepagents")
    except PackageNotFoundError:
        deepagents_ver = "unknown"

    major_minor = [0, 0]
    for i, part in enumerate(deepagents_ver.split(".")[:2]):
        digits = "".join(ch for ch in part if ch.isdigit())
        major_minor[i] = int(digits) if digits else 0

    if major_minor[0] > 0 or major_minor[1] >= 6:
        permissions_status = (
            f"DeepAgents {deepagents_ver}: version may support finer-grained filesystem permissions"
        )
        permissions_level = "info"
    else:
        permissions_status = (
            f"DeepAgents {deepagents_ver}: using `virtual_mode` sandbox "
            "(fine-grained per-subagent permissions not integrated here)"
        )
        permissions_level = "warning"

    return {
        "cache_status": cache_status,
        "cache_level": cache_level,
        "permissions_status": permissions_status,
        "permissions_level": permissions_level,
        "deepagents_version": deepagents_ver,
    }

# Page configuration
st.set_page_config(
    page_title="LocalAiLab Orchestrator",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS for a cleaner, more professional product UI
st.markdown("""
    <style>
    :root {
        --app-bg: #0b1120;
        --panel-bg: #111827;
        --panel-border: rgba(148, 163, 184, 0.18);
        --text-primary: #f8fafc;
        --text-secondary: #cbd5e1;
        --text-muted: #94a3b8;
        --accent: #60a5fa;
        --accent-soft: rgba(96, 165, 250, 0.16);
    }

    html, body, [class*="stApp"] {
        font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Oxygen,
            Ubuntu, Cantarell, "Helvetica Neue", sans-serif;
        color: var(--text-primary);
        background:
            radial-gradient(circle at top left, rgba(96, 165, 250, 0.14), transparent 24%),
            radial-gradient(circle at bottom right, rgba(15, 118, 110, 0.12), transparent 22%),
            linear-gradient(180deg, #0f172a 0%, var(--app-bg) 100%);
    }

    .block-container {
        padding-top: 1.35rem;
        padding-bottom: 2rem;
        max-width: 1280px;
    }

    section[data-testid="stSidebar"] .block-container {
        padding-top: 1rem;
    }

    [data-testid="stHeader"] {
        background: transparent;
    }

    .page-shell,
    .section-card,
    .feature-card {
        background: var(--panel-bg);
        border: 1px solid var(--panel-border);
        border-radius: 20px;
        box-shadow: 0 16px 40px rgba(2, 6, 23, 0.35);
    }

    .page-shell {
        padding: 1.2rem 1.4rem;
        margin-bottom: 1rem;
        color: var(--text-primary);
    }

    .page-kicker {
        display: inline-block;
        margin-bottom: 0.35rem;
        color: var(--accent);
        font-size: 0.72rem;
        font-weight: 700;
        letter-spacing: 0.12em;
        text-transform: uppercase;
    }

    .page-title {
        margin: 0;
        font-size: 1.85rem;
        font-weight: 750;
        line-height: 1.15;
        letter-spacing: -0.03em;
        color: var(--text-primary);
    }

    .page-subtitle {
        margin-top: 0.45rem;
        color: var(--text-secondary);
        font-size: 0.96rem;
        line-height: 1.55;
    }

    .section-card,
    .feature-card {
        padding: 1rem 1.05rem;
    }

    .section-header {
        display: flex;
        align-items: baseline;
        justify-content: space-between;
        gap: 0.75rem;
        margin-bottom: 0.65rem;
    }

    .section-title {
        margin: 0;
        font-size: 0.98rem;
        font-weight: 700;
        color: var(--text-primary);
    }

    .section-note {
        margin: 0;
        color: var(--text-muted);
        font-size: 0.82rem;
    }

    .summary-row {
        display: flex;
        flex-wrap: wrap;
        gap: 0.5rem;
        margin-top: 0.75rem;
    }

    .summary-pill {
        display: inline-flex;
        align-items: center;
        gap: 0.35rem;
        padding: 0.34rem 0.7rem;
        border-radius: 999px;
        border: 1px solid rgba(96, 165, 250, 0.24);
        background: rgba(96, 165, 250, 0.10);
        color: var(--text-primary);
        font-size: 0.8rem;
        font-weight: 600;
    }

    .clean-caption {
        color: var(--text-muted);
        font-size: 0.85rem;
        line-height: 1.45;
    }

    .stChatMessage {
        background-color: transparent;
    }

    hr {
        margin: 1.25rem 0;
        border: none;
        border-top: 1px solid rgba(148, 163, 184, 0.22);
    }

    .stButton > button {
        border-radius: 10px;
        font-weight: 600;
        transition: transform 0.15s ease, box-shadow 0.15s ease, border-color 0.15s ease;
    }

    .stButton > button:hover {
        transform: translateY(-1px);
        box-shadow: 0 8px 18px rgba(15, 23, 42, 0.12);
    }

    .hero-shell {
        background:
            linear-gradient(180deg, rgba(37, 99, 235, 0.16), rgba(17, 24, 39, 1)),
            var(--panel-bg);
        border-color: rgba(96, 165, 250, 0.22);
    }

    .hero-shell h1,
    .hero-shell .page-subtitle,
    .hero-shell .hero-meta {
        color: var(--text-primary);
    }

    .hero-shell .hero-meta {
        margin-top: 0.75rem;
        color: var(--text-secondary);
        font-size: 0.84rem;
        font-weight: 600;
    }

    .feature-card {
        background: #0f172a;
    }

    .feature-card h3 {
        color: var(--text-primary);
        font-size: 1rem;
        margin: 0 0 0.35rem;
        font-weight: 700;
    }

    .feature-card p {
        color: var(--text-secondary);
        font-size: 0.92rem;
        line-height: 1.55;
        margin: 0;
    }

    .stExpander {
        border-color: var(--panel-border);
    }

    .stExpander [data-testid="stExpanderDetails"] {
        background: rgba(15, 23, 42, 0.28);
        border-radius: 0 0 16px 16px;
    }

    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #0f172a 0%, #111827 100%);
    }

    [data-testid="stSidebar"] * {
        color: var(--text-primary);
    }

    [data-testid="stSidebar"] .stCaption,
    [data-testid="stSidebar"] .stMarkdown p,
    [data-testid="stSidebar"] label {
        color: var(--text-secondary);
    }
    </style>
    """, unsafe_allow_html=True)

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []  # Recent messages only (trimmed to last 50 for faster serialization)
if "messages_archive_count" not in st.session_state:
    st.session_state.messages_archive_count = 0  # Track total count for UI display
if "thread_id" not in st.session_state:
    st.session_state.thread_id = "streamlit-chat-thread"
if "user_id" not in st.session_state:
    st.session_state.user_id = resolve_user_id(default_user_id="streamlit-user")
if "current_model" not in st.session_state:
    st.session_state.current_model = get_agent_model("main")
if "selected_agent" not in st.session_state:
    st.session_state.selected_agent = AUTO_AGENT_DISPLAY_NAME
if "voice_chat_enabled" not in st.session_state:
    st.session_state.voice_chat_enabled = False
if "voice_auto_speak" not in st.session_state:
    st.session_state.voice_auto_speak = True
if "voice_speech_rate" not in st.session_state:
    st.session_state.voice_speech_rate = 1.0
if "data_scientist_model" not in st.session_state:
    st.session_state.data_scientist_model = get_agent_model_override("data_scientist")
if "last_execution_time" not in st.session_state:
    st.session_state.last_execution_time = None
if "total_messages" not in st.session_state:
    st.session_state.total_messages = 0
if "uploaded_file_path" not in st.session_state:
    st.session_state.uploaded_file_path = None
if "rag_uploaded_files" not in st.session_state:
    st.session_state.rag_uploaded_files = []
if "main_chat_attachment_mode" not in st.session_state:
    st.session_state.main_chat_attachment_mode = "Ask in chat"
if "rag_project_mode" not in st.session_state:
    st.session_state.rag_project_mode = "All Projects"
if "rag_query_themes" not in st.session_state:
    st.session_state.rag_query_themes = []
if "rag_project_filter_selector" not in st.session_state:
    st.session_state.rag_project_filter_selector = ""
if "rag_query_modalities" not in st.session_state:
    st.session_state.rag_query_modalities = []
if "rag_knowledge_only" not in st.session_state:
    st.session_state.rag_knowledge_only = False
_RAG_QA_PRESET_NAME = "💬 Q&A"
_RAG_QA_PRESET = {
    "mode": "Top-K Globally",
    "top_k": 5,
    "fetch_k": 80,
    "max_files": 5,
    "min_rerank": 0.10,
}
_RAG_DEFAULTS_VERSION = 4
if "rag_retrieval_mode" not in st.session_state:
    st.session_state.rag_retrieval_mode = _RAG_QA_PRESET["mode"]
if "rag_top_k" not in st.session_state:
    st.session_state.rag_top_k = _RAG_QA_PRESET["top_k"]
if "rag_fetch_k" not in st.session_state:
    st.session_state.rag_fetch_k = _RAG_QA_PRESET["fetch_k"]
if "rag_max_files" not in st.session_state:
    st.session_state.rag_max_files = _RAG_QA_PRESET["max_files"]
if "rag_min_rerank_score" not in st.session_state:
    st.session_state.rag_min_rerank_score = _RAG_QA_PRESET["min_rerank"]
# Slider widget keys — initialized independently so hot-reloads and existing
# sessions pick up the correct values even if the data keys were already set.
if "rag_top_k_slider_v2" not in st.session_state:
    st.session_state.rag_top_k_slider_v2 = st.session_state.rag_top_k
if "rag_fetch_k_slider_v2" not in st.session_state:
    st.session_state.rag_fetch_k_slider_v2 = st.session_state.rag_fetch_k
if "rag_max_files_slider_v2" not in st.session_state:
    st.session_state.rag_max_files_slider_v2 = st.session_state.rag_max_files
if "rag_min_rerank_score_slider_v2" not in st.session_state:
    st.session_state.rag_min_rerank_score_slider_v2 = st.session_state.rag_min_rerank_score
if "rag_active_preset" not in st.session_state:
    st.session_state.rag_active_preset = _RAG_QA_PRESET_NAME
if "rag_ui_initialized" not in st.session_state:
    st.session_state.rag_ui_initialized = False
if "rag_chat_target" not in st.session_state:
    st.session_state.rag_chat_target = "ragsub"
if "rag_main_agent_live_progress" not in st.session_state:
    st.session_state.rag_main_agent_live_progress = True
if "main_chat_use_direct_model" not in st.session_state:
    st.session_state.main_chat_use_direct_model = st.session_state.get("main_agent_fast_mode", True)
if st.session_state.get("rag_defaults_version") != _RAG_DEFAULTS_VERSION:
    for _rag_widget_key in (
        "rag_top_k_slider",
        "rag_fetch_k_slider",
        "rag_max_files_slider",
        "rag_min_rerank_score_slider",
        "rag_top_k_slider_v2",
        "rag_fetch_k_slider_v2",
        "rag_max_files_slider_v2",
        "rag_min_rerank_score_slider_v2",
    ):
        st.session_state.pop(_rag_widget_key, None)
    st.session_state.rag_retrieval_mode = _RAG_QA_PRESET["mode"]
    st.session_state.rag_top_k = _RAG_QA_PRESET["top_k"]
    st.session_state.rag_fetch_k = _RAG_QA_PRESET["fetch_k"]
    st.session_state.rag_max_files = _RAG_QA_PRESET["max_files"]
    st.session_state.rag_min_rerank_score = _RAG_QA_PRESET["min_rerank"]
    st.session_state.rag_top_k_slider_v2 = _RAG_QA_PRESET["top_k"]
    st.session_state.rag_fetch_k_slider_v2 = _RAG_QA_PRESET["fetch_k"]
    st.session_state.rag_max_files_slider_v2 = _RAG_QA_PRESET["max_files"]
    st.session_state.rag_min_rerank_score_slider_v2 = _RAG_QA_PRESET["min_rerank"]
    st.session_state.rag_active_preset = _RAG_QA_PRESET_NAME
    st.session_state.rag_ui_initialized = False
    st.session_state.rag_defaults_version = _RAG_DEFAULTS_VERSION
if "rag_enable_main_chat" not in st.session_state:
    st.session_state.rag_enable_main_chat = False
if "main_agent_fast_mode" not in st.session_state:
    st.session_state.main_agent_fast_mode = True
if "disable_agent_cache" not in st.session_state:
    st.session_state.disable_agent_cache = False
if "force_write_todos" not in st.session_state:
    st.session_state.force_write_todos = False
if "self_learning_enabled" not in st.session_state:
    st.session_state.self_learning_enabled = False
if "pipeline_steps" not in st.session_state:
    st.session_state.pipeline_steps = []  # ordered list of step display names
if "active_mode" not in st.session_state:
    st.session_state.active_mode = "main"
if "scan_to_text_result" not in st.session_state:
    st.session_state.scan_to_text_result = ""
if "scan_to_text_sources" not in st.session_state:
    st.session_state.scan_to_text_sources = []
if "rag_file_summary_selected_file" not in st.session_state:
    st.session_state.rag_file_summary_selected_file = ""
if "rag_file_summary_context" not in st.session_state:
    st.session_state.rag_file_summary_context = ""
if "rag_ocr_engine_mode" not in st.session_state:
    st.session_state.rag_ocr_engine_mode = "Auto"
if "rag_image_ingest_mode" not in st.session_state:
    st.session_state.rag_image_ingest_mode = "OCR only"
if "model_main" not in st.session_state:
    st.session_state.model_main = get_agent_model("main")
if "model_rag" not in st.session_state:
    st.session_state.model_rag = get_agent_model_override("ragsub")
if "model_rag_embed" not in st.session_state:
    st.session_state.model_rag_embed = os.getenv("RAG_EMBED_MODEL", "llama_server:")
if "model_rag_vision" not in st.session_state:
    st.session_state.model_rag_vision = os.getenv("RAG_VISION_MODEL", "llama_cpp:")
if "model_rag_context" not in st.session_state:
    st.session_state.model_rag_context = os.getenv("RAG_CONTEXT_LLM_MODEL", "")
if "model_data" not in st.session_state:
    st.session_state.model_data = get_agent_model_override("data_scientist")
if "model_literature" not in st.session_state:
    st.session_state.model_literature = os.getenv("LITERATURE_MODEL", "")
for _agent_name in AGENT_MODEL_ENV_KEYS:
    if _agent_name == "main":
        continue
    _session_key = f"model_override_{_agent_name}"
    if _session_key not in st.session_state:
        st.session_state[_session_key] = get_agent_model_override(_agent_name)

os.environ["SELF_LEARNING_ENABLED"] = "true" if st.session_state.self_learning_enabled else "false"
if "llm_provider" not in st.session_state:
    st.session_state.llm_provider = get_default_llm_provider()

CHAT_INDEX_EVERY_N_ASSISTANT_TURNS = 4

# ── Chat History Persistence ──
CHAT_HISTORY_DIR = PROJECT_DIR / "chat_history"
CHAT_HISTORY_DIR.mkdir(exist_ok=True)


def _load_chat_sessions(merge_only: bool = False):
    """Load saved chat sessions from disk.

    Args:
        merge_only: If True, only add sessions not already present in session state
                    (avoids overwriting in-memory state on repeated calls).
    """
    import json
    CHAT_HISTORY_DIR.mkdir(exist_ok=True)
    for fp in sorted(CHAT_HISTORY_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True):
        try:
            sid = fp.stem
            if merge_only and sid in st.session_state.chat_sessions:
                continue
            data = json.loads(fp.read_text(encoding="utf-8"))
            st.session_state.chat_sessions[sid] = {
                "name": data.get("name", sid[:20]),
                "messages": data.get("messages", []),
                "created": data.get("created", ""),
                "model": data.get("model", ""),
            }
        except Exception:
            pass


def _trim_message_history(max_messages: int = 50) -> None:
    """Archive older messages to reduce session state size (speeds up Streamlit reruns).
    
    When session state grows large, Streamlit must serialize/deserialize it on every rerun,
    which adds latency. By keeping only the most recent messages in memory and persisting
    full history to disk, we keep the UI responsive without losing conversation data.
    """
    if len(st.session_state.messages) > max_messages:
        archived = st.session_state.messages[: len(st.session_state.messages) - max_messages]
        st.session_state.messages_archive_count += len(archived)
        st.session_state.messages = st.session_state.messages[-max_messages:]
        # Full history is still persisted to disk, so no data loss

def _save_current_session(*, force_index: bool = False):
    """Save the current chat to disk, deferring full-text indexing when possible."""
    sid = st.session_state.thread_id
    msgs = st.session_state.messages
    if not msgs:
        return
    
    # Trim memory before saving to speed up next rerun
    _trim_message_history(max_messages=50)
    
    assistant_turns = sum(1 for m in msgs if m.get("role") == "assistant")
    should_index = force_index or assistant_turns <= 1 or (
        assistant_turns % CHAT_INDEX_EVERY_N_ASSISTANT_TURNS == 0
    )
    _persist_kwargs = {
        "session_id": sid,
        "messages": [{"role": m["role"], "content": _export_message_content(m)} for m in msgs],
        "model": st.session_state.current_model,
    }
    try:
        session_data = conversation_memory.persist_session(
            **_persist_kwargs,
            update_index=should_index,
        )
    except TypeError:
        # Streamlit hot reload can leave an older ConversationMemoryStore instance
        # alive briefly; fall back to the legacy signature for compatibility.
        session_data = conversation_memory.persist_session(**_persist_kwargs)
    if session_data:
        st.session_state.chat_sessions[sid] = session_data


def _new_chat_session():
    """Start a brand new chat session."""
    import uuid
    if st.session_state.messages:
        _save_current_session(force_index=True)
    new_id = f"chat-{uuid.uuid4().hex[:8]}"
    st.session_state.messages = []
    st.session_state.thread_id = new_id
    st.session_state.active_session_id = new_id


def _switch_to_session(session_id: str):
    """Switch to an existing saved session."""
    if st.session_state.messages:
        _save_current_session(force_index=True)
    session = st.session_state.chat_sessions.get(session_id, {})
    st.session_state.messages = list(session.get("messages", []))
    st.session_state.thread_id = session_id
    st.session_state.active_session_id = session_id


def _delete_session(session_id: str):
    """Delete a saved session from disk and memory."""
    fp = CHAT_HISTORY_DIR / f"{session_id}.json"
    if fp.exists():
        fp.unlink()
    st.session_state.chat_sessions.pop(session_id, None)


if "chat_sessions" not in st.session_state:
    st.session_state.chat_sessions = {}  # {session_id: {name, messages, created, model}}
    _load_chat_sessions()  # populate from disk — but do NOT auto-restore messages
if "active_session_id" not in st.session_state:
    st.session_state.active_session_id = st.session_state.thread_id


def _render_welcome_screen() -> None:
    """Render a professional landing screen shown when no conversation is active."""
    logo_path = PROJECT_DIR / "assets" / "localailab-logo.png"
    if logo_path.exists():
        st.image(str(logo_path), width=220)

    st.markdown(
        """
        <div class="page-shell hero-shell">
            <div class="page-kicker">Main Agent</div>
            <h1>LocalAiLab Orchestrator</h1>
            <p class="page-subtitle">
                A focused workspace for chat, RAG, analysis, coding, and specialist routing.
                The controls on this page shape how the main agent behaves.
            </p>
            <p class="hero-meta">Built with LangGraph · Local-first · Privacy-conscious</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.divider()




def get_status_info() -> dict:
    """Get current system status information.

    Returns:
        Dictionary with status metrics
    """
    plots_dir = PROJECT_DIR / "generated_plots"
    plot_count = len(list(plots_dir.glob("*.png"))) if plots_dir.exists() else 0

    uploaded_file = None
    file_size = None
    if st.session_state.uploaded_file_path:
        file_path = Path(st.session_state.uploaded_file_path)
        if file_path.exists():
            uploaded_file = file_path.name
            file_size = file_path.stat().st_size / 1024  # KB

    return {
        "plot_count": plot_count,
        "message_count": len(st.session_state.messages),
        "user_messages": sum(1 for m in st.session_state.messages if m["role"] == "user"),
        "assistant_messages": sum(1 for m in st.session_state.messages if m["role"] == "assistant"),
        "uploaded_file": uploaded_file,
        "file_size_kb": file_size,
        "execution_time": st.session_state.last_execution_time,
    }


def render_mission_control_dashboard():
    """Render the mission control dashboard at the top of the page."""
    st.markdown("## [GOAL] Mission Control Dashboard")

    # Get status info
    status = get_status_info()

    # Create columns for metrics
    col1, col2, col3, col4, col5 = st.columns(5)

    with col1:
        st.metric(
            label="💬 Total Messages",
            value=status["message_count"],
            delta=status["user_messages"],
            delta_color="normal"
        )

    with col2:
        st.metric(
            label="🤖 Agent Responses",
            value=status["assistant_messages"],
            delta=None,
            delta_color="off"
        )

    with col3:
        st.metric(
            label="📊 Generated Plots",
            value=status["plot_count"],
            delta=None,
            delta_color="off"
        )

    with col4:
        if status["execution_time"]:
            st.metric(
                label="⏱️ Last Response Time",
                value=f"{status['execution_time']:.1f}s",
                delta=None,
                delta_color="off"
            )
        else:
            st.metric(
                label="⏱️ Last Response Time",
                value="—",
                delta=None,
                delta_color="off"
            )

    with col5:
        if status["uploaded_file"]:
            st.metric(
                label="📁 File Loaded",
                value=status["uploaded_file"].split('.')[0],
                delta=f"{status['file_size_kb']:.0f} KB" if status["file_size_kb"] else None,
                delta_color="off"
            )
        else:
            st.metric(
                label="📁 File Loaded",
                value="None",
                delta=None,
                delta_color="off"
            )

    # Current configuration summary
    st.markdown("---")

    cfg_col1 = st.columns(1)[0]

    with cfg_col1:
        st.markdown("**🤖 Active Agent**")
        agent_display = st.session_state.selected_agent
        if agent_display == AUTO_AGENT_DISPLAY_NAME:
            st.markdown("🎯 **Auto** (Main decides)")
        else:
            st.markdown(f"🎯 {agent_display}")

    st.markdown("---")




def _build_rag_messages(user_input: str) -> list[tuple[str, str]]:
    """Build a message list for the RAG sub-agent respecting the history mode selector.
    
    Filters out file block references and attachment metadata to avoid Claude API errors.
    """
    def _clean_message_content(content: str) -> str:
        """Remove file block references and attachment metadata from message content."""
        # Remove attachment metadata tags that could be misinterpreted as file blocks
        content = re.sub(r'\[CHAT ATTACHMENTS:.*?\]', '', content)
        content = re.sub(r'\[RAG UPLOADED FILES:.*?\]', '', content)
        content = re.sub(r'\[UPLOADED FILES:.*?\]', '', content)
        content = re.sub(r'\[UPLOADED FILE:.*?\]', '', content)
        content = re.sub(r'\[ATTACHMENT MODE:.*?\]', '', content)
        content = re.sub(r'\[ATTACHMENT NAMES:.*?\]', '', content)
        content = re.sub(r'\[ATTACHMENT GOAL:.*?\]', '', content)
        content = re.sub(r'\[ORIGINAL USER REQUEST:.*?\]', '', content)
        # Clean up excess whitespace
        content = re.sub(r'\n\n+', '\n\n', content).strip()
        return content
    
    mode = st.session_state.get("rag_history_mode", "No history")
    if mode == "Full history":
        messages = [
            (msg["role"], _clean_message_content(msg["content"])) 
            for msg in st.session_state.messages
            if msg.get("content")
        ]
    elif mode == "Current question only":
        messages = []
    else:  # "No history"
        messages = []
    
    # Clean the user input to remove attachment metadata
    cleaned_input = _clean_message_content(user_input)
    messages.append(("user", cleaned_input))
    return messages


def run_agent_async(user_input: str) -> dict | None:
    """Run the agent with the given user input using registry (non-streaming).

    Used for direct-routed agents (Data Scientist, RAG, Presenter).
    """
    selected_agent_name = st.session_state.selected_agent
    session_id = st.session_state.thread_id

    print(f"[DEBUG] selected_agent = {selected_agent_name}")

    agent_display_options = registry.get_all_agents()
    agent_type_str, _auto_fast_route = _streamlit_chat_execution.resolve_chat_agent_type(
        user_message=user_input,
        selected_agent=selected_agent_name,
        agent_display_options=agent_display_options,
        rag_knowledge_only=st.session_state.rag_knowledge_only,
        force_write_todos=False,  # Agent autonomy: agents decide planning independently
    )

    messages = [(msg["role"], msg["content"]) for msg in st.session_state.messages]
    messages.append(("user", user_input))

    return _streamlit_chat_execution.invoke_agent_route(
        **_build_agent_route_kwargs(
            agent_type=agent_type_str,
            auto_fast_route=_auto_fast_route,
            user_message=user_input,
            history=messages,
            session_id=session_id,
        )
    )


def _run_main_agent(user_input: str, messages: list) -> dict:
    """Run main orchestrator via ainvoke (fallback for non-streaming contexts)."""
    safe_user_input = coerce_message_content_to_text(user_input)
    safe_messages = sanitize_history_pairs(messages)
    _check_llama_server_available()
    agent_obj, _ = _select_main_agent_for_request(safe_user_input)
    session_id = st.session_state.thread_id
    cfg = cast(RunnableConfig, build_agent_config(
        thread_id=session_id,
        user_id=st.session_state.user_id,
        recursion_limit=20,
        default_user_id="streamlit-user",
    ))

    try:
        return _run_async(
            agent_obj.ainvoke(
                {"messages": safe_messages},
                config=cfg,
            )
        )
    except Exception as ex:
        if _is_llama_server_provider_active() and "connection attempts failed" in str(ex).lower():
            raise RuntimeError(
                "llama-server request failed because the server is unreachable. "
                "Start run-llama-server.bat and retry."
            ) from ex
        # When middleware overhead blows past context limits, retry with only
        # the user message against the configured main model.
        if _should_fallback_to_raw_llm(ex):
            return _invoke_raw_main_llm(safe_user_input)
        raise


def _run_direct_main_chat(user_input: str) -> dict:
    """Run the main chat as a plain model call without Deep Agent orchestration."""
    safe_user_input = coerce_message_content_to_text(user_input)
    _check_llama_server_available()
    return _invoke_raw_main_llm(safe_user_input)


def _select_main_agent_for_request(user_input: str) -> tuple[Any, str]:
    """Return the selected main agent and a short mode label for this request."""
    main_agent, trivial, standard, complex_a, lean, _, select_by_complexity = get_cached_agent_bundle(_agent_bundle_config_signature())
    
    track = infer_task_track(user_input, st.session_state.selected_agent, registry.get_all_agents())
    
    # Complete agent autonomy: agents choose complexity independently
    agent_obj, tier = select_by_complexity(
        user_input,
        track=track,
        force_complex=False,  # Disabled: agent chooses complexity
        has_attachments=bool(st.session_state.get("chat_attachment_paths", [])),
        rag_enabled=(st.session_state.rag_enable_main_chat and not st.session_state.rag_knowledge_only),
    )

    # Keep llama.cpp prompts small when context is configured low.
    if os.getenv("DEEPAGENT_LLM_PROVIDER", "").strip().lower() == "llama_cpp":
        try:
            n_ctx = int(os.getenv("LLAMA_CPP_N_CTX", "8192"))
        except ValueError:
            n_ctx = 8192
        if n_ctx <= 16384 and tier == "standard":
            return lean, "lean"
    
    return agent_obj, tier


def _short_preview(value: Any, limit: int = 120) -> str:
    text = str(value or "").strip().replace("\n", " ")
    if len(text) <= limit:
        return text
    return text[: limit - 3] + "..."


def _estimate_tokens_per_second(text: str, exec_time: float) -> tuple[int, float]:
    """Return a rough token estimate and throughput for caption display."""
    token_estimate = max(0, len(str(text or "")) // 4)
    if exec_time <= 0:
        return token_estimate, 0.0
    return token_estimate, token_estimate / exec_time


def _render_response_metrics(
    exec_time: float,
    text: str,
    model_label: str,
    prefix: str = "",
    timing_breakdown: dict[str, float] | None = None,
) -> None:
    """Render a compact response metrics footer."""
    token_estimate, tok_per_sec = _estimate_tokens_per_second(text, exec_time)
    metric_prefix = f"{prefix} · " if prefix else ""
    breakdown_text = ""
    if timing_breakdown:
        breakdown_text = (
            f" · pre {timing_breakdown.get('pre', 0.0):.1f}s"
            f" · agent {timing_breakdown.get('agent', 0.0):.1f}s"
            f" · post {timing_breakdown.get('post', 0.0):.1f}s"
        )
    st.caption(
        f"⏱️ {exec_time:.1f}s{breakdown_text} · ~{token_estimate} tokens · {tok_per_sec:.1f} tok/s · {metric_prefix}{model_label}"
    )


def _build_agent_route_kwargs(
    *,
    agent_type: str,
    auto_fast_route: str | None,
    user_message: str,
    history: list[tuple[str, str]],
    session_id: str,
    status_container: Any = None,
    status_steps: list[str] | None = None,
) -> dict[str, Any]:
    """Build the shared keyword arguments for invoke_agent_route()."""
    route_kwargs = {
        "agent_type": agent_type,
        "auto_fast_route": auto_fast_route,
        "user_message": user_message,
        "history": history,
        "session_id": session_id,
        "user_id": st.session_state.user_id,
        "run_async": _run_async,
        "get_cached_data_scientist_bundle": get_cached_data_scientist_bundle,
        "get_cached_ragsub_bundle": get_cached_ragsub_bundle,
        "get_cached_presenter_bundle": get_cached_presenter_bundle,
        "get_cached_specialist_router_bundle": get_cached_specialist_router_bundle,
        "build_rag_messages": _build_rag_messages,
        "register_analysis_result": (
            lambda kind, payload, agent_name, sid: metadata_store.register_analysis_result(
                kind, payload, agent_name, session_id=sid
            )
        ),
        "run_main_agent_streaming_fn": run_main_agent_streaming,
        "run_main_agent_fn": _run_main_agent,
        "run_direct_main_chat_fn": _run_direct_main_chat,
        "main_chat_use_direct_model": st.session_state.main_chat_use_direct_model,
        "status_container": status_container,
        "status_steps": status_steps,
    }

    try:
        supported_params = set(inspect.signature(_streamlit_chat_execution.invoke_agent_route).parameters)
    except (TypeError, ValueError):
        return route_kwargs

    return {key: value for key, value in route_kwargs.items() if key in supported_params}


def _prepare_artifact_tracking(agent_type: str) -> dict[str, Any]:
    """Capture artifact directory state so new outputs can be detected after execution."""
    plots_dir = PROJECT_DIR / "generated_plots"
    presentations_dir = PROJECT_DIR / "presentations"
    artifact_agent_types = {
        AgentType.DATA_SCIENTIST.value,
        AgentType.RAG_SUB.value,
        AgentType.PRESENTER.value,
    }
    should_scan = agent_type in artifact_agent_types
    presentation_exts = ("*.pptx", "*.pdf", "*.html")
    if not should_scan:
        return {
            "should_scan": False,
            "plots_dir": plots_dir,
            "presentations_dir": presentations_dir,
            "presentation_exts": presentation_exts,
            "before_plots": set(),
            "before_presentations": set(),
        }

    plots_dir.mkdir(exist_ok=True)
    presentations_dir.mkdir(exist_ok=True)
    return {
        "should_scan": True,
        "plots_dir": plots_dir,
        "presentations_dir": presentations_dir,
        "presentation_exts": presentation_exts,
        "before_plots": {p.name for p in plots_dir.glob("*.png")},
        "before_presentations": {
            p.name
            for pattern in presentation_exts
            for p in presentations_dir.glob(pattern)
        },
    }


def _collect_new_artifacts(tracker: dict[str, Any]) -> tuple[list[str], list[str]]:
    """Return newly created plot paths and presentation filenames."""
    if not tracker["should_scan"]:
        return [], []

    new_plot_paths = [
        str(tracker["plots_dir"] / name)
        for name in sorted({p.name for p in tracker["plots_dir"].glob("*.png")} - tracker["before_plots"])
    ]
    new_presentation_files = sorted(
        {
            p.name
            for pattern in tracker["presentation_exts"]
            for p in tracker["presentations_dir"].glob(pattern)
        }
        - tracker["before_presentations"]
    )
    return new_plot_paths, new_presentation_files


def _extract_rag_raw_chunks(
    *,
    agent_type: str,
    result: dict | Any,
    selected_rag_context: str,
) -> str:
    """Extract raw retrieved chunk text for the selected-chunks UI."""
    rag_raw_chunks = selected_rag_context if selected_rag_context.startswith("Retrieved ") else ""
    if agent_type != AgentType.RAG_SUB.value or not isinstance(result, dict):
        return rag_raw_chunks

    from langchain_core.messages import ToolMessage as _ToolMessage

    for msg in result.get("messages", []):
        if (
            isinstance(msg, _ToolMessage)
            and isinstance(msg.content, str)
            and msg.content.startswith("Retrieved ")
        ):
            return msg.content
    return rag_raw_chunks


def _render_generated_artifacts(
    *,
    session_id: str,
    plot_paths: list[str],
    presentation_files: list[str],
    presentations_dir: Path,
) -> None:
    """Render plots/presentations created during the current agent run."""
    if plot_paths:
        st.markdown(f"**📊 Generated Visualizations ({len(plot_paths)}):**")
        for plot_path in plot_paths:
            st.image(plot_path, width="stretch")
            metadata_store.register_plot(
                plot_name=Path(plot_path).stem,
                plot_path=plot_path,
                agent_name="Agent",
                description="Visualization",
                tags=["analysis"],
                session_id=session_id,
            )

    if presentation_files:
        st.markdown("**📊 Generated Presentations:**")
        for artifact_name in presentation_files:
            artifact_path = presentations_dir / artifact_name
            mime = {
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".pdf": "application/pdf",
                ".html": "text/html",
            }.get(artifact_path.suffix.lower(), "application/octet-stream")
            with open(artifact_path, "rb") as artifact_f:
                st.download_button(
                    label=f"📥 {artifact_name}",
                    data=artifact_f.read(),
                    file_name=artifact_name,
                    mime=mime,
                )


def _persist_assistant_response(
    *,
    session_id: str,
    final_text: str,
    plot_paths: list[str],
    rag_asset_paths: list[str],
) -> None:
    """Persist the assistant message and render feedback controls."""
    st.session_state.messages.append(
        {
            "role": "assistant",
            "content": final_text,
            "plots": plot_paths,
            "rag_assets": rag_asset_paths,
        }
    )
    _save_current_session()

    msg_index = len(st.session_state.messages) - 1
    prior_user = next(
        (
            st.session_state.messages[i]["content"]
            for i in range(msg_index - 1, -1, -1)
            if st.session_state.messages[i]["role"] == "user"
        ),
        "",
    )
    _render_feedback_buttons(
        msg_index=msg_index,
        user_prompt=prior_user,
        agent_response=final_text,
    )


def _infer_main_agent_status_steps(user_input: str) -> list[str]:
    steps = ["Reviewing the request", "Choosing whether to answer directly or delegate"]
    track = infer_task_track(user_input, st.session_state.selected_agent, registry.get_all_agents())
    shared_steps = execution_steps_for_track(track, is_auto_mode=True)
    steps.append(
        shared_steps[2].replace("⏳ ", "").replace("...", "")
        if len(shared_steps) > 2
        else "Coordinating the next action across available tools"
    )
    # REMOVED: force_write_todos — agent decides task planning autonomously
    if st.session_state.pipeline_steps:
        steps.append("Following the configured pipeline sequence")
    if st.session_state.rag_enable_main_chat and not st.session_state.rag_knowledge_only:
        steps.append("Using RAG-augmented context if retrieval applies")
    steps.append("Synthesizing the final response")
    return steps


def _format_progress_timestamp(started_at: float, now: float | None = None) -> str:
    current = time.time() if now is None else now
    elapsed = current - started_at
    return f"[{elapsed:6.1f}s]"


def run_main_agent_streaming(
    user_input: str,
    status_container,
    messages: list[tuple[str, str]] | None = None,
) -> dict:
    """Run main orchestrator with real-time progress via astream_events.

    Displays tool calls, subagent delegations, and todo updates
    inside ``status_container`` (a ``st.status`` widget).

    Returns the final agent result dict.
    """
    safe_user_input = coerce_message_content_to_text(user_input)
    safe_messages: list[tuple[str, str]]
    if messages is None:
        safe_messages = sanitize_history_pairs(
            [(msg["role"], msg["content"]) for msg in st.session_state.messages]
            + [("user", safe_user_input)]
        )
    else:
        safe_messages = sanitize_history_pairs(messages)

    _check_llama_server_available()

    # Debug: show which provider and model are being used
    from model_config import resolve_provider_and_model
    active_provider, active_model = resolve_provider_and_model(get_main_model())
    st.sidebar.caption(f"🔧 Using **{active_provider}** provider (model: {active_model or 'server-default'})")

    agent_obj, agent_mode = _select_main_agent_for_request(safe_user_input)
    session_id = st.session_state.thread_id

    # ── Tool label mapping for progress display ──
    _TOOL_LABELS = {
        "task": ("🤖", "Delegating to"),
        "write_todos": ("📋", "Planning tasks"),
        "session_search": ("🧠", "Searching past chats"),
        "record_learning": ("📝", "Saving reusable learning"),
        "execute_python_code": ("🐍", "Running Python code"),
        "install_package": ("📦", "Installing package"),
        "tavily_search": ("🌐", "Searching the web"),
        "generate_presentation": ("📊", "Generating slides"),
        "ingest_rag_documents": ("📚", "Ingesting documents"),
        "rag_retrieve": ("📚", "Retrieving documents"),
        "execute": ("⚡", "Running command"),
    }
    _MILESTONE_TOOLS = set(_TOOL_LABELS.keys())

    active_subagent = None
    progress_items: list[str] = [
        f"⚙️ **Main agent started ({agent_mode} path)**",
        f"📝 Request: `{_short_preview(safe_user_input, 100)}`",
        "🧭 Reviewing the request and deciding the next action",
    ]
    final_state = None
    last_status_update = 0.0
    started_at = time.time()
    model_wait_started_at: float | None = None
    model_heartbeat_active = False
    heartbeat_task: asyncio.Task | None = None
    with status_container:
        progress_placeholder = st.empty()

    def _flush_status(force: bool = False) -> None:
        nonlocal last_status_update
        now = time.time()
        if not force and now - last_status_update < 0.25:
            return
        try:
            progress_placeholder.markdown("\n\n".join(progress_items[-14:]))
        except StopException:
            # Streamlit uses StopException for rerun/stop control flow; treat it
            # as a normal shutdown signal for background progress updates.
            raise
        last_status_update = now

    def _upsert_progress_line(prefix: str, line: str) -> None:
        for idx in range(len(progress_items) - 1, -1, -1):
            existing = progress_items[idx]
            if existing.startswith(prefix) or f"] {prefix}" in existing:
                progress_items[idx] = line
                return
        progress_items.append(line)

    def _add_progress(line: str, *, now: float | None = None) -> None:
        progress_items.append(f"{_format_progress_timestamp(started_at, now)} {line}")

    def _upsert_timed_progress(prefix: str, line: str, *, now: float | None = None) -> None:
        _upsert_progress_line(prefix, f"{_format_progress_timestamp(started_at, now)} {line}")

    async def _model_heartbeat() -> None:
        nonlocal model_wait_started_at
        try:
            while model_heartbeat_active and model_wait_started_at is not None:
                now = time.time()
                elapsed = int(now - model_wait_started_at)
                _upsert_timed_progress(
                    "⏳ Model still thinking",
                    f"⏳ Model still thinking ({elapsed}s elapsed)",
                    now=now,
                )
                _flush_status(force=True)
                await asyncio.sleep(2)
        except (StopException, asyncio.CancelledError):
            return

    async def _stream_and_display():
        nonlocal active_subagent, final_state, model_wait_started_at, model_heartbeat_active, heartbeat_task
        event_stream = agent_obj.astream_events(
            {"messages": safe_messages},
            config=cast(RunnableConfig, build_agent_config(
                thread_id=session_id,
                user_id=st.session_state.user_id,
                recursion_limit=20,
                default_user_id="streamlit-user",
            )),
            version="v2",
        )
        try:
            async for event in event_stream:
                ev_type = event.get("event", "")
                ev_name = event.get("name", "")
                event_time = time.time()

                if ev_type == "on_chain_end" and ev_name == "LangGraph":
                    final_state = event.get("data", {}).get("output")
                    _add_progress("✅ **Main orchestration finished**", now=event_time)
                    _flush_status(force=True)

                elif ev_type == "on_chat_model_start":
                    model_name = event.get("metadata", {}).get("ls_model_name", "") or ev_name
                    model_wait_started_at = event_time
                    model_heartbeat_active = True
                    _add_progress(
                        f"🧠 Thinking with `{_short_preview(model_name or 'main model', 60)}`",
                        now=event_time,
                    )
                    _flush_status()
                    if heartbeat_task is None or heartbeat_task.done():
                        heartbeat_task = asyncio.create_task(_model_heartbeat())

                elif ev_type == "on_chat_model_stream":
                    if model_wait_started_at is not None:
                        elapsed = int(event_time - model_wait_started_at)
                        _upsert_timed_progress(
                            "⏳ Model still thinking",
                            f"⏳ Generating response ({elapsed}s elapsed)",
                            now=event_time,
                        )
                        _flush_status()

                elif ev_type == "on_chat_model_end":
                    model_heartbeat_active = False
                    model_wait_started_at = None
                    _add_progress("✍️ Synthesizing the answer for the chat", now=event_time)
                    _flush_status()

                elif ev_type == "on_chain_start" and ev_name != "LangGraph":
                    _add_progress(f"🔄 Entering `{_short_preview(ev_name, 60)}`", now=event_time)
                    _flush_status()

                if ev_type == "on_tool_start":
                    model_heartbeat_active = False
                    model_wait_started_at = None
                    tool_input = event.get("data", {}).get("input", {})
                    if ev_name == "task":
                        subagent_name = (
                            tool_input.get("subagent_type", "subagent")
                            if isinstance(tool_input, dict)
                            else "subagent"
                        )
                        active_subagent = subagent_name
                        desc_text = (
                            tool_input.get("description", "")[:80]
                            if isinstance(tool_input, dict)
                            else ""
                        )
                        line = f"🤖 **Delegating to {subagent_name}**"
                        if desc_text:
                            line += f" — _{desc_text}_"
                        _add_progress(line, now=event_time)
                    elif ev_name == "write_todos":
                        _add_progress("📋 **Updating task plan...**", now=event_time)
                    elif ev_name in _MILESTONE_TOOLS:
                        emoji, label = _TOOL_LABELS[ev_name]
                        ctx = f" (inside {active_subagent})" if active_subagent else ""
                        _add_progress(f"{emoji} {label}{ctx}...", now=event_time)
                    else:
                        continue
                    _flush_status()

                elif ev_type == "on_tool_end":
                    if ev_name == "task":
                        output = event.get("data", {}).get("output", "")
                        _add_progress(f"✅ **{active_subagent or 'Subagent'}** finished", now=event_time)
                        output_preview = _short_preview(output, 180)
                        if output_preview:
                            _add_progress(f"↳ {output_preview}", now=event_time)
                        active_subagent = None
                    elif ev_name == "write_todos":
                        output = event.get("data", {}).get("output", "")
                        if isinstance(output, str) and output.strip():
                            preview = output.strip()[:200]
                            _add_progress(f"✅ Tasks updated\n```\n{preview}\n```", now=event_time)
                        else:
                            _add_progress("✅ Tasks updated", now=event_time)
                    elif ev_name in _MILESTONE_TOOLS:
                        output = event.get("data", {}).get("output", "")
                        output_preview = _short_preview(output, 180)
                        if output_preview:
                            _add_progress(f"✅ {ev_name} finished — {output_preview}", now=event_time)
                        else:
                            _add_progress(f"✅ {ev_name} finished", now=event_time)
                    else:
                        continue
                    _flush_status(force=True)
        except (StopException, asyncio.CancelledError, GeneratorExit):
            return
        finally:
            model_heartbeat_active = False
            model_wait_started_at = None
            if heartbeat_task is not None:
                heartbeat_task.cancel()
                with suppress(asyncio.CancelledError, StopException, RuntimeError):
                    await heartbeat_task
                heartbeat_task = None
            aclose = getattr(event_stream, "aclose", None)
            if callable(aclose):
                with suppress(asyncio.CancelledError, RuntimeError):
                    await aclose()

    _flush_status(force=True)

    # Run streaming in the event loop - this handles async execution properly
    # without blocking Streamlit's UI responsiveness thanks to astream_events progress updates
    try:
        _run_async(_stream_and_display())
    except Exception as exc:
        if _is_llama_server_provider_active() and "connection attempts failed" in str(exc).lower():
            raise RuntimeError(
                "llama-server request failed because the server is unreachable. "
                "Start run-llama-server.bat and retry."
            ) from exc
        if _should_fallback_to_raw_llm(exc):
            return _invoke_raw_main_llm(safe_user_input)
        message_preview = [
            f"{idx + 1}. {role}: {_short_preview(content, 120)}"
            for idx, (role, content) in enumerate(safe_messages[-8:])
        ]
        raise RuntimeError(
            "Main agent streaming failed. "
            f"Original error: {type(exc).__name__}: {exc}\n"
            "Recent sanitized messages:\n"
            + "\n".join(message_preview)
        ) from exc

    # If astream_events didn't capture the final state via on_chain_end,
    # it means the agent finished but didn't emit the completion event.
    # In that case, just return what we have (the stream should have captured the output).
    if final_state is None:
        st.warning("⚠️ Agent finished but response capture incomplete. Attempting recovery...")
        # Minimal fallback: build a response dict from what we captured
        final_state = {
            "messages": [],
            "output": "",
            "status": "streaming_fallback"
        }

    return final_state


# ── Model Utilities ───────────────────────────────────────────────────────────

@st.cache_data(ttl=30)
def get_llama_cpp_models(models_dir: str | None = None) -> dict[str, str]:
    """Discover GGUF models for llama.cpp and return {display_label: model_id}."""
    base_dir = (models_dir or os.getenv("LLAMA_CPP_MODELS_DIR", "G:/llama_cpp/models")).strip()
    base_path = Path(base_dir)
    models: dict[str, str] = {}

    if base_path.exists() and base_path.is_dir():
        for gguf_path in sorted(base_path.rglob("*.gguf")):
            rel = gguf_path.relative_to(base_path).as_posix()
            models[f"📦 {rel}"] = f"llama_cpp:{gguf_path.as_posix()}"

    # Always include explicitly configured path even if outside models dir.
    configured_model_path = os.getenv("LLAMA_CPP_MODEL_PATH", "").strip()
    if configured_model_path:
        configured_path = Path(configured_model_path)
        configured_id = f"llama_cpp:{configured_path.as_posix()}"
        if configured_id not in models.values():
            models[f"📦 {configured_path.name}"] = configured_id

    return models


def get_huggingface_models() -> dict[str, str]:
    """Return configured Hugging Face repos as {display_label: model_id}."""
    models: dict[str, str] = {}

    configured_repo = os.getenv("HUGGINGFACE_REPO_ID", "").strip()
    if configured_repo:
        models[f"🤗 {configured_repo}"] = f"huggingface:{configured_repo}"

    # Comma/newline/semicolon separated repo IDs.
    configured_list = os.getenv("HUGGINGFACE_MODELS", "").strip()
    if configured_list:
        repos = [
            item.strip()
            for item in re.split(r"[,;\n]", configured_list)
            if item.strip()
        ]
        for repo in repos:
            models.setdefault(f"🤗 {repo}", f"huggingface:{repo}")

    return models


def get_llama_server_models() -> dict[str, str]:
    """Return llama-server model ids as {display_label: model_id}."""
    configured = os.getenv("LLAMA_SERVER_MODEL", "local").strip() or "local"
    return {f"🖧 {configured}": configured}


def _build_model_options(provider: str) -> dict[str, str]:
    """Return ordered {display_label: model_id} for the selected provider only."""
    if provider == "ollama":
        models: dict[str, str] = {}
        for name, mid in sorted(get_ollama_models().items()):
            models[f"🔵 {name}"] = mid
        return models

    if provider == "llama_server":
        return get_llama_server_models()

    if provider == "llama_cpp":
        return get_llama_cpp_models()

    if provider == "huggingface":
        return get_huggingface_models()

    return {}


def _build_all_model_options() -> dict[str, str]:
    """Return ordered {display_label: model_id} across all supported providers."""
    models: dict[str, str] = {}

    for display, mid in get_ollama_models().items():
        models[f"🔵 Ollama · {display}"] = mid

    for display, mid in get_llama_server_models().items():
        models[f"🖧 llama-server · {display}"] = mid

    for display, mid in get_llama_cpp_models().items():
        models[f"📦 llama.cpp · {display}"] = mid

    for display, mid in get_huggingface_models().items():
        models[f"🤗 Hugging Face · {display}"] = mid

    return dict(sorted(models.items(), key=lambda item: item[0].lower()))


def _reload_model_modules() -> None:
    """Drop cached agent modules after model changes."""
    reload_prefixes = (
        "agent",
        "agent_registry",
        "data_scientist_agent",
        "ragsub_agent",
        "presentation_agent",
        "specialist_router",
    )
    for module_name in list(sys.modules):
        if module_name in reload_prefixes or module_name.startswith(tuple(f"{p}." for p in reload_prefixes)):
            sys.modules.pop(module_name, None)
    get_cached_agent_bundle.clear()
    get_cached_data_scientist_bundle.clear()
    get_cached_ragsub_bundle.clear()
    get_cached_specialist_router_bundle.clear()
    get_cached_presenter_bundle.clear()


def _render_provider_selector() -> str:
    """Render a provider selectbox, persist to .env."""
    providers = {
        "🖧 llama-server (OpenAI-compatible)": "llama_server",
        "🔵 Ollama (local)": "ollama",
        "📦 llama.cpp (GGUF)": "llama_cpp",
        "🤗 Hugging Face": "huggingface",
    }
    
    current_provider = st.session_state.get("llm_provider") or get_default_llm_provider()
    current_display = next(
        (k for k, v in providers.items() if v == current_provider),
        "🔵 Ollama (local)"
    )
    
    selected_display = st.selectbox(
        "LLM Provider",
        list(providers.keys()),
        index=list(providers.keys()).index(current_display),
        key="provider_selector",
        help="Switch between llama-server, Ollama, llama.cpp, or Hugging Face providers",
    )
    selected = providers[selected_display]
    
    if selected != current_provider:
        set_env_value("DEEPAGENT_LLM_PROVIDER", selected)
        st.session_state["llm_provider"] = selected

        # Keep provider + model linked by switching to a model from this provider.
        provider_models = _build_model_options(selected)
        if provider_models:
            first_model = next(iter(provider_models.values()))
            if selected == "llama_cpp":
                _sync_llama_cpp_model_env(first_model)
            else:
                set_env_value("DEEPAGENT_MODEL", first_model)
                st.session_state["current_model"] = first_model

        reload_env_file()
        _reload_model_modules()
        st.rerun()

    if selected == "llama_server":
        ok, status_message = _probe_llama_server()
        if ok:
            st.caption(f"🟢 {status_message}")
        else:
            st.warning(status_message)
            if st.button("Start llama-server", key="start_llama_server_btn", use_container_width=True):
                try:
                    _launch_llama_server()
                except Exception as exc:
                    st.error(f"Failed to launch llama-server: {exc}")
                else:
                    st.success("Launching llama-server in a new window. Wait a few seconds, then retry.")
    
    return selected


def _render_local_context_preset_selector(active_provider: str) -> None:
    """Render a simple context preset selector for local llama providers."""
    if active_provider not in {"llama_server", "llama_cpp"}:
        return

    preset_map = {
        "8k (safe)": "8192",
        "16k (balanced)": "16384",
        "32k (tight)": "32768",
    }
    current_value = get_env_value(
        "LLAMA_SERVER_N_CTX" if active_provider == "llama_server" else "LLAMA_CPP_N_CTX",
        get_env_value("LLAMA_CPP_N_CTX", "8192"),
    ) or "8192"
    current_label = next((label for label, value in preset_map.items() if value == current_value), "Custom")
    option_labels = list(preset_map.keys())
    if current_label == "Custom":
        option_labels = ["Custom", *option_labels]

    selected_label = st.selectbox(
        "Local context preset",
        option_labels,
        index=option_labels.index(current_label),
        key=f"local_context_preset_{active_provider}",
        help="Higher context allows larger prompts and RAG/tool usage, but uses more VRAM/system RAM and can be slower.",
    )

    if selected_label != "Custom":
        selected_value = preset_map[selected_label]
        changed = False
        if get_env_value("LLAMA_CPP_N_CTX", "8192") != selected_value:
            set_env_value("LLAMA_CPP_N_CTX", selected_value)
            changed = True
        if get_env_value("LLAMA_SERVER_N_CTX", "") != selected_value:
            set_env_value("LLAMA_SERVER_N_CTX", selected_value)
            changed = True
        if changed:
            reload_env_file()
            _reload_model_modules()
            st.rerun()

    active_ctx = get_env_value(
        "LLAMA_SERVER_N_CTX" if active_provider == "llama_server" else "LLAMA_CPP_N_CTX",
        "8192",
    )
    st.caption(f"Active context window: `{active_ctx}` tokens")
    if active_provider == "llama_server":
        st.caption("Restart llama-server after changing the preset so the new context window takes effect.")


def _normalize_llama_cpp_model_path(model_path: str) -> str:
    """Normalize a llama.cpp GGUF path for storage in env/session state."""
    candidate = Path(model_path.strip()).expanduser()
    if not candidate.is_absolute():
        candidate = candidate.resolve(strict=False)
    return candidate.as_posix()


def _get_llama_cpp_models_root() -> Path:
    """Return the configured llama.cpp models root directory."""
    base_dir = os.getenv("LLAMA_CPP_MODELS_DIR", "G:/llama_cpp/models").strip()
    return Path(base_dir).expanduser()


def _list_llama_cpp_browser_dirs(root: Path) -> list[Path]:
    """Return browsable llama.cpp directories under the configured root."""
    dirs = [root]
    if root.exists() and root.is_dir():
        for child in sorted((p for p in root.rglob("*") if p.is_dir()), key=lambda p: p.as_posix().lower()):
            dirs.append(child)
    return dirs


def _list_llama_cpp_browser_files(folder: Path) -> list[Path]:
    """Return GGUF files in a selected llama.cpp folder."""
    if not folder.exists() or not folder.is_dir():
        return []
    return sorted(folder.glob("*.gguf"), key=lambda p: p.name.lower())


def _llama_cpp_folder_display_label(folder: Path, root_dir: Path) -> str:
    """Return the human-friendly label used for a llama.cpp browser folder."""
    if folder != root_dir and folder.is_relative_to(root_dir):
        display_path = folder.relative_to(root_dir).as_posix()
    else:
        display_path = folder.as_posix()
    return f"📁 {display_path}"


def _sync_llama_cpp_model_env(model_id: str) -> str:
    """Persist a llama.cpp model selection to both model env vars."""
    provider, resolved_model = resolve_provider_and_model(model_id)
    if provider != "llama_cpp" or not resolved_model:
        return resolved_model

    normalized_path = _normalize_llama_cpp_model_path(resolved_model)
    set_env_value("LLAMA_CPP_MODEL_PATH", normalized_path)
    set_env_value("DEEPAGENT_MODEL", f"llama_cpp:{normalized_path}")
    st.session_state["current_model"] = f"llama_cpp:{normalized_path}"
    return normalized_path


def _render_llama_cpp_model_path_selector() -> None:
    """Render a llama.cpp GGUF path picker backed by LLAMA_CPP_MODEL_PATH."""
    current_model = st.session_state.get("current_model") or get_main_model()
    current_provider, current_path = resolve_provider_and_model(current_model)
    if current_provider != "llama_cpp":
        current_path = get_env_value("LLAMA_CPP_MODEL_PATH", "")

    current_path = _normalize_llama_cpp_model_path(current_path) if current_path else ""
    root_dir = _get_llama_cpp_models_root()
    browser_dirs = _list_llama_cpp_browser_dirs(root_dir)
    current_dir = Path(st.session_state.get("llama_cpp_browser_dir", "")).expanduser() if st.session_state.get("llama_cpp_browser_dir") else (Path(current_path).parent if current_path else root_dir)
    if not current_dir.exists() or not current_dir.is_dir():
        current_dir = root_dir

    if current_dir not in browser_dirs:
        browser_dirs = [current_dir, *browser_dirs]

    folder_options = {_llama_cpp_folder_display_label(folder, root_dir): folder.as_posix() for folder in browser_dirs}
    current_folder_label = next(
        (label for label, folder in folder_options.items() if Path(folder) == current_dir),
        next(iter(folder_options)),
    )

    def _on_folder_change() -> None:
        selected_label = st.session_state.get("llama_cpp_browser_folder_selector", current_folder_label)
        selected_folder_path = folder_options.get(selected_label, current_dir.as_posix())
        st.session_state["llama_cpp_browser_dir"] = selected_folder_path
        st.rerun()

    def _go_to_parent_folder() -> None:
        parent_dir = current_dir.parent
        st.session_state["llama_cpp_browser_dir"] = parent_dir.as_posix()
        st.session_state["llama_cpp_browser_folder_selector"] = _llama_cpp_folder_display_label(parent_dir, root_dir)
        st.rerun()

    def _on_model_change() -> None:
        selected_display = st.session_state.get("llama_cpp_model_path_selector", "Choose a GGUF model...")
        selected_path = file_options.get(selected_display, "")
        if selected_path:
            st.session_state["llama_cpp_browser_dir"] = Path(selected_path).parent.as_posix()
            _sync_llama_cpp_model_env(f"llama_cpp:{selected_path}")
            reload_env_file()
            _reload_model_modules()
            st.rerun()

    def _apply_custom_path() -> None:
        custom_path_value = st.session_state.get("llama_cpp_custom_model_path", "").strip()
        if not custom_path_value:
            st.warning("Enter a llama.cpp GGUF path first.")
            return

        normalized_path = _normalize_llama_cpp_model_path(custom_path_value)
        if not Path(normalized_path).exists():
            st.error(f"llama.cpp model not found: {normalized_path}")
            return

        st.session_state["llama_cpp_model_path_selector"] = f"🟡 {Path(normalized_path).name}"
        st.session_state["llama_cpp_browser_dir"] = Path(normalized_path).parent.as_posix()
        _sync_llama_cpp_model_env(f"llama_cpp:{normalized_path}")
        reload_env_file()
        _reload_model_modules()
        st.rerun()

    selected_folder_label = st.selectbox(
        "Browse folder",
        list(folder_options.keys()),
        index=list(folder_options.keys()).index(current_folder_label),
        key="llama_cpp_browser_folder_selector",
        help="Pick the folder that contains your GGUF files.",
        on_change=_on_folder_change,
    )
    selected_folder = Path(folder_options[selected_folder_label])

    st.caption(f"Browsing: `{selected_folder.as_posix()}`")
    if selected_folder != root_dir and root_dir in selected_folder.parents:
        st.button(
            "⬆ Parent folder",
            key="llama_cpp_browser_parent_btn",
            use_container_width=True,
            on_click=_go_to_parent_folder,
        )

    gguf_files = _list_llama_cpp_browser_files(selected_folder)
    file_options = {"Choose a GGUF model...": ""}
    for file_path in gguf_files:
        file_options[f"📦 {file_path.name}"] = file_path.as_posix()

    if len(file_options) == 1:
        st.caption("No GGUF files found. Put models in the configured llama.cpp models folder or enter a custom path below.")
    else:
        current_display = next(
            (label for label, path in file_options.items() if path == current_path),
            "Choose a GGUF model...",
        )
        selected_display = st.selectbox(
            "llama.cpp model path",
            list(file_options.keys()),
            index=list(file_options.keys()).index(current_display),
            key="llama_cpp_model_path_selector",
            help="Choose a local GGUF file to load with llama.cpp. The selected path is written to LLAMA_CPP_MODEL_PATH and DEEPAGENT_MODEL.",
            on_change=_on_model_change,
        )
        selected_path = file_options[selected_display]
        if selected_path:
            st.caption(f"Active llama.cpp path: `{selected_path}`")

    custom_path = st.text_input(
        "Custom GGUF path",
        value=current_path,
        key="llama_cpp_custom_model_path",
        help="Paste an absolute path to a .gguf file that is not in the discovered models folder.",
    )
    st.button(
        "Use custom path",
        key="apply_llama_cpp_custom_model_path",
        use_container_width=True,
        on_click=_apply_custom_path,
    )


def _render_model_selector(
    label: str,
    session_key: str,
    env_key: str,
    widget_key: str,
    allow_inherit: bool = False,
    provider_override: str | None = None,
) -> str:
    """Render an inline model selectbox, persist to .env, return selected model id."""
    provider = provider_override or st.session_state.get("llm_provider") or get_default_llm_provider()
    if provider == "all":
        all_models = _build_all_model_options()
    else:
        all_models = _build_model_options(provider)
    if allow_inherit:
        all_models["↔️ Use Main Model"] = ""

    if not all_models:
        provider_hint = {
            "llama_server": "Start run-llama-server.bat and set LLAMA_SERVER_BASE_URL/LLAMA_SERVER_MODEL if needed",
            "ollama": "Start Ollama or pull a model (e.g., ollama pull gemma4:26b)",
            "llama_cpp": "Put .gguf files in G:/llama_cpp/models or set LLAMA_CPP_MODEL_PATH",
            "huggingface": "Set HUGGINGFACE_REPO_ID (e.g., meta-llama/Meta-Llama-3.1-8B-Instruct)",
            "all": "Configure at least one model for Ollama, llama-server, llama.cpp, or Hugging Face",
        }.get(provider, "Configure a model for this provider")
        st.caption(f"No models available for `{provider}`. {provider_hint}.")
        return st.session_state.get(session_key, "")

    current = st.session_state.get(session_key, "")
    current_provider, _ = resolve_provider_and_model(current)
    current_display: str | None = None
    for display, mid in all_models.items():
        if mid == current:
            current_display = display
            break
    if current_display is None:
        if current:
            if provider == "all" or current_provider == provider:
                short = current.split(":", 1)[1] if ":" in current else current
                current_display = f"🟡 {short}"
                all_models = {current_display: current, **all_models}
            else:
                # Ignore stale model values from other providers.
                current = ""
                current_display = next(iter(all_models))
        elif allow_inherit:
            # Keep "Use Main Model" available, but do not auto-select it when
            # no explicit context model has been chosen yet.
            current_display = next(iter(all_models))
        else:
            current_display = next(iter(all_models))

    selected_display = st.selectbox(
        label,
        list(all_models.keys()),
        index=list(all_models.keys()).index(current_display),
        key=widget_key,
    )
    selected = all_models[selected_display]

    if selected != st.session_state.get(session_key, ""):
        set_env_value(env_key, selected)
        st.session_state[session_key] = selected
        if env_key == "DEEPAGENT_MODEL" and resolve_provider_and_model(selected)[0] == "llama_cpp":
            _sync_llama_cpp_model_env(selected)
        reload_env_file()
        _reload_model_modules()
        st.rerun()

    return selected


def _render_ollama_tag_selector(
    label: str,
    session_key: str,
    env_key: str,
    widget_key: str,
) -> str:
    """Render an Ollama-only selector that stores the raw model tag in .env."""
    ollama_models = get_ollama_models()
    all_models = {f"🔵 {name}": name for name in sorted(ollama_models)}

    current = st.session_state.get(session_key, "") or os.getenv(env_key, "").strip()
    if current.startswith("ollama:"):
        current = current[len("ollama:"):]
    if current and st.session_state.get(session_key, "") != current:
        st.session_state[session_key] = current
    current_display: str | None = None
    for display, model_tag in all_models.items():
        if model_tag == current:
            current_display = display
            break
    if current_display is None:
        if current:
            current_display = f"🟡 {current}"
            all_models = {current_display: current, **all_models}
        else:
            current_display = next(iter(all_models), "🟡 no llama.cpp models found")
            all_models.setdefault(current_display, current_display.replace("🟡 ", ""))

    selected_display = st.selectbox(
        label,
        list(all_models.keys()),
        index=list(all_models.keys()).index(current_display),
        key=widget_key,
    )
    selected = all_models[selected_display]

    if selected != current:
        set_env_value(env_key, selected)
        st.session_state[session_key] = selected
        reload_env_file()
        _reload_model_modules()
        st.rerun()

    return selected


def _render_agent_model_overrides() -> None:
    """Render per-agent model override controls backed by the shared config helper."""
    st.markdown("**Agent Models**")
    agent_order = [
        "websearch",
        "writer",
        "coder",
        "planner",
        "reviewer",
        "presenter",
        "data_scientist",
        "ragsub",
    ]
    left_col, right_col = st.columns(2)
    for idx, agent_name in enumerate(agent_order):
        col = left_col if idx % 2 == 0 else right_col
        with col:
            _render_model_selector(
                f"{AGENT_MODEL_LABELS[agent_name]} Model",
                session_key=f"model_override_{agent_name}",
                env_key=AGENT_MODEL_ENV_KEYS[agent_name],
                widget_key=f"model_override_selector_{agent_name}",
                allow_inherit=True,
            )


# ── Chat Utilities ────────────────────────────────────────────────────────────

def _render_feedback_buttons(msg_index: int, user_prompt: str, agent_response: str) -> None:
    """Render 👍 ❤️ 👎 feedback buttons for an assistant message.

    Feedback is persisted to agent_learnings.md so the agent learns from it.
    Once rated, buttons are replaced with a confirmation label.
    """
    fb_key = f"feedback_{msg_index}"
    bad_reason_key = f"bad_reason_{msg_index}"
    bad_pending_key = f"bad_pending_{msg_index}"

    if not st.session_state.get("self_learning_enabled", True):
        st.caption("Self-learning is off")
        return

    # Already rated — show confirmation only
    if st.session_state.get(fb_key):
        rating = st.session_state[fb_key]
        label = {"very_good": "❤️ Saved", "good": "👍 Saved", "bad": "👎 Saved"}.get(rating, "Saved")
        st.caption(label)
        return

    # Bad feedback pending — show optional reason box
    if st.session_state.get(bad_pending_key):
        reason = st.text_input(
            "What was wrong? (optional)",
            key=f"bad_reason_input_{msg_index}",
            placeholder="e.g. wrong language, too long, off-topic...",
            label_visibility="collapsed",
        )
        c1, c2 = st.columns([1, 4])
        with c1:
            if st.button("Submit", key=f"bad_submit_{msg_index}", width="stretch"):
                conversation_memory.record_feedback(
                    rating="bad",
                    user_prompt=user_prompt,
                    agent_response=agent_response,
                    reason=reason,
                )
                st.session_state[fb_key] = "bad"
                st.session_state[bad_pending_key] = False
                st.rerun()
        with c2:
            if st.button("Skip", key=f"bad_skip_{msg_index}", width="stretch"):
                conversation_memory.record_feedback(
                    rating="bad",
                    user_prompt=user_prompt,
                    agent_response=agent_response,
                )
                st.session_state[fb_key] = "bad"
                st.session_state[bad_pending_key] = False
                st.rerun()
        return

    # Default — show three rating buttons
    fb_col1, fb_col2, fb_col3, _ = st.columns([0.6, 0.6, 0.6, 6])
    with fb_col1:
        if st.button("👍", key=f"fb_good_{msg_index}", help="Good response"):
            conversation_memory.record_feedback(
                rating="good",
                user_prompt=user_prompt,
                agent_response=agent_response,
            )
            st.session_state[fb_key] = "good"
            st.rerun()
    with fb_col2:
        if st.button("❤️", key=f"fb_vgood_{msg_index}", help="Very good response"):
            conversation_memory.record_feedback(
                rating="very_good",
                user_prompt=user_prompt,
                agent_response=agent_response,
            )
            st.session_state[fb_key] = "very_good"
            st.rerun()
    with fb_col3:
        if st.button("👎", key=f"fb_bad_{msg_index}", help="Bad response"):
            st.session_state[bad_pending_key] = True
            st.rerun()


def _render_chat_history() -> None:
    """Render all messages stored in session state."""
    messages = st.session_state.messages
    for idx, message in enumerate(messages):
        with st.chat_message(message["role"]):
            render_message_with_react(message["content"])
            if message["role"] == "user":
                _render_chat_attachment_block(message)
            if message["role"] == "assistant" and message.get("plots"):
                for plot_path in message["plots"]:
                    if Path(plot_path).exists():
                        st.image(plot_path, width="stretch")
            if message["role"] == "assistant" and message.get("rag_assets"):
                for asset_path in message["rag_assets"]:
                    if Path(asset_path).exists():
                        st.image(asset_path, width="stretch")
            if message["role"] == "assistant":
                # Find the preceding user message for context
                prior_user = next(
                    (messages[i]["content"] for i in range(idx - 1, -1, -1) if messages[i]["role"] == "user"),
                    "",
                )
                _render_feedback_buttons(
                    msg_index=idx,
                    user_prompt=prior_user,
                    agent_response=message["content"],
                )
                _render_qa_export_buttons(messages=messages, assistant_index=idx)


def _attachment_export_lines(message: dict[str, Any]) -> list[str]:
    """Return export-friendly attachment summary lines for a message."""
    attachment_paths = [str(p) for p in message.get("attachment_paths", []) if str(p).strip()]
    if not attachment_paths:
        return []

    names = ", ".join(Path(path).name for path in attachment_paths)
    lines = [f"Attachments: {names}"]
    attachment_mode = str(message.get("attachment_mode", "") or "").strip()
    attachment_route = str(message.get("attachment_route", "") or "").strip()
    if attachment_mode:
        lines.append(f"Attachment mode: {attachment_mode}")
    if attachment_route:
        lines.append(f"Attachment route: {attachment_route}")
    return lines


def _export_message_content(message: dict[str, Any]) -> str:
    """Return message content augmented with attachment metadata for exports."""
    content = str(message.get("content", "")).strip()
    extra_lines = _attachment_export_lines(message)
    if not extra_lines:
        return content
    extra_block = "\n".join(extra_lines)
    return f"{content}\n\n{extra_block}" if content else extra_block


def _get_exportable_qa_pair(messages: list[dict[str, Any]], assistant_index: int) -> list[dict[str, Any]]:
    """Return the preceding user message plus the selected assistant reply."""
    if assistant_index < 0 or assistant_index >= len(messages):
        return []
    assistant_message = messages[assistant_index]
    if str(assistant_message.get("role", "")).strip() != "assistant":
        return []

    assistant_content = _export_message_content(assistant_message)
    if not assistant_content:
        return []

    prior_user_message: dict[str, Any] | None = None
    for i in range(assistant_index - 1, -1, -1):
        if str(messages[i].get("role", "")).strip() == "user":
            prior_user_message = messages[i]
            break

    pair: list[dict[str, Any]] = []
    if prior_user_message is not None:
        prior_user_content = _export_message_content(prior_user_message)
        if prior_user_content:
            pair.append({"role": "user", "content": prior_user_content, **{
                k: prior_user_message[k] for k in ("attachment_paths", "attachment_mode", "attachment_route")
                if k in prior_user_message
            }})
    pair.append({"role": "assistant", "content": assistant_content})
    return pair


def _build_chat_export_basename(prefix: str = "chat_output", sequence_number: int | None = None) -> str:
    """Return a stable export filename stem for a chat export."""
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    if sequence_number is None:
        return f"{prefix}_{timestamp}"
    return f"{prefix}_qna_{sequence_number}_{timestamp}"


def _build_chat_export_title(sequence_number: int | None = None) -> str:
    """Return a human-readable export title."""
    label = "Q&A Export" if sequence_number is not None else "Chat Output Export"
    if sequence_number is not None:
        label = f"Q&A Export {sequence_number}"
    return f"{label} - {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"


def _normalize_markdown_text(text: str) -> str:
    """Strip markdown markers that are not rendered natively in exports."""
    cleaned = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    cleaned = cleaned.replace("**", "").replace("__", "")
    cleaned = cleaned.replace("`", "")
    return cleaned


def _parse_inline_markdown(text: str) -> list[tuple[str, bool]]:
    """Parse simple inline markdown, preserving bold segments."""
    normalized = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    parts: list[tuple[str, bool]] = []
    cursor = 0
    for match in re.finditer(r"\*\*(.+?)\*\*", normalized):
        if match.start() > cursor:
            plain = normalized[cursor:match.start()].replace("`", "")
            if plain:
                parts.append((plain, False))
        bold_text = match.group(1).replace("`", "")
        if bold_text:
            parts.append((bold_text, True))
        cursor = match.end()
    if cursor < len(normalized):
        tail = normalized[cursor:].replace("`", "")
        if tail:
            parts.append((tail, False))
    return parts or [("", False)]


def _markdown_to_blocks(text: str) -> list[dict[str, Any]]:
    """Convert markdown-ish text to simple render blocks."""
    blocks: list[dict[str, Any]] = []
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading_match:
            blocks.append(
                {
                    "type": "heading",
                    "level": len(heading_match.group(1)),
                    "segments": _parse_inline_markdown(heading_match.group(2).strip()),
                }
            )
            continue

        bullet_match = re.match(r"^[-*+]\s+(.*)$", stripped)
        if bullet_match:
            blocks.append(
                {
                    "type": "bullet",
                    "level": 0,
                    "segments": _parse_inline_markdown(bullet_match.group(1).strip()),
                }
            )
            continue

        blocks.append(
            {
                "type": "paragraph",
                "level": 0,
                "segments": _parse_inline_markdown(stripped),
            }
        )
    return blocks or [{"type": "paragraph", "level": 0, "segments": [("", False)]}]


def _plain_text_from_segments(segments: list[tuple[str, bool]]) -> str:
    """Flatten inline segments to plain text."""
    return "".join(text for text, _ in segments).strip()


def _add_docx_runs(paragraph, segments: list[tuple[str, bool]]) -> None:
    """Append inline runs with bold styling to a DOCX paragraph."""
    for text, is_bold in segments:
        run = paragraph.add_run(text)
        run.bold = bool(is_bold)


def _build_chat_docx_bytes(messages: list[dict[str, str]], sequence_number: int | None = None) -> bytes:
    """Build a DOCX transcript for a Q&A export."""
    from docx import Document

    doc = Document()
    doc.add_heading(_build_chat_export_title(sequence_number), level=0)
    for idx, message in enumerate(messages, start=1):
        doc.add_heading(f"{idx}. {message['role'].title()}", level=1)
        for block in _markdown_to_blocks(message["content"]):
            if block["type"] == "heading":
                paragraph = doc.add_paragraph()
                _add_docx_runs(paragraph, block["segments"])
                paragraph.style = f"Heading {min(block['level'] + 1, 4)}"
            elif block["type"] == "bullet":
                paragraph = doc.add_paragraph(style="List Bullet")
                _add_docx_runs(paragraph, block["segments"])
            else:
                paragraph = doc.add_paragraph()
                _add_docx_runs(paragraph, block["segments"])
    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _segments_to_pdf_markup(segments: list[tuple[str, bool]]) -> str:
    """Convert inline segments to reportlab Paragraph markup."""
    from xml.sax.saxutils import escape

    parts: list[str] = []
    for text, is_bold in segments:
        escaped = escape(text)
        parts.append(f"<b>{escaped}</b>" if is_bold else escaped)
    return "".join(parts) or " "


def _build_chat_pdf_bytes(messages: list[dict[str, str]], sequence_number: int | None = None) -> bytes:
    """Build a PDF transcript for a Q&A export."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, leftMargin=48, rightMargin=48, topMargin=48, bottomMargin=48)
    styles = getSampleStyleSheet()
    title_style = styles["Title"]
    section_style = styles["Heading2"]
    body_style = styles["BodyText"]
    heading_style = ParagraphStyle("ExportHeading", parent=styles["Heading3"], textColor=colors.black)
    bullet_style = ParagraphStyle("ExportBullet", parent=styles["BodyText"], leftIndent=18, bulletIndent=6)
    story = [Paragraph(_build_chat_export_title(sequence_number), title_style), Spacer(1, 12)]

    for idx, message in enumerate(messages, start=1):
        story.append(Paragraph(f"{idx}. {message['role'].title()}", section_style))
        story.append(Spacer(1, 6))
        for block in _markdown_to_blocks(message["content"]):
            if block["type"] == "heading":
                story.append(Paragraph(_segments_to_pdf_markup(block["segments"]), heading_style))
            elif block["type"] == "bullet":
                story.append(Paragraph(_segments_to_pdf_markup(block["segments"]), bullet_style, bulletText="•"))
            else:
                story.append(Paragraph(_segments_to_pdf_markup(block["segments"]), body_style))
            story.append(Spacer(1, 4))
        story.append(Spacer(1, 8))

    doc.build(story)
    return buffer.getvalue()


def _build_chat_pptx_bytes(messages: list[dict[str, str]], sequence_number: int | None = None) -> bytes:
    """Build a PPTX transcript for a Q&A export."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Q&A Export" if sequence_number is not None else "Chat Output Export"
    title_slide.placeholders[1].text = _build_chat_export_title(sequence_number)

    for idx, message in enumerate(messages, start=1):
        blocks = _markdown_to_blocks(message["content"])
        if not blocks:
            blocks = [{"type": "paragraph", "level": 0, "segments": [(" ", False)]}]

        for chunk_start in range(0, len(blocks), 6):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{idx}. {message['role'].title()}"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for block_idx, block in enumerate(blocks[chunk_start:chunk_start + 6]):
                paragraph = text_frame.paragraphs[0] if block_idx == 0 else text_frame.add_paragraph()
                paragraph.text = ""
                paragraph.level = 0
                run = paragraph.add_run()
                if block["type"] == "bullet":
                    run.text = f"• {_plain_text_from_segments(block['segments'])}"
                    run.font.size = Pt(18)
                elif block["type"] == "heading":
                    run.text = _plain_text_from_segments(block["segments"])
                    run.font.bold = True
                    run.font.size = Pt(22 if block["level"] <= 2 else 18)
                else:
                    paragraph.level = 0
                    for seg_idx, (text, is_bold) in enumerate(block["segments"]):
                        segment_run = run if seg_idx == 0 else paragraph.add_run()
                        segment_run.text = _normalize_markdown_text(text)
                        segment_run.font.bold = bool(is_bold)
                        segment_run.font.size = Pt(18)
                    continue
                run.font.bold = block["type"] == "heading"

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _render_qa_export_buttons(messages: list[dict[str, Any]], assistant_index: int) -> None:
    """Render PPTX/DOCX/PDF download buttons for one Q&A sequence."""
    export_messages = _get_exportable_qa_pair(messages, assistant_index)
    if not export_messages:
        return

    sequence_number = sum(
        1
        for i, message in enumerate(messages[: assistant_index + 1])
        if str(message.get("role", "")).strip() == "assistant"
    )
    base_name = _build_chat_export_basename(prefix="rag_chat_output", sequence_number=sequence_number)
    pptx_bytes = _build_chat_pptx_bytes(export_messages, sequence_number=sequence_number)
    docx_bytes = _build_chat_docx_bytes(export_messages, sequence_number=sequence_number)
    pdf_bytes = _build_chat_pdf_bytes(export_messages, sequence_number=sequence_number)

    st.caption("Download this Q&A")
    col_pptx, col_docx, col_pdf = st.columns(3)
    with col_pptx:
        st.download_button(
            label="📥 PPTX",
            data=pptx_bytes,
            file_name=f"{base_name}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"rag_chat_export_{assistant_index}_pptx",
            width="stretch",
        )
    with col_docx:
        st.download_button(
            label="📥 DOCX",
            data=docx_bytes,
            file_name=f"{base_name}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            key=f"rag_chat_export_{assistant_index}_docx",
            width="stretch",
        )
    with col_pdf:
        st.download_button(
            label="📥 PDF",
            data=pdf_bytes,
            file_name=f"{base_name}.pdf",
            mime="application/pdf",
            key=f"rag_chat_export_{assistant_index}_pdf",
            width="stretch",
        )


def _chat_execution_block(
    user_message: str,
    user_input_raw: str,
    force_agent_type: str | None = None,
    force_main_agent: bool = False,
    selected_rag_context: str = "",
    user_message_meta: dict[str, Any] | None = None,
    started_at: float | None = None,
) -> None:
    """Append user message, call appropriate agent, display response + artifacts.

    Args:
        user_message:     Augmented message sent to the LLM (may contain file tags, RAG context).
        user_input_raw:   Original typed input stored in history for clean display.
        force_agent_type: AgentType value to override the sidebar selection (e.g. RAG_SUB).
    """
    request_started_at = started_at if started_at is not None else time.time()
    session_id = st.session_state.thread_id
    model_short = (
        st.session_state.current_model.split(":", 1)[1]
        if ":" in st.session_state.current_model
        else st.session_state.current_model
    )

    _agent_display_opts = registry.get_all_agents()
    _agent_type, _auto_fast_route = _streamlit_chat_execution.resolve_chat_agent_type(
        user_message=user_message,
        selected_agent=st.session_state.selected_agent,
        agent_display_options=_agent_display_opts,
        rag_knowledge_only=st.session_state.rag_knowledge_only,
        force_write_todos=st.session_state.force_write_todos,
        force_agent_type=force_agent_type,
        force_main_agent=force_main_agent,
    )

    _non_streaming = (AgentType.DATA_SCIENTIST.value, AgentType.RAG_SUB.value, AgentType.PRESENTER.value)
    _use_streaming = _agent_type not in _non_streaming

    # Append raw message to history + display it
    _user_message_record: dict[str, Any] = {"role": "user", "content": user_input_raw}
    if user_message_meta:
        _user_message_record.update(user_message_meta)
    st.session_state.messages.append(_user_message_record)
    with st.chat_message("user"):
        st.write(user_input_raw)
        if user_message_meta:
            _render_chat_attachment_block(_user_message_record)

    # Fast local reply for pure social turns (e.g., "hi", "thanks") so
    # llama.cpp does not spin up full orchestration for trivial chat.
    _has_attachments = bool(user_message_meta and user_message_meta.get("attachments"))
    _has_selected_rag_context = bool((selected_rag_context or "").strip())
    if _streamlit_chat_execution.should_bypass_agent_for_trivial_turn(
        user_text=user_input_raw,
        has_attachments=_has_attachments,
        has_selected_rag_context=_has_selected_rag_context,
    ):
        instant_reply = _streamlit_chat_execution.get_instant_chat_response(user_input_raw)
        if instant_reply:
            with st.chat_message("assistant"):
                render_message_with_react(instant_reply)
                if st.session_state.voice_chat_enabled:
                    _render_speak_button(
                        instant_reply,
                        key=f"live-assistant-{len(st.session_state.messages)}",
                        auto_speak=st.session_state.voice_auto_speak,
                    )
                _persist_assistant_response(
                    session_id=session_id,
                    final_text=instant_reply,
                    plot_paths=[],
                    rag_asset_paths=[],
                )
                total_time = max(0.0, time.time() - request_started_at)
                st.session_state.last_execution_time = total_time
                _render_response_metrics(total_time, instant_reply, model_short, prefix="instant")
            return

    # Build message history for the agent.
    # Use all prior messages, then replace the just-appended raw user turn
    # with the augmented version (contains RAG context, presets, etc.).
    # This avoids sending consecutive HumanMessages which some LLMs reject.
    _history = [(m["role"], m["content"]) for m in st.session_state.messages[:-1]]
    _history.append(("user", user_message))
    _history = sanitize_history_pairs(_history)

    with st.chat_message("assistant"):
        try:
            artifact_tracker = _prepare_artifact_tracking(_agent_type)
            route_kwargs = _build_agent_route_kwargs(
                agent_type=_agent_type,
                auto_fast_route=_auto_fast_route,
                user_message=user_message,
                history=_history,
                session_id=session_id,
            )
            
            # Show immediate progress indicator to prevent UI freeze perception
            _progress_placeholder = st.empty()
            _progress_placeholder.info("⏳ Processing your request...", icon="ℹ️")
            agent_started_at = time.time()

            # ── Route to agent ──
            if _agent_type in {
                AgentType.DATA_SCIENTIST.value,
                AgentType.RAG_SUB.value,
                AgentType.PRESENTER.value,
            }:
                # Non-streaming agents (specialized tools)
                spinner_text = {
                    AgentType.DATA_SCIENTIST.value: "🔬 Analysing...",
                    AgentType.RAG_SUB.value: "📚 Retrieving...",
                    AgentType.PRESENTER.value: "📊 Building presentation...",
                }[_agent_type]
                with st.spinner(spinner_text):
                    result = _streamlit_chat_execution.invoke_agent_route(**route_kwargs)
            elif _agent_type in {
                AgentType.WEBSEARCH.value,
                AgentType.WRITER.value,
                AgentType.CODER.value,
                AgentType.PLANNER.value,
                AgentType.REVIEWER.value,
            }:
                # Specialist agents (usually delegate to main eventually)
                spinner_text = "Working with specialist..."
                if _auto_fast_route:
                    spinner_text = f"Fast-routing to {_agent_type.replace('_', ' ')}..."
                with st.spinner(spinner_text):
                    result = _streamlit_chat_execution.invoke_agent_route(**route_kwargs)
            else:
                status_label = (
                    "💬 Direct chat model"
                    if st.session_state.main_chat_use_direct_model
                    else "⚙️ Full Deep Agent working..."
                )
                status_expanded = not st.session_state.main_chat_use_direct_model
                
                with st.status(status_label, expanded=status_expanded) as _status:
                    if st.session_state.main_chat_use_direct_model:
                        result = _run_direct_main_chat(user_message)
                    else:
                        result = _streamlit_chat_execution.invoke_agent_route(
                            **_build_agent_route_kwargs(
                                agent_type=_agent_type,
                                auto_fast_route=_auto_fast_route,
                                user_message=user_message,
                                history=_history,
                                session_id=session_id,
                                status_container=_status,
                            )
                        )
                    _status.update(label="✅ Complete", state="complete", expanded=False)
            
            # Clear progress indicator
            _progress_placeholder.empty()

            agent_finished_at = time.time()
            pre_agent_time = max(0.0, agent_started_at - request_started_at)
            agent_time = max(0.0, agent_finished_at - agent_started_at)
            final_text = extract_agent_display_text(result)

            new_plot_paths, new_presentation_files = _collect_new_artifacts(artifact_tracker)
            rag_raw_chunks = _extract_rag_raw_chunks(
                agent_type=_agent_type,
                result=result,
                selected_rag_context=selected_rag_context,
            )

            if final_text:
                chunk_details = _extract_selected_chunk_details(rag_raw_chunks)
                render_message_with_react(final_text)
                if st.session_state.voice_chat_enabled:
                    _render_speak_button(
                        final_text,
                        key=f"live-assistant-{len(st.session_state.messages)}",
                        auto_speak=st.session_state.voice_auto_speak,
                    )
                rag_asset_paths = _render_rag_asset_previews(chunk_details) if chunk_details else []
                _render_selected_chunks_expander(chunk_details)
                _render_generated_artifacts(
                    session_id=session_id,
                    plot_paths=new_plot_paths,
                    presentation_files=new_presentation_files,
                    presentations_dir=artifact_tracker["presentations_dir"],
                )
                _persist_assistant_response(
                    session_id=session_id,
                    final_text=final_text,
                    plot_paths=new_plot_paths,
                    rag_asset_paths=rag_asset_paths,
                )
                total_time = max(0.0, time.time() - request_started_at)
                post_agent_time = max(0.0, total_time - pre_agent_time - agent_time)
                st.session_state.last_execution_time = total_time
                _render_response_metrics(
                    total_time,
                    final_text,
                    model_short,
                    timing_breakdown={
                        "pre": pre_agent_time,
                        "agent": agent_time,
                        "post": post_agent_time,
                    },
                )
            else:
                st.error("No response from the agent.")
                # Append placeholder so the next retry doesn't create
                # consecutive user messages (which many LLMs reject).
                st.session_state.messages.append(
                    {"role": "assistant", "content": "(No response from the agent.)"}
                )
                _save_current_session(force_index=True)
                with st.expander("Debug", expanded=False):
                    st.code(str(result)[:2000], language="text")

        except Exception as ex:
            exc_type = type(ex).__name__
            exc_msg = str(ex)
            exc_trace = traceback.format_exc()
            st.error(f"Error ({exc_type}): {exc_msg}")
            # Same guard: prevent orphan user message on exception.
            st.session_state.messages.append(
                {"role": "assistant", "content": f"(Error {exc_type}: {exc_msg})"}
            )
            _save_current_session(force_index=True)
            with st.expander("Debug (full traceback)", expanded=True):
                st.code(exc_trace, language="python")


def _render_main_chat_helpers() -> None:
    """Render main-mode voice and attachment helper UI above the chat input."""
    if st.session_state.voice_chat_enabled:
        _render_browser_voice_chat_button()
        with st.expander("🎙️ Voice Chat", expanded=False):
            st.caption(
                "Use the mic button inside the chat box for free browser-based dictation. "
                "Best support is typically in Chrome or Edge."
            )
            st.caption(
                "No API key is required for the browser mic path. "
                "Speak into the chat box, then send as normal."
            )

        


def _render_main_chat_history_or_empty_state() -> None:
    """Render chat history when available, otherwise the main-mode empty state."""
    if st.session_state.messages:
        st.divider()
        _render_chat_history()
        return

    st.info(
        "💡 **Get started** — type a message below, or click a previous chat in the sidebar to continue where you left off.",
        icon=None,
    )


def _get_main_chat_submission() -> tuple[str, list[Any]] | None:
    """Read the main chat input widget and normalize its value."""
    chat_input_value = st.chat_input(
        "Ask me to analyze data, research topics, write content, or convert to PPTX...",
        accept_file="multiple",
        file_type=CHAT_ATTACHMENT_FILE_TYPES,
    )
    if not chat_input_value:
        return None

    if isinstance(chat_input_value, str):
        user_input = chat_input_value
        uploaded_chat_files: list[Any] = []
    else:
        user_input = chat_input_value.text
        uploaded_chat_files = list(chat_input_value.files or [])

    if not user_input.strip() and not uploaded_chat_files:
        return None
    return user_input, uploaded_chat_files


def _augment_main_chat_user_message(
    *,
    user_input: str,
    chat_relpaths: list[str],
    attachment_mode: str,
) -> tuple[str, str, str, str | None, bool, str]:
    """Build the attachment-aware main chat message and routing metadata."""
    image_only_attachments = _chat_attachments_are_image_only(chat_relpaths)
    single_pdf_attachment = _chat_attachments_are_single_pdf(chat_relpaths)
    direct_image_context = ""
    direct_pdf_context = ""
    normalized_user_input = user_input

    if not normalized_user_input.strip() and chat_relpaths:
        normalized_user_input = (
            "Please index and summarize the attached files."
            if attachment_mode == "Add to RAG"
            else "Please analyze the attached files."
        )

    user_message = normalized_user_input.strip()
    display_user_input = _build_attachment_display_text(
        normalized_user_input.strip(),
        chat_relpaths,
        attachment_mode,
    )

    forced_agent_type, force_main_agent = _attachment_route_for_chat(
        relpaths=chat_relpaths,
        attachment_mode=attachment_mode,
        selected_agent=st.session_state.selected_agent,
    )
    attachment_route = _attachment_route_label(
        force_agent_type=forced_agent_type,
        force_main_agent=force_main_agent,
        selected_agent=st.session_state.selected_agent,
        attachment_mode=attachment_mode,
    )

    if chat_relpaths:
        if attachment_mode == "Ask in chat" and image_only_attachments:
            with st.spinner("Extracting text from image..."):
                direct_image_context = _build_direct_image_context(chat_relpaths)
        elif attachment_mode == "Ask in chat" and single_pdf_attachment:
            with st.spinner("Extracting text from PDF..."):
                direct_pdf_context = _build_direct_pdf_context(chat_relpaths)
        user_message = _build_attachment_augmented_message(
            user_message=user_message,
            relpaths=chat_relpaths,
            attachment_mode=attachment_mode,
            direct_image_context=direct_image_context,
            direct_pdf_context=direct_pdf_context,
        )

    return (
        user_message,
        display_user_input,
        attachment_route,
        forced_agent_type,
        force_main_agent,
        normalized_user_input,
    )


def _apply_main_chat_context_overrides(user_message: str, user_input: str) -> str:
    """Apply uploaded-file, RAG, workflow, and agent-preference wrappers to main chat input."""
    if st.session_state.uploaded_file_path:
        fp = Path(st.session_state.uploaded_file_path)
        rel = _safe_project_relative_path(fp)
        if fp.exists() and rel:
            user_message = (
                f"[UPLOADED FILE: {rel}]\n"
                f"[DATA CHECK: Use execute_python_code for exact dataset shape/statistics.]\n\n"
                f"{user_message}"
            )

    if st.session_state.rag_uploaded_files:
        valid = [
            _safe_project_relative_path(Path(p))
            for p in st.session_state.rag_uploaded_files
            if Path(p).exists() and _safe_project_relative_path(Path(p))
        ]
        valid_strs = [v for v in valid if v is not None]
        if valid_strs:
            user_message = (
                f"[RAG UPLOADED FILES: {', '.join(valid_strs)}]\n"
                f"[RAG WORKFLOW: Ingest files first, then retrieve and rerank.]\n\n"
                f"{user_message}"
            )

    _themes = st.session_state.rag_query_themes
    if _themes and _themes != "select_all":
        user_message = f"[RAG THEMES: {', '.join(_themes)}]\n{user_message}"

    if (
        st.session_state.rag_enable_main_chat
        and not st.session_state.rag_knowledge_only
        and "Data Scientist" not in st.session_state.selected_agent
        and "RAG SubAgent" not in st.session_state.selected_agent
        and "Presenter" not in st.session_state.selected_agent
    ):
        rag_context = get_main_chat_rag_context(user_input.strip())
        if (
            (not rag_context)
            or ("No usable chunks returned from retrieval" in rag_context)
            or ("⚠️ No usable chunks returned from retrieval" in rag_context)
        ):
            current_project_filter = st.session_state.get("rag_project_filter_selector", "")
            if current_project_filter and current_project_filter != "All Projects":
                fallback_context = ""
                try:
                    st.session_state.rag_project_filter_selector = "All Projects"
                    fallback_context = get_main_chat_rag_context(user_input.strip())
                finally:
                    st.session_state.rag_project_filter_selector = current_project_filter
                if fallback_context and fallback_context != rag_context:
                    rag_context = (
                        fallback_context
                        + "\n\n[Fallback retrieval scope: All Projects]"
                    )
        if rag_context:
            user_message = (
                "[MAIN CHAT RAG: enabled]\n"
                "[MAIN CHAT RAG RULE: Use retrieved context. If insufficient, say so.]\n"
                f"{rag_context}\n\n{user_message}"
            )

    if st.session_state.rag_knowledge_only:
        user_message = (
            "[RAG KNOWLEDGE ONLY: true]\n"
            "[RAG RULE: Answer strictly from retrieved RAG context.]\n"
            f"{user_message}"
        )

    if st.session_state.force_write_todos:
        user_message = (
            "[EXECUTION CONTROL: Call write_todos before taking other actions.]\n"
            f"{user_message}"
        )

    workflow_text = build_pipeline_prompt(st.session_state.pipeline_steps)
    if workflow_text:
        user_message = (
            "[WORKFLOW SEQUENCE: required]\n"
            "[WORKFLOW RULE: Follow the ordered sequence below exactly. Pass output of each step as input to the next.]\n"
            f"{workflow_text}\n\n{user_message}"
        )
    elif looks_like_rag_presentation_request(user_input):
        user_message = (
            "[WORKFLOW SEQUENCE: required]\n"
            "[WORKFLOW RULE: First use the RAG specialist to retrieve and synthesize grounded evidence, "
            "then hand a slide-ready grounded outline to the Presenter to create the PPTX. "
            "Do not have the RAG specialist generate the final presentation file directly.]\n"
            "1. 📚 RAG Search (RAG SubAgent)\n"
            "2. 📊 Create Presentation (Presenter)\n\n"
            f"{user_message}"
        )

    if st.session_state.selected_agent != AUTO_AGENT_DISPLAY_NAME and not st.session_state.rag_knowledge_only:
        agent_name = re.sub(r"^\[[^\]]+\]\s*", "", st.session_state.selected_agent.split("(")[0].strip())
        user_message = f"[AGENT PREFERENCE: Use {agent_name} for this task]\n\n{user_message}"

    return user_message


def _render_main_settings() -> None:
    """Render the main-mode settings organized into logical sections with collapsible expanders."""
    with st.expander("Settings", expanded=False):
        # ──────────────────────────────────────────────────────────────────
        # Section 1: Agent & Model Selection
        # ──────────────────────────────────────────────────────────────────
        with st.expander("Agent & Model", expanded=True):
            ctrl_col1, ctrl_col2, ctrl_col3 = st.columns(3)
            with ctrl_col1:
                agent_display_options = registry.get_all_agents()
                agent_options = list(agent_display_options.keys())
                selected_agent = st.selectbox(
                    "Agent",
                    agent_options,
                    index=(
                        agent_options.index(st.session_state.selected_agent)
                        if st.session_state.selected_agent in agent_options
                        else 0
                    ),
                    key="main_agent_selector",
                    help="Auto routes to the best specialist.",
                )
                st.session_state.selected_agent = selected_agent
            with ctrl_col2:
                _render_provider_selector()
                _render_local_context_preset_selector(
                    st.session_state.get("llm_provider") or get_default_llm_provider()
                )
            with ctrl_col3:
                _render_model_selector(
                    "Model", session_key="current_model", env_key="DEEPAGENT_MODEL",
                    widget_key="main_model_selector",
                )

            _render_agent_model_overrides()

            status = get_status_info()
            parts = [f"**Agent:** {selected_agent.split('(')[0].strip()}"]
            if st.session_state.pipeline_steps:
                parts.append("**Pipeline:** " + " → ".join(get_pipeline_step_agents(st.session_state.pipeline_steps)))
            if status["uploaded_file"]:
                parts.append(f"**File:** {status['uploaded_file']}")
            if status["execution_time"]:
                parts.append(f"**Last:** {status['execution_time']:.1f}s")
            st.caption(" · ".join(parts))

        # ──────────────────────────────────────────────────────────────────
        # Section 2: Execution & Processing
        # ──────────────────────────────────────────────────────────────────
        # llama.cpp performance controls (only shown when provider is llama_cpp)
        _active_provider = st.session_state.get("llm_provider") or get_default_llm_provider()
        if _active_provider == "llama_cpp":
            with st.expander("llama.cpp", expanded=True):
                _render_llama_cpp_model_path_selector()
                st.divider()
                lc1, lc2 = st.columns(2)
                with lc1:
                    _n_ctx = st.number_input(
                        "Context window (n_ctx)",
                        min_value=512,
                        max_value=262144,
                        value=int(os.getenv("LLAMA_CPP_N_CTX", "8192")),
                        step=512,
                        help="Number of tokens the model can see at once. Higher = more memory.",
                    )
                    if str(_n_ctx) != os.getenv("LLAMA_CPP_N_CTX", "8192"):
                        set_env_value("LLAMA_CPP_N_CTX", str(_n_ctx))
                        reload_env_file()
                        _reload_model_modules()
                with lc2:
                    _n_gpu = st.number_input(
                        "GPU layers (n_gpu_layers)",
                        min_value=-1,
                        max_value=200,
                        value=int(os.getenv("LLAMA_CPP_N_GPU_LAYERS", "-1")),
                        step=1,
                        help="-1 = all layers on GPU. 0 = CPU only. Partial offload for limited VRAM.",
                    )
                    if str(_n_gpu) != os.getenv("LLAMA_CPP_N_GPU_LAYERS", "-1"):
                        set_env_value("LLAMA_CPP_N_GPU_LAYERS", str(_n_gpu))
                        reload_env_file()
                        _reload_model_modules()
                st.caption(
                    "ℹ️ Flash Attention is now controlled by `LLAMA_CPP_FLASH_ATTN` "
                    "(default: enabled). Set it to `0` to disable. If unavailable in the "
                    "installed wheel/build, llama.cpp will ignore it and log a warning."
                )

        with st.expander("Execution", expanded=False):
            exec_col1, exec_col2 = st.columns(2)
            with exec_col1:
                st.session_state.main_chat_use_direct_model = st.toggle(
                    "Direct chat model for main chat",
                    value=st.session_state.main_chat_use_direct_model,
                    help="On: send main chat directly to the current model. Off: use the full Deep Agent orchestration stack.",
                )
                st.session_state.force_write_todos = st.checkbox(
                    "Require write_todos planning",
                    value=st.session_state.force_write_todos,
                    help="Forces the agent to plan complex tasks before execution.",
                )
            with exec_col2:
                st.session_state.disable_agent_cache = st.checkbox(
                    "Disable agent cache",
                    value=st.session_state.disable_agent_cache,
                    help="Force fresh computation by clearing agent cache before processing. Useful when cached results are outdated or incorrect.",
                )

        # ──────────────────────────────────────────────────────────────────
        # Section 3: RAG Integration
        # ──────────────────────────────────────────────────────────────────
        with st.expander("RAG", expanded=False):
            st.session_state.rag_enable_main_chat = st.checkbox(
                "Augment main-agent chat with RAG",
                value=st.session_state.rag_enable_main_chat,
                help="Enable RAG context retrieval to augment agent responses with knowledge base information.",
            )
            st.caption("Adds retrieval context from the RAG knowledge base to main-chat replies.")

        # ──────────────────────────────────────────────────────────────────
        # Section 4: Voice & Audio
        # ──────────────────────────────────────────────────────────────────
        with st.expander("Voice", expanded=False):
            voice_col1, voice_col2 = st.columns(2)
            with voice_col1:
                st.session_state.voice_chat_enabled = st.checkbox(
                    "Enable voice chat in main mode",
                    value=st.session_state.voice_chat_enabled,
                    help="Adds microphone input and speech playback for assistant replies.",
                )
                st.session_state.voice_auto_speak = st.checkbox(
                    "Auto-speak new agent replies",
                    value=st.session_state.voice_auto_speak,
                    disabled=not st.session_state.voice_chat_enabled,
                    help="Automatically play audio when the agent responds.",
                )
            with voice_col2:
                st.session_state.voice_speech_rate = st.slider(
                    "Speech rate",
                    min_value=0.7,
                    max_value=1.3,
                    value=float(st.session_state.voice_speech_rate),
                    step=0.05,
                    disabled=not st.session_state.voice_chat_enabled,
                    help="Adjust the playback speed of agent speech (0.7x to 1.3x)",
                )
            
            _tts_ok, _tts_msg = _coqui_tts_available()
            if _tts_ok:
                st.caption("Coqui TTS is available.")
            else:
                st.warning(f"⚠️ Coqui TTS unavailable: {_tts_msg}")

        # ──────────────────────────────────────────────────────────────────
        # Section 5: Learning & Memory
        # ──────────────────────────────────────────────────────────────────
        with st.expander("Memory", expanded=False):
            learn_col1, learn_col2 = st.columns(2)
            with learn_col1:
                st.session_state.self_learning_enabled = st.checkbox(
                    "Enable self-learning",
                    value=st.session_state.self_learning_enabled,
                    help="Controls whether feedback and reusable learnings are written to agent memory.",
                )
                os.environ["SELF_LEARNING_ENABLED"] = (
                    "true" if st.session_state.self_learning_enabled else "false"
                )
            with learn_col2:
                if st.session_state.self_learning_enabled:
                    st.caption("✅ Feedback and `record_learning` can write to memory.")
                else:
                    st.caption("⏸️ Feedback and `record_learning` writes are disabled.")

        # ──────────────────────────────────────────────────────────────────
        # Section 6: System Info
        # ──────────────────────────────────────────────────────────────────
        with st.expander("System", expanded=False):
            runtime_status = get_runtime_feature_status()
            info_col1, info_col2 = st.columns(2)
            with info_col1:
                if runtime_status["cache_level"] == "success":
                    st.caption(f"✅ {runtime_status['cache_status']}")
                else:
                    st.caption(f"⚠️ {runtime_status['cache_status']}")
            with info_col2:
                st.caption(f"📦 DeepAgents: v{runtime_status['deepagents_version']}")
        with st.expander("Attachments", expanded=False):
            _attach_col1, _attach_col2 = st.columns([1, 3])
            with _attach_col1:
                st.radio(
                    "Attachment mode",
                    options=["Ask in chat", "Add to RAG"],
                    horizontal=True,
                    key="main_chat_attachment_mode",
                    label_visibility="collapsed",
                )
            with _attach_col2:
                st.caption("Attach images, PDFs, spreadsheets, documents, CSV, JSON, TXT, or Markdown files.")
                if st.session_state.main_chat_attachment_mode == "Add to RAG":
                    st.caption("Route: ingest through RAG before answering.")
                else:
                    st.caption("Route: spreadsheets to Data Scientist; documents, images, and decks stay in Main Agent unless delegated.")

# ── Mode Renderers ────────────────────────────────────────────────────────────

def _mode_main() -> None:
    """Main Agent mode — orchestrator with agent + model selectors on the main panel."""
    # Show welcome header first when there are no messages
    if not st.session_state.messages:
        _render_welcome_screen()

    _render_main_settings()
    _render_main_chat_helpers()
    _render_main_chat_history_or_empty_state()

    main_chat_started_at = time.time()
    submission = _get_main_chat_submission()
    if not submission:
        return
    user_input, uploaded_chat_files = submission

    # Clear cache if disabled
    if st.session_state.disable_agent_cache:
        get_cached_agent_bundle.clear()
        get_cached_data_scientist_bundle.clear()
        get_cached_ragsub_bundle.clear()
        get_cached_specialist_router_bundle.clear()
        get_cached_presenter_bundle.clear()

    saved_chat_paths = (
        _save_uploaded_streamlit_files(
            uploaded_chat_files,
            target_dir=PROJECT_DIR / "temp_chat_upload",
        )
        if uploaded_chat_files
        else []
    )
    chat_relpaths = _attachment_paths_to_relpaths(saved_chat_paths)
    attachment_mode = st.session_state.main_chat_attachment_mode
    (
        user_message,
        display_user_input,
        attachment_route,
        forced_agent_type,
        force_main_agent,
        normalized_user_input,
    ) = _augment_main_chat_user_message(
        user_input=user_input,
        chat_relpaths=chat_relpaths,
        attachment_mode=attachment_mode,
    )
    user_message = _apply_main_chat_context_overrides(user_message, normalized_user_input)

    # Full DeepAgents autonomy: no forced routing regardless of attachments
    _chat_execution_block(
        user_message,
        display_user_input,
        force_agent_type=forced_agent_type,
        force_main_agent=force_main_agent,
        user_message_meta={
            "attachment_paths": chat_relpaths,
            "attachment_mode": attachment_mode,
            "attachment_route": attachment_route,
        } if chat_relpaths else None,
        started_at=main_chat_started_at,
    )


def _mode_scan_to_text() -> None:
    """Scan to Text mode — OCR scanned images and PDFs into copyable text."""
    st.markdown("### 📝 Scan to Text")
    st.caption("Upload scanned PDFs or images, run OCR, and export the extracted text.")

    try:
        ocr_status = _get_rag_tools().get_multimodal_ocr_status()
    except Exception:
        ocr_status = None

    if isinstance(ocr_status, dict):
        if ocr_status.get("enabled"):
            st.caption(f"✅ OCR: {ocr_status.get('message', '')}")
        else:
            st.caption(f"⚠️ OCR: {ocr_status.get('message', '')}")
            configured_path = str(ocr_status.get("configured_path") or "").strip()
            detected_path = str(ocr_status.get("detected_path") or "").strip()
            remediation = str(ocr_status.get("remediation") or "").strip()
            if configured_path:
                st.caption(f"Configured `TESSERACT_CMD`: `{configured_path}`")
            elif detected_path:
                st.caption(f"Detected Tesseract binary: `{detected_path}`")
            elif remediation:
                st.caption(remediation)

    uploaded_files = st.file_uploader(
        "Upload scanned PDFs or images",
        type=["pdf", "png", "jpg", "jpeg", "webp", "bmp", "gif", "tif", "tiff"],
        accept_multiple_files=True,
        help="Scanned PDFs and image files are supported.",
        key="scan_to_text_uploader",
    )

    action_col1, action_col2 = st.columns([1, 1])
    with action_col1:
        convert_clicked = st.button(
            "🔎 Convert to Text",
            width="stretch",
            key="scan_to_text_convert_btn",
            disabled=not uploaded_files,
        )
    with action_col2:
        if st.button("🗑️ Clear Output", width="stretch", key="scan_to_text_clear_btn"):
            st.session_state.scan_to_text_result = ""
            st.session_state.scan_to_text_sources = []
            st.rerun()

    if convert_clicked and uploaded_files:
        saved_paths = _save_uploaded_streamlit_files(
            list(uploaded_files),
            target_dir=PROJECT_DIR / "tmp" / "scan_to_text_upload",
        )
        extracted_blocks: list[str] = []
        failed_files: list[str] = []

        with st.spinner("Running OCR on uploaded files..."):
            for uploaded, saved_path in zip(uploaded_files, saved_paths):
                extracted_text = extract_text_from_path(saved_path).strip()
                if extracted_text:
                    extracted_blocks.append(
                        f"===== {uploaded.name} =====\n{extracted_text}"
                    )
                else:
                    failed_files.append(uploaded.name)

        combined_text = "\n\n".join(extracted_blocks).strip()
        st.session_state.scan_to_text_result = combined_text
        st.session_state.scan_to_text_sources = [file.name for file in uploaded_files]

        if combined_text:
            st.success(f"Extracted text from {len(extracted_blocks)} file(s).")
        else:
            st.warning("No readable text was extracted from the uploaded files.")

        if failed_files:
            st.caption(
                "No text detected in: " + ", ".join(f"`{name}`" for name in failed_files)
            )

    result_text = str(st.session_state.scan_to_text_result or "").strip()
    if result_text:
        st.divider()
        st.subheader("Extracted Text")
        st.text_area("OCR output", value=result_text, height=420, label_visibility="collapsed")
        output_name = "scanned_text.txt"
        st.download_button(
            label="⬇️ Download .txt",
            data=result_text,
            file_name=output_name,
            mime="text/plain",
            width="stretch",
            key="scan_to_text_download_btn",
        )
        source_count = len(st.session_state.scan_to_text_sources or [])
        st.caption(f"{len(result_text):,} characters from {source_count} uploaded file(s).")
    else:
        st.info("Upload one or more scanned files, then click Convert to Text.")


def _mode_rag_databases() -> None:
    """RAG Databases mode — quick access to the stored RAG index browser."""
    st.markdown("### 🗄️ RAG Databases")
    st.caption("Browse indexed projects, chunk counts, themes, and file-level storage details.")

    try:
        _rag_tools = _get_rag_tools()
        get_rag_index_summary = _rag_tools.get_rag_index_summary
        get_rag_projects = _rag_tools.get_rag_projects
        get_rag_themes = _rag_tools.get_rag_available_themes
        delete_rag_documents = _rag_tools.delete_rag_documents
    except Exception as exc:
        st.error(f"Unable to load RAG database tools: {exc}")
        return

    try:
        rag_projects = get_rag_projects() if get_rag_projects else []
    except Exception:
        rag_projects = []

    project_options = ["All Projects", "Default"] + [p for p in rag_projects if p != "Default"]
    current_project = st.session_state.get("rag_store_project_filter", st.session_state.get("rag_project_mode", "All Projects"))
    if current_project not in project_options:
        current_project = "All Projects"

    top_col1, top_col2 = st.columns([2, 1])
    with top_col1:
        selected_project = st.selectbox(
            "Project",
            options=project_options,
            index=project_options.index(current_project),
            key="rag_store_project_filter",
            help="Select a project to inspect its stored chunks and themes.",
        )
    with top_col2:
        if st.button("📚 Open RAG workspace", width="stretch", key="open_rag_workspace_btn"):
            st.session_state.active_mode = "rag"
            st.rerun()

    try:
        if selected_project == "All Projects":
            all_files: dict[str, dict[str, Any]] = {}
            total_chunks = 0
            for proj in project_options:
                if proj == "All Projects":
                    continue
                try:
                    summary = get_rag_index_summary(project=proj)
                    total_chunks += summary.get("total_chunks", 0)
                    for src, info in summary.get("files", {}).items():
                        prefixed_src = f"{proj}/{src}" if proj != "Default" else src
                        all_files[prefixed_src] = info
                except Exception:
                    continue
            summary = {"total_chunks": total_chunks, "files": all_files}
        else:
            summary = get_rag_index_summary(project=selected_project)
    except Exception as exc:
        st.error(f"Unable to load index summary: {exc}")
        return

    try:
        themes = get_rag_themes(project=selected_project) if get_rag_themes else []
    except Exception:
        themes = []

    files = summary.get("files", {})
    total_chunks = summary.get("total_chunks", 0)

    stat_col1, stat_col2, stat_col3 = st.columns(3)
    stat_col1.metric("Projects", len(rag_projects) or 0)
    stat_col2.metric("Files", len(files) or 0)
    stat_col3.metric("Chunks", total_chunks or 0)

    if themes:
        st.caption("Themes: " + ", ".join(f"`{theme}`" for theme in themes[:12]))

    if not files:
        st.info("No documents indexed yet. Use the RAG workspace to ingest files first.")
        return

    import pandas as pd

    rows: list[dict[str, Any]] = []
    for src, info in sorted(files.items()):
        if "/" in src and selected_project == "All Projects":
            proj_name, file_name = src.split("/", 1)
        elif selected_project == "All Projects":
            proj_name = "Default"
            file_name = src
        else:
            proj_name = selected_project
            file_name = src

        rows.append(
            {
                "Project": proj_name,
                "File": file_name,
                "Theme": ", ".join(sorted(info.get("themes", []))),
                "Chunks": info.get("chunks", 0),
                "Modalities": ", ".join(
                    f"{name}:{count}" for name, count in sorted(info.get("modalities", {}).items())
                ) or "text",
                "Table Extraction": ", ".join(
                    f"{name}:{count}" for name, count in sorted(info.get("table_extraction_methods", {}).items())
                ) or "—",
                "Vision Captioned": info.get("vision_captioned_chunks", 0),
                "Added": info.get("date_added", "—"),
            }
        )

    st.dataframe(pd.DataFrame(rows), width="stretch", hide_index=True)

    with st.expander("Delete stored chunks", expanded=False):
        st.caption("Delete entire projects or narrow the deletion by theme/file inside the main RAG workspace if needed.")
        project_scope = st.selectbox(
            "Delete project",
            options=[""] + [p for p in project_options if p not in {"All Projects"}],
            key="rag_db_delete_project_scope",
            help="Choose a specific project to delete.",
        )
        if not project_scope:
            st.warning("Choose a project to enable deletion.")
        elif st.button("🧹 Delete project", width="stretch", key="rag_db_delete_project_btn"):
            deleted = delete_rag_documents(project=project_scope, allow_delete_all=True) if delete_rag_documents else 0
            st.success(f"Deleted {deleted} chunks from project '{project_scope}'")
            st.rerun()


def _mode_data_analysis() -> None:
    """Data Analysis mode — file upload + dedicated Data Scientist agent."""
    st.markdown("### 📊 Data Analysis")

    ctrl_col1, ctrl_col2 = st.columns(2)
    with ctrl_col1:
        _render_model_selector(
            "Data Scientist Model", session_key="model_data", env_key="DATA_SCIENTIST_MODEL",
            widget_key="data_model_selector", allow_inherit=True,
        )
    with ctrl_col2:
        status = get_status_info()
        if status["uploaded_file"]:
            st.metric(
                "Loaded File", status["uploaded_file"],
                delta=f"{status['file_size_kb']:.0f} KB" if status["file_size_kb"] else None,
                delta_color="off",
            )
        else:
            st.info("No dataset loaded yet")

    # Contextual file upload — only shown in this mode
    with st.expander("📁 Dataset Upload", expanded=st.session_state.uploaded_file_path is None):
        uploaded_file = st.file_uploader(
            "Upload CSV or Excel file",
            type=["csv", "xlsx", "xls"],
            help="Upload a dataset for the Data Scientist agent to analyse",
            key="ds_file_uploader",
        )
        if uploaded_file is not None:
            file_path = PROJECT_DIR / "temp_upload" / uploaded_file.name
            file_path.parent.mkdir(exist_ok=True)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            st.session_state.uploaded_file_path = str(file_path)
            st.success(f"Uploaded: {uploaded_file.name}")
            try:
                import pandas as pd
                df = (
                    pd.read_csv(file_path)
                    if uploaded_file.name.endswith(".csv")
                    else pd.read_excel(file_path)
                )
                st.dataframe(df.head(3), width="stretch")
                st.caption(f"{df.shape[0]} rows × {df.shape[1]} cols | {uploaded_file.size / 1024:.1f} KB")
            except Exception as e:
                st.error(f"Preview failed: {e}")

    st.divider()
    _render_chat_history()

    user_input = st.chat_input(
        "Describe the analysis you want — EDA, regression, clustering, or 'generate a .docx report'..."
    )
    if not user_input:
        return

    user_message = user_input.strip()

    if st.session_state.uploaded_file_path:
        fp = Path(st.session_state.uploaded_file_path)
        rel = _safe_project_relative_path(fp)
        if fp.exists() and rel:
            user_message = (
                f"[UPLOADED FILE: {rel}]\n"
                f"[DATA CHECK: Use execute_python_code for exact dataset shape/statistics.]\n\n"
                f"{user_message}"
            )

    if st.session_state.force_write_todos:
        user_message = (
            "[EXECUTION CONTROL: Call write_todos before taking other actions.]\n"
            f"{user_message}"
        )

    # Full DeepAgents autonomy: no forced routing
    _chat_execution_block(user_message, user_input.strip(), force_agent_type=None)


# ── RAG Retrieval Presets ─────────────────────────────────────────────────────

_RAG_PRESETS: dict[str, dict] = {
    _RAG_QA_PRESET_NAME: {
        **_RAG_QA_PRESET,
        "agent_label": "RAG SubAgent",
        "desc": "Best chunks globally — RAG SubAgent answers directly from retrieved context.",
    },
    "🖼️ Presentation": {
        "mode": "Top-K Per File", "top_k": 8, "fetch_k": 150, "max_files": 6, "min_rerank": 0.05,
        "agent_label": "Presenter",
        "desc": "Best chunks per document — Presenter builds slides from pre-retrieved content.",
    },
    "📂 Project summary": {
        "mode": "MMR", "top_k": 10, "fetch_k": 200, "max_files": 8, "min_rerank": 0.00,
        "agent_label": "RAG SubAgent",
        "desc": "Diverse chunks across all files — RAG SubAgent summarises the project directly.",
    },
    "🏷️ Theme summary": {
        "mode": "MMR", "top_k": 8, "fetch_k": 150, "max_files": 6, "min_rerank": 0.05,
        "agent_label": "RAG SubAgent",
        "desc": "Diverse chunks scoped to the selected theme — RAG SubAgent summarises by theme.",
    },
    "📄 File summary": {
        "mode": "Top-K Per File", "top_k": 8, "fetch_k": 120, "max_files": 5, "min_rerank": 0.00,
        "agent_label": "RAG SubAgent",
        "desc": "Best chunks per file — RAG SubAgent summarises each document.",
    },
    "📝 Write report": {
        "mode": "MMR", "top_k": 10, "fetch_k": 200, "max_files": 8, "min_rerank": 0.00,
        "agent_label": "Writer (Quarto → docx)",
        "desc": "Broad, diverse context — Writer produces a structured .docx report via Quarto.",
    },
    "🔬 Literature review": {
        "mode": "Top-K Globally", "top_k": 12, "fetch_k": 250, "max_files": 8, "min_rerank": 0.00,
        "agent_label": "Writer (Quarto → docx)",
        "desc": "High recall globally — Writer produces a citation-heavy academic .docx via Quarto.",
    },
}
_RAG_PRESETS[_RAG_QA_PRESET_NAME]["min_rerank"] = 0.0
_DEFAULT_RAG_PRESET_NAME = _RAG_QA_PRESET_NAME


def _apply_rag_preset(
    preset_name: str,
    mode: str,
    top_k: int,
    fetch_k: int,
    max_files: int,
    min_rerank: float,
) -> None:
    """Callback: apply a retrieval preset to all relevant session state keys."""
    st.session_state.rag_active_preset = preset_name
    st.session_state.rag_retrieval_mode = mode
    st.session_state.rag_top_k = top_k
    st.session_state.rag_fetch_k = fetch_k
    st.session_state.rag_max_files = max_files
    st.session_state.rag_min_rerank_score = min_rerank
    st.session_state.rag_top_k_slider_v2 = top_k
    st.session_state.rag_fetch_k_slider_v2 = fetch_k
    st.session_state.rag_max_files_slider_v2 = max_files
    st.session_state.rag_min_rerank_score_slider_v2 = min_rerank


def _list_rag_file_chunks(
    rag_tools: Any,
    *,
    project: str,
    file_name: str,
) -> list[dict[str, Any]]:
    """Return all stored chunks for a file without using chat-style retrieval."""
    list_docs = getattr(rag_tools, "list_rag_documents", None)
    if not callable(list_docs):
        return []

    project = (project or "").strip()
    file_name = (file_name or "").strip()

    call_attempts: list[dict[str, Any]] = [
        {"project": project, "source": file_name},
        {"project": project, "file": file_name},
        {"project": project, "filename": file_name},
        {"project": project},
        {},
    ]

    docs: Any = None
    for kwargs in call_attempts:
        try:
            docs = list_docs(**{k: v for k, v in kwargs.items() if v})
            break
        except TypeError:
            continue
        except Exception:
            continue

    if docs is None:
        return []

    if isinstance(docs, dict):
        if "documents" in docs and isinstance(docs["documents"], list):
            iterable: list[Any] = docs["documents"]
        elif "items" in docs and isinstance(docs["items"], list):
            iterable = docs["items"]
        elif "chunks" in docs and isinstance(docs["chunks"], list):
            iterable = docs["chunks"]
        else:
            iterable = list(docs.values())
    else:
        iterable = list(docs)

    wanted_name = Path(file_name).name
    chunks: list[dict[str, Any]] = []
    for item in iterable:
        if isinstance(item, dict):
            text = str(item.get("text") or item.get("content") or item.get("page_content") or "").strip()
            metadata = item.get("metadata") if isinstance(item.get("metadata"), dict) else {}
            if not metadata and "source" in item:
                metadata = item
        else:
            text = str(getattr(item, "page_content", "") or getattr(item, "text", "") or item).strip()
            metadata = getattr(item, "metadata", {}) if isinstance(getattr(item, "metadata", {}), dict) else {}

        source_value = str(metadata.get("source") or metadata.get("file") or metadata.get("filename") or "")
        source_name = Path(source_value).name if source_value else ""
        if file_name and source_name and source_name != wanted_name and file_name not in source_value:
            continue
        if file_name and not source_value and wanted_name not in text:
            continue

        chunks.append(
            {
                "text": text,
                "metadata": metadata,
            }
        )

    if chunks:
        return chunks

    # Fallback: inspect the persistent Chroma store directly and filter by source metadata.
    try:
        import chromadb
    except Exception:
        return []

    chroma_dir = PROJECT_DIR / "rag-chroma"
    if not chroma_dir.exists():
        return []

    wanted_name = Path(file_name).name.lower()
    wanted_source = (file_name or "").strip().lower()
    wanted_project = (project or "").strip().lower()
    metadata_keys = ("source", "file", "filename", "path", "source_path", "chunk_source")

    try:
        client = chromadb.PersistentClient(path=str(chroma_dir))
        collections = list(client.list_collections())
    except Exception:
        return []

    def _collection_name(col: Any) -> str:
        if isinstance(col, str):
            return col
        return str(getattr(col, "name", "") or (col["name"] if isinstance(col, dict) and "name" in col else ""))

    for collection_info in collections:
        try:
            collection_name = _collection_name(collection_info)
            collection = client.get_collection(collection_name)
            data = collection.get(include=["documents", "metadatas", "ids"])
        except Exception:
            continue

        docs = list(data.get("documents") or [])
        metas = list(data.get("metadatas") or [])
        ids = list(data.get("ids") or [])

        for idx, doc_text in enumerate(docs):
            metadata = metas[idx] if idx < len(metas) and isinstance(metas[idx], dict) else {}

            metadata_blob = json.dumps(metadata, ensure_ascii=False, default=str).lower()
            source_value = ""
            for key in metadata_keys:
                raw_value = metadata.get(key)
                if raw_value:
                    source_value = str(raw_value)
                    break

            source_value_l = source_value.lower()
            doc_text_l = str(doc_text).lower()
            project_match = True
            if wanted_project:
                meta_project = str(metadata.get("project") or metadata.get("collection") or "").strip().lower()
                project_match = (
                    meta_project == wanted_project
                    or wanted_project in source_value_l
                    or wanted_project in metadata_blob
                    or wanted_project in collection_name.lower()
                )
                if not project_match and wanted_project != "default":
                    continue

            file_match = False
            if wanted_source:
                file_match = (
                    wanted_source in source_value_l
                    or wanted_name in Path(source_value).name.lower()
                    or wanted_source in doc_text_l
                    or wanted_name in doc_text_l
                    or wanted_source in metadata_blob
                    or wanted_name in metadata_blob
                )
            else:
                file_match = True

            if not file_match and not project_match:
                continue
            if not file_match and wanted_source:
                continue

            chunks.append(
                {
                    "text": str(doc_text or "").strip(),
                    "metadata": metadata,
                    "id": ids[idx] if idx < len(ids) else "",
                    "collection": collection_name,
                }
            )

    return chunks


def _mode_rag() -> None:
    """RAG mode — document management, ingestion controls, and knowledge retrieval chat."""
    st.markdown(
        """
        <div class="page-shell">
            <div class="page-kicker">Knowledge Base</div>
            <h1 class="page-title" style="font-size: 1.55rem; margin-bottom: 0;">RAG</h1>
        </div>
        """,
        unsafe_allow_html=True,
    )

    if not st.session_state.get("rag_ui_initialized"):
        _apply_rag_preset(
            _DEFAULT_RAG_PRESET_NAME,
            _RAG_QA_PRESET["mode"],
            _RAG_QA_PRESET["top_k"],
            _RAG_QA_PRESET["fetch_k"],
            _RAG_QA_PRESET["max_files"],
            _RAG_QA_PRESET["min_rerank"],
        )
        st.session_state.rag_ui_initialized = True


    try:
        _rag_tools = _get_rag_tools()
        _raw_ingest_rag_paths = _rag_tools.ingest_rag_paths
        ingest_web_search_results = _rag_tools.ingest_web_search_results
        _get_rag_idx = _rag_tools.get_rag_index_summary
        _get_rag_projects = _rag_tools.get_rag_projects
        _get_rag_themes = _rag_tools.get_rag_available_themes
        delete_rag_documents = _rag_tools.delete_rag_documents
        get_last_rag_query_diagnostics = _rag_tools.get_last_rag_query_diagnostics
        get_multimodal_ocr_status = _rag_tools.get_multimodal_ocr_status
        get_multimodal_vision_status = _rag_tools.get_multimodal_vision_status
    except Exception:
        _raw_ingest_rag_paths = _get_rag_idx = _get_rag_projects = _get_rag_themes = delete_rag_documents = None
        ingest_web_search_results = None
        get_last_rag_query_diagnostics = None
        get_multimodal_ocr_status = None
        get_multimodal_vision_status = None

    def ingest_rag_paths(*args: Any, **kwargs: Any):
        """Wrap ingest so scanned PDFs/images get OCR text sidecars indexed."""
        if _raw_ingest_rag_paths is None:
            raise RuntimeError("RAG ingest is unavailable.")

        prepared_args = list(args)
        if prepared_args:
            first_arg = prepared_args[0]
            if isinstance(first_arg, (list, tuple)):
                prepared_args[0] = _prepare_rag_ingest_paths([Path(p) for p in first_arg])
        prepared_kwargs = dict(kwargs)

        for key in ("paths", "file_paths", "files", "documents"):
            value = prepared_kwargs.get(key)
            if isinstance(value, (list, tuple)):
                prepared_kwargs[key] = _prepare_rag_ingest_paths([Path(p) for p in value])
                break

        return _raw_ingest_rag_paths(*prepared_args, **prepared_kwargs)

    def _queue_rag_reingest() -> None:
        """Queue a rerun that will auto-run ingestion with the current staged files."""
        if not st.session_state.rag_uploaded_files:
            return
        st.session_state["rag_auto_ingest_pending"] = True
        st.rerun()

    retrieval_settings_tab, transcript_tab, ingest_tab, stored_docs_tab = st.tabs([
        "Retrieval",
        "YouTube",
        "Ingest",
        "Documents",
    ])
    user_input = None

    with transcript_tab:
        yt_url = st.text_input(
            "YouTube video URL",
            key="rag_youtube_url_mode",
            placeholder="https://www.youtube.com/watch?v=...",
        )
        if st.button("Download", key="rag_youtube_transcript_btn", width="stretch"):
            if not yt_url.strip():
                st.warning("Enter a YouTube URL first.")
            else:
                try:
                    out_path, video_title = _download_youtube_transcript_to_txt(
                        yt_url.strip(),
                        PROJECT_DIR / "youtube_transcripts",
                    )
                    st.success(f"Saved transcript: {video_title}")
                    st.caption(f"File: `{out_path}`")
                except Exception as exc:
                    st.error(f"Unable to download transcript: {exc}")

        st.markdown("---")
        yt_playlist_url = st.text_input(
            "YouTube playlist URL",
            key="rag_youtube_playlist_url_mode",
            placeholder="https://www.youtube.com/playlist?list=...",
        )
        if st.button("Download playlist", key="rag_youtube_playlist_transcript_btn", width="stretch"):
            if not yt_playlist_url.strip():
                st.warning("Enter a YouTube playlist URL first.")
            else:
                try:
                    playlist_result = _download_youtube_playlist_transcripts_to_txt(
                        yt_playlist_url.strip(),
                        PROJECT_DIR / "youtube_transcripts",
                    )
                    saved_count = len(playlist_result["saved_files"])
                    failed_count = len(playlist_result["failures"])
                    st.success(
                        f"Saved {saved_count}/{playlist_result['total_videos']} transcripts "
                        f"for playlist: {playlist_result['playlist_title']}"
                    )
                    st.caption(f"Folder: `{playlist_result['playlist_dir']}`")
                    if failed_count:
                        with st.expander(f"Show failures ({failed_count})", expanded=False):
                            for failure in playlist_result["failures"]:
                                st.caption(failure)
                except Exception as exc:
                    st.error(f"Unable to download playlist transcripts: {exc}")

    # Contextual ingest panel — only in this mode
    with ingest_tab:
        ingest_top_col1, ingest_top_col2, ingest_top_col3 = st.columns(3)
        
        current_embed_config = os.getenv("RAG_EMBED_MODEL", "llama_server:").strip()
        current_embed_provider, current_embed_model = resolve_provider_and_model(current_embed_config)

        def _apply_rag_embed_config(provider: str, model: str, *, rerun: bool = True) -> None:
            full_config = f"{provider}:{model.strip()}" if model.strip() else f"{provider}:"
            set_env_value("RAG_EMBED_MODEL", full_config)
            st.session_state["model_rag_embed"] = full_config
            reload_env_file()
            st.success(f"Embedding config updated: {full_config}")
            if rerun:
                st.rerun()

        def _pick_fallback_ollama_embedding() -> str | None:
            ollama_models = get_ollama_models()
            ollama_model_names = sorted(ollama_models.keys())
            return ollama_model_names[0] if ollama_model_names else None
        
        with ingest_top_col1:
            # Provider selector - fully flexible
            embed_providers = ["llama_server", "ollama", "llama_cpp", "custom"]
            if current_embed_provider not in embed_providers:
                embed_providers.insert(0, current_embed_provider)
            
            selected_embed_provider = st.selectbox(
                "Provider",
                embed_providers,
                index=embed_providers.index(current_embed_provider) if current_embed_provider in embed_providers else 0,
                key="rag_embed_provider_selector",
                help="Select the embedding provider. You can use any provider: llama_server, ollama, llama_cpp, or custom.",
            )
        
        with ingest_top_col2:
            if selected_embed_provider == "llama_cpp":
                placeholder_text = "C:\\models\\embedding.gguf or /home/user/model.gguf"
                help_text = "Full path to embedding model file (any format: .gguf, .bin, etc.)"
                embed_model_input = st.text_input(
                    "Model",
                    value=current_embed_model or "",
                    placeholder=placeholder_text,
                    key="rag_embed_model_input",
                    help=help_text,
                )
            elif selected_embed_provider == "ollama":
                ollama_models = get_ollama_models()
                ollama_model_names = sorted(ollama_models.keys())
                if ollama_model_names:
                    default_ollama_model = (
                        current_embed_model
                        if current_embed_provider == "ollama" and current_embed_model in ollama_model_names
                        else ollama_model_names[0]
                    )
                    embed_model_input = st.selectbox(
                        "Model",
                        ollama_model_names,
                        index=ollama_model_names.index(default_ollama_model),
                        key="rag_embed_model_selector",
                        help="Installed Ollama embedding models fetched from `http://localhost:11434/api/tags`.",
                    )
                else:
                    st.warning("No Ollama models found. Pull one with `ollama pull <model>` and refresh.")
                    if st.button("Refresh Ollama models", key="rag_refresh_ollama_models", use_container_width=True):
                        get_ollama_models.clear()
                        st.rerun()
                    placeholder_text = "mistral, nomic-embed-text, etc."
                    help_text = "Ollama model name from ollama.ai registry"
                    embed_model_input = st.text_input(
                        "Model",
                        value=current_embed_model or "",
                        placeholder=placeholder_text,
                        key="rag_embed_model_input",
                        help=help_text,
                    )
            elif selected_embed_provider == "llama_server":
                placeholder_text = "model_name (auto-detected from server)"
                help_text = "Leave blank for auto-detection, or specify a model name"
                embed_model_input = st.text_input(
                    "Model",
                    value=current_embed_model or "",
                    placeholder=placeholder_text,
                    key="rag_embed_model_input",
                    help=help_text,
                )
            else:
                placeholder_text = "any-model-identifier"
                help_text = "Model identifier for your custom embedding provider"
                embed_model_input = st.text_input(
                    "Model",
                    value=current_embed_model or "",
                    placeholder=placeholder_text,
                    key="rag_embed_model_input",
                    help=help_text,
                )
            
            # Apply configuration when changed
            if embed_model_input != current_embed_model or selected_embed_provider != current_embed_provider:
                if embed_model_input.strip() or selected_embed_provider in ("llama_server", "ollama"):
                    _apply_rag_embed_config(selected_embed_provider, embed_model_input)
        
        with ingest_top_col3:
            st.caption("Current embedding config")
            st.code(current_embed_config, language="text")
            
            if st.button("Reset", key="rag_embed_reset", use_container_width=True):
                _apply_rag_embed_config("llama_server", "")
        
        if selected_embed_provider == "llama_server":
            _emb_ok, _emb_msg = _probe_llama_server_embeddings(current_embed_model or os.getenv("LLAMA_SERVER_MODEL", "local"))
            if not _emb_ok:
                st.error(_emb_msg)
                _ollama_models = get_ollama_models()
                _ollama_model_names = sorted(_ollama_models.keys())
                if _ollama_model_names:
                    _switch_to = _ollama_model_names[0]
                    st.button(
                        f"Switch to Ollama embeddings: {_switch_to}",
                        key="rag_switch_to_ollama_embeddings_btn",
                        use_container_width=True,
                        on_click=_apply_rag_embed_config,
                        args=("ollama", _switch_to),
                    )
                else:
                    st.info("Install an Ollama embedding model or restart llama-server with `--embeddings` enabled.")

        with st.expander("Vision & OCR", expanded=False):
            vision_col1, vision_col2 = st.columns(2)
            with vision_col1:
                st.session_state.rag_image_ingest_mode = st.selectbox(
                    "Image handling",
                    ["OCR only", "OCR + vision caption", "Multimodal vision"],
                    index=[
                        "OCR only",
                        "OCR + vision caption",
                        "Multimodal vision",
                    ].index(st.session_state.rag_image_ingest_mode)
                    if st.session_state.rag_image_ingest_mode in [
                        "OCR only",
                        "OCR + vision caption",
                        "Multimodal vision",
                    ]
                    else 0,
                    key="rag_image_ingest_mode_selector",
                    help=(
                        "OCR only stores extracted text. OCR + vision caption stores OCR text plus a generated description. "
                        "Multimodal vision uses the vision model first, but the current RAG store still indexes text."
                    ),
                )
                _active_vision = _render_model_selector(
                    "Vision model",
                    session_key="model_rag_vision",
                    env_key="RAG_VISION_MODEL",
                    widget_key="rag_vision_model_selector",
                    allow_inherit=False,
                )
                st.caption("Images are indexed as text.")
            
            with vision_col2:
                if get_multimodal_ocr_status:
                    _vision_model_current = st.session_state.get("model_rag_vision") or os.getenv("RAG_VISION_MODEL", "llama_cpp:")
                    _vision_model_current = str(_vision_model_current).strip()
                    _vision_provider, _vision_model = resolve_provider_and_model(_vision_model_current)
                    _vision_provider_options = ["llama_server", "ollama", "llama_cpp", "custom"]
                    if _vision_provider not in _vision_provider_options:
                        _vision_provider_options.insert(0, _vision_provider)
                    _ocr_engine_options = ["Auto", "Tesseract", "Vision"]
                    _ocr_engine_current = str(st.session_state.get("rag_ocr_engine_mode", "Auto") or "Auto")
                    _vision_col1, _vision_col2, _vision_col3, _vision_col4 = st.columns([1.25, 1.25, 1, 1])
                    with _vision_col1:
                        _vision_provider_selected = st.selectbox(
                            "Vision provider",
                            _vision_provider_options,
                            index=_vision_provider_options.index(_vision_provider) if _vision_provider in _vision_provider_options else 0,
                            key="rag_vision_provider_selector",
                            help="Provider used for vision-based document understanding.",
                        )
                    with _vision_col2:
                        _vision_placeholder = (
                            "model_name, e.g. llama-vision"
                            if _vision_provider_selected != "llama_cpp"
                            else "C:\\models\\vision.gguf or /home/user/vision.gguf"
                        )
                        _vision_model_input = st.text_input(
                            "Model",
                            value=_vision_model,
                            placeholder=_vision_placeholder,
                            key="rag_vision_model_input",
                            help="Model name or path used for vision-based extraction and OCR.",
                        )
                    with _vision_col3:
                        st.session_state.rag_ocr_engine_mode = st.selectbox(
                            "OCR",
                            options=_ocr_engine_options,
                            index=_ocr_engine_options.index(_ocr_engine_current) if _ocr_engine_current in _ocr_engine_options else 0,
                            key="rag_ocr_engine_selector",
                            help="OCR backend used when ingesting scanned images or PDFs.",
                        )
                    with _vision_col4:
                        _ocr_status = get_multimodal_ocr_status()
                        st.markdown(
                            f"**{'✅' if bool(_ocr_status.get('enabled')) else '⚠️'} OCR**  \n"
                            f"`{st.session_state.rag_ocr_engine_mode}`"
                        )
                    _tesseract_current = str(os.getenv("TESSERACT_CMD", "") or "").strip()
                    _tesseract_detected = str(_ocr_status.get("detected_path") or "").strip()
                    if not _tesseract_current and _tesseract_detected:
                        _tesseract_current = _tesseract_detected
                    _tesseract_cmd = st.text_input(
                        "Tesseract binary",
                        value=_tesseract_current,
                        placeholder=r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                        key="rag_tesseract_cmd_input",
                        help="Full path to tesseract.exe. Leave blank to use PATH-based auto-detection.",
                    )
                    _tesseract_cmd = _tesseract_cmd.strip()
                    if _tesseract_detected and not os.getenv("TESSERACT_CMD", "").strip():
                        if st.button("Use detected Tesseract", key="rag_use_detected_tesseract_btn"):
                            set_env_value("TESSERACT_CMD", _tesseract_detected)
                            reload_env_file()
                            st.session_state["rag_tesseract_cmd_input"] = _tesseract_detected
                            st.success(f"✅ Using detected Tesseract at {_tesseract_detected}")
                            st.rerun()
                    if st.button("Detect Tesseract", key="rag_detect_tesseract_btn"):
                        _candidate_paths = [
                            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                            r"C:\Program Files\Tesseract\tesseract.exe",
                            r"C:\Program Files (x86)\Tesseract\tesseract.exe",
                        ]
                        _detected_tesseract = ""
                        for _candidate in _candidate_paths:
                            if Path(_candidate).exists():
                                _detected_tesseract = _candidate
                                break
                        if not _detected_tesseract:
                            import shutil

                            _which_tesseract = shutil.which("tesseract")
                            if _which_tesseract:
                                _detected_tesseract = _which_tesseract

                        if _detected_tesseract:
                            set_env_value("TESSERACT_CMD", _detected_tesseract)
                            reload_env_file()
                            st.session_state["rag_tesseract_cmd_input"] = _detected_tesseract
                            st.success(f"✅ Found Tesseract at {_detected_tesseract}")
                            st.rerun()
                        else:
                            st.warning("No Tesseract installation was detected in the common Windows locations or PATH.")
                    if _tesseract_cmd != _tesseract_current:
                        set_env_value("TESSERACT_CMD", _tesseract_cmd)
                        reload_env_file()
                        st.success("✅ Updated TESSERACT_CMD")
                        st.rerun()
                    _vision_full_config = f"{_vision_provider_selected}:{_vision_model_input.strip()}" if _vision_model_input.strip() else f"{_vision_provider_selected}:"
                    if _vision_full_config != _vision_model_current and _vision_model_input is not None:
                        set_env_value("RAG_VISION_MODEL", _vision_full_config)
                        st.session_state["model_rag_vision"] = _vision_full_config
                        reload_env_file()
                        st.success(f"✅ Vision config: {_vision_full_config}")
                        st.rerun()

                    with st.popover("OCR details"):
                        _ocr_configured = str(_ocr_status.get("configured_path") or "").strip()
                        _ocr_detected = str(_ocr_status.get("detected_path") or "").strip()
                        _ocr_remediation = str(_ocr_status.get("remediation") or "").strip()
                        if _ocr_configured:
                            st.caption(f"Configured `TESSERACT_CMD`: `{_ocr_configured}`")
                        elif _ocr_detected:
                            st.caption(f"Detected Tesseract binary: `{_ocr_detected}`")
                        elif _ocr_remediation:
                            st.caption(_ocr_remediation)
                    if _ocr_status.get("enabled"):
                        st.caption(f"OCR: {_ocr_status.get('message', '')}")
                    else:
                        st.caption(f"OCR: {_ocr_status.get('message', '')}")
                        _ocr_configured = str(_ocr_status.get("configured_path") or "").strip()
                        _ocr_detected = str(_ocr_status.get("detected_path") or "").strip()
                        _ocr_remediation = str(_ocr_status.get("remediation") or "").strip()
                        if _ocr_configured:
                            st.caption(f"Configured `TESSERACT_CMD`: `{_ocr_configured}`")
                        elif _ocr_detected:
                            st.caption(f"Detected Tesseract binary: `{_ocr_detected}`")
                        elif _ocr_remediation:
                            st.code(
                                "TESSERACT_CMD=C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
                                language="text",
                            )
                    if st.session_state.rag_image_ingest_mode == "OCR only":
                        st.caption("Path: image -> OCR text -> embeddings.")
                    elif st.session_state.rag_image_ingest_mode == "OCR + vision caption":
                        st.caption("Path: image -> OCR text + vision caption -> embeddings.")
                    else:
                        st.caption("Path: image -> vision extraction -> embeddings.")

        col_proj, col_theme = st.columns(2)
        with col_proj:
            # Display "All Projects" as "Default" in the ingest input for clarity
            _display_proj = "Default" if st.session_state.rag_project_mode == "All Projects" else st.session_state.rag_project_mode
            rag_project = st.text_input(
                "Project",
                value=_display_proj,
                key="rag_project_ingest_mode",
            )
        with col_theme:
            rag_theme = st.text_input("Theme", value="", key="rag_theme_mode")

        _CHUNK_LABELS = {
            "semantic": "Semantic (recommended)",
            "recursive": "Recursive",
            "contextual": "Contextual — Semantic + LLM context (slow)",
            "late_chunking": "Late Chunking — Semantic + neighbour context window",
            "semantic_contextual_late": "Semantic + Contextual + Late — all three combined (slowest)",
        }
        rag_chunk_method = st.selectbox(
            "Chunking",
            list(_CHUNK_LABELS.keys()),
            index=list(_CHUNK_LABELS.keys()).index("semantic_contextual_late"),
            format_func=lambda x: _CHUNK_LABELS.get(x, x),
            key="rag_chunk_method_mode",
        )
        _needs_semantic = rag_chunk_method in (
            "semantic", "contextual", "late_chunking", "semantic_contextual_late"
        )
        _needs_llm = rag_chunk_method in ("contextual", "semantic_contextual_late")
        if _needs_semantic:
            rag_bp = st.slider("Breakpoint threshold", 50, 99, 95, 1, key="rag_bp_mode")
            rag_cs, rag_co = 1500, 300
            if _needs_llm:
                _ctx_model = _render_model_selector(
                    "Context LLM",
                    session_key="model_rag_context",
                    env_key="RAG_CONTEXT_LLM_MODEL",
                    widget_key="rag_context_model_selector",
                    allow_inherit=True,
                    provider_override="all",
                )
                _ctx_model_label = _ctx_model or os.environ.get(
                    AGENT_MODEL_ENV_KEYS["main"],
                    get_agent_model("main"),
                )
                st.caption(f"Each chunk calls the LLM (`{_ctx_model_label}`).")
            if rag_chunk_method in ("late_chunking", "semantic_contextual_late"):
                st.caption("Late chunking adds neighboring context before embedding.")
        else:
            rag_bp = 95
            cs_col, co_col = st.columns(2)
            with cs_col:
                rag_cs = st.number_input("Chunk size", 100, 4000, 1500, 100, key="rag_cs_mode")
            with co_col:
                rag_co = st.number_input("Overlap", 0, 1000, 300, 50, key="rag_co_mode")

        rag_files = st.file_uploader(
            "Upload documents",
            type=["pdf", "docx", "doc", "odt", "pptx", "ppt", "xlsx", "xls", "txt", "md", "csv", "json", "png", "jpg", "jpeg", "webp", "bmp", "gif", "tif", "tiff", "mp3", "wav", "m4a", "aac", "flac", "ogg", "opus", "mp4", "mov", "mkv", "avi", "webm"],
            accept_multiple_files=True,
            key="rag_uploader_mode",
        )
        if rag_files:
            rag_dir = PROJECT_DIR / "temp_rag_upload"
            rag_dir.mkdir(exist_ok=True)
            saved = []
            file_kinds = []
            for up in rag_files:
                fp = rag_dir / up.name
                fp.write_bytes(up.getbuffer())
                saved.append(str(fp))
                file_kinds.append((up.name, _classify_rag_upload_type(up.name)))
            st.session_state.rag_uploaded_files = saved
            st.caption(f"{len(saved)} file(s) ready")
            with st.expander("Preview upload", expanded=False):
                for filename, file_kind in file_kinds:
                    st.caption(f"`{filename}` -> {file_kind}")

        with st.expander("Web ingest", expanded=False):
            web_col1, web_col2 = st.columns(2)
            with web_col1:
                rag_web_query = st.text_input(
                    "Query",
                    value="",
                    key="rag_web_query_mode",
                    placeholder="e.g. ASEAN AI policy updates 2026",
                )
            with web_col2:
                rag_web_topic = st.selectbox(
                    "Topic",
                    options=["general", "news", "finance"],
                    index=0,
                    key="rag_web_topic_mode",
                )
            web_col3, web_col4 = st.columns(2)
            with web_col3:
                rag_web_project = st.text_input(
                    "Project",
                    value=rag_project or "Default",
                    key="rag_web_project_mode",
                )
            with web_col4:
                rag_web_theme = st.text_input(
                    "Theme",
                    value=rag_theme,
                    key="rag_web_theme_mode",
                )
            web_col5, web_col6 = st.columns(2)
            with web_col5:
                rag_web_max_results = st.slider(
                    "Results",
                    min_value=1,
                    max_value=10,
                    value=3,
                    step=1,
                    key="rag_web_max_results_mode",
                )
            with web_col6:
                rag_web_chunk_method = st.selectbox(
                    "Chunking",
                    list(_CHUNK_LABELS.keys()),
                    index=list(_CHUNK_LABELS.keys()).index("recursive"),
                    format_func=lambda x: _CHUNK_LABELS.get(x, x),
                    key="rag_web_chunk_method_mode",
                )
            if st.button("Search and ingest", key="rag_web_ingest_btn", width="stretch"):
                if not rag_web_query.strip():
                    st.warning("Enter a web search query first.")
                elif ingest_web_search_results is None:
                    st.error("Web-search ingestion unavailable.")
                else:
                    with st.spinner("Searching the web and ingesting results into RAG..."):
                        web_result = ingest_web_search_results.invoke(
                            {
                                "query": rag_web_query.strip(),
                                "project": rag_web_project.strip() or "Default",
                                "theme": rag_web_theme.strip(),
                                "max_results": int(rag_web_max_results),
                                "topic": rag_web_topic,
                                "chunking_method": rag_web_chunk_method,
                                "chunk_size": int(rag_cs),
                                "chunk_overlap": int(rag_co),
                                "breakpoint_threshold": float(rag_bp),
                            }
                        )
                    if isinstance(web_result, str) and web_result.startswith("[OK]"):
                        st.success(web_result)
                    elif isinstance(web_result, str) and web_result.startswith("[WARN]"):
                        st.warning(web_result)
                    else:
                        st.error(str(web_result))

        _ingest_col, _stop_col = st.columns([3, 1])
        with _ingest_col:
            _do_ingest = st.button("Ingest files", key="rag_ingest_mode", width="stretch")
        _auto_ingest = bool(st.session_state.get("rag_auto_ingest_pending"))
        if _auto_ingest:
            st.session_state["rag_auto_ingest_pending"] = False
        with _stop_col:
            if st.button("Stop", key="rag_stop_ingest_btn", width="stretch"):
                st.session_state["rag_stop_ingest"] = True

        if _do_ingest or _auto_ingest:
            if not st.session_state.rag_uploaded_files:
                st.warning("No files uploaded yet.")
            elif ingest_rag_paths is None:
                st.error("RAG ingestion unavailable.")
            else:
                _embed_provider_now, _embed_model_now = resolve_provider_and_model(
                    os.getenv("RAG_EMBED_MODEL", "llama_server:")
                )
                if _embed_provider_now == "llama_server":
                    _emb_ok, _emb_msg = _probe_llama_server_embeddings(_embed_model_now or os.getenv("LLAMA_SERVER_MODEL", "local"))
                    if not _emb_ok:
                        _fallback_ollama_model = _pick_fallback_ollama_embedding()
                        if _fallback_ollama_model:
                            st.warning(
                                f"{_emb_msg} Falling back to Ollama embedding model `{_fallback_ollama_model}` for this ingest."
                            )
                            _apply_rag_embed_config("ollama", _fallback_ollama_model, rerun=False)
                            st.session_state["rag_auto_ingest_pending"] = True
                            st.rerun()
                        else:
                            st.error(_emb_msg)
                            st.info(
                                "Pull an Ollama embedding model or restart llama-server with `--embeddings`, then ingest again."
                            )
                            st.stop()
                import threading as _threading
                _stop_ev = _threading.Event()
                st.session_state["rag_stop_ingest"] = False
                _n_files = len(st.session_state.rag_uploaded_files)
                _prog_bar = st.progress(0, text=f"Ingesting 0/{_n_files} files…")
                _prog_status = st.empty()

                def _ingest_progress(step: str, idx: int, total: int, msg: str) -> None:
                    if st.session_state.get("rag_stop_ingest"):
                        _stop_ev.set()
                    if total > 0:
                        _prog_bar.progress(min(idx / max(total, 1), 1.0), text=msg)
                    _prog_status.caption(msg)

                res = ingest_rag_paths(
                    paths=[Path(p) for p in st.session_state.rag_uploaded_files],
                    project=rag_project or "Default",
                    theme=rag_theme,
                    chunking_method=rag_chunk_method,
                    chunk_size=int(rag_cs),
                    chunk_overlap=int(rag_co),
                    breakpoint_threshold=float(rag_bp),
                    on_progress=_ingest_progress,
                    stop_event=_stop_ev,
                )
                _prog_bar.empty()
                _prog_status.empty()

                if res.get("stopped"):
                    st.warning(
                        f"🛑 Ingestion stopped — {res['loaded_files']}/{res['total_files']} files "
                        f"· {res['added_chunks']} chunks indexed before stop."
                    )
                elif res["loaded_files"] > 0:
                    # Preserve current project mode unless user specified a different one
                    active_rag_project = rag_project.strip() or "Default"
                    st.session_state.rag_project_mode = active_rag_project
                    st.session_state.rag_project_filter_selector = active_rag_project
                    _actual_method = res.get("method_used", rag_chunk_method)
                    _semantic_family = ("semantic", "contextual", "late_chunking", "semantic_contextual_late")
                    st.success(
                        f"✅ {res['loaded_files']} files · {res['added_chunks']} chunks "
                        f"[{_actual_method}]"
                        + (" *(fell back from contextual/semantic/late)*" if _actual_method == "recursive" and rag_chunk_method in _semantic_family else "")
                    )
                    # Show only the chunks from this ingest batch so the preview reflects
                    # the files the user just uploaded rather than arbitrary older rows.
                    try:
                        _rag_tools = _get_rag_tools()
                        import chromadb as _pv_chroma
                        _pv_proj = rag_project or "Default"
                        _pv_cname = _rag_tools._get_rag_collection_name(_pv_proj)
                        _pv_client = _pv_chroma.PersistentClient(path=str(Path(__file__).parent / "rag-chroma"))
                        _uploaded_names = {
                            Path(_uploaded_path).name
                            for _uploaded_path in st.session_state.rag_uploaded_files
                        }
                        _preview_chunks: list[dict[str, str]] = []
                        try:
                            _pv_coll = _pv_client.get_collection(name=_pv_cname)
                            _pv_count = _pv_coll.count()
                            if _pv_count > 0:
                                _pv_data = _pv_coll.get(include=["documents", "metadatas"])
                                _pv_docs = _pv_data.get("documents") or []
                                _pv_metas = _pv_data.get("metadatas") or []
                                for _pdoc, _pmeta in zip(_pv_docs, _pv_metas):
                                    _pmeta = _pmeta or {}
                                    _source_name = str(_pmeta.get("source") or "").strip()
                                    _file_name = Path(str(_pmeta.get("file_path") or "")).name
                                    if not (_source_name or _file_name):
                                        continue
                                    if _source_name not in _uploaded_names and _file_name not in _uploaded_names:
                                        continue
                                    if not _pdoc or not _pdoc.strip():
                                        continue
                                    _preview_chunks.append(
                                        {
                                            "source": _source_name or _file_name,
                                            "chunking_method": str(_pmeta.get("chunking_method") or "").strip(),
                                            "content": _pdoc.strip(),
                                        }
                                    )
                                    if len(_preview_chunks) >= 5:
                                        break
                        except Exception:
                            pass  # collection missing or unreadable
                        if _preview_chunks:
                            with st.expander("🔎 Preview stored chunks (verify chunking)", expanded=False):
                                def _chunk_marker_tags(content: str) -> list[str]:
                                    tags: list[str] = []
                                    if "[Context:" in content:
                                        tags.append("Context")
                                    if "[Topics:" in content:
                                        tags.append("Topics")
                                    if "[LateCtx:" in content:
                                        tags.append("LateCtx")
                                    return tags

                                def _marker_badge(label: str) -> str:
                                    styles = {
                                        "Context": "background:#1f7a4f;color:#ffffff;border:1px solid #2fb36d;",
                                        "Topics": "background:#8a5a00;color:#ffffff;border:1px solid #d79b2e;",
                                        "LateCtx": "background:#1558b0;color:#ffffff;border:1px solid #4d8fe8;",
                                        "none": "background:#4a4a4a;color:#ffffff;border:1px solid #777;",
                                    }
                                    style = styles.get(label, styles["none"])
                                    return (
                                        f"<span style=\"display:inline-block;padding:0.14rem 0.5rem;"
                                        f"border-radius:999px;font-size:0.78rem;font-weight:700;"
                                        f"line-height:1.2;{style}\">{label}</span>"
                                    )

                                _has_context = any("[Context:" in c["content"] for c in _preview_chunks)
                                _has_late = any("[LateCtx:" in c["content"] for c in _preview_chunks)
                                _has_topics = any("[Topics:" in c["content"] for c in _preview_chunks)
                                st.markdown(
                                    " ".join([
                                        _marker_badge("Context"),
                                        _marker_badge("Topics"),
                                        _marker_badge("LateCtx"),
                                    ]),
                                    unsafe_allow_html=True,
                                )
                                if _actual_method in ("contextual", "semantic_contextual_late"):
                                    if _has_context:
                                        st.success(
                                            f"✅ `[Context: ...]` prefix found — LLM enrichment active "
                                            f"(`{_context_llm_status.get('model', '?')}`)"
                                        )
                                    elif _has_topics:
                                        st.warning(
                                            "⚠️ No `[Context: ...]` prefix found. "
                                            "Context generation fell back to topic tags for these chunks."
                                        )
                                        _model_used = _context_llm_status.get("model") or _strip_ollama_provider_prefix(
                                            os.environ.get(
                                                "RAG_CONTEXT_LLM_MODEL",
                                                get_agent_model("main"),
                                            )
                                        )
                                        st.caption(f"Model attempted: `{_model_used}`")
                                        _ctx_state = _context_llm_status.get("state", "")
                                        if _ctx_state == "empty-response":
                                            st.caption(
                                                "The model returned no usable text, so ingestion used "
                                                "`[Topics: ...]` fallback tags instead of LLM-generated context."
                                            )
                                        elif _ctx_state == "llm-error":
                                            st.caption(
                                                "The context model raised an error, so ingestion used "
                                                "`[Topics: ...]` fallback tags instead."
                                            )
                                        _ctx_err = _context_llm_status.get("error", "")
                                        if _ctx_err:
                                            st.error(f"**LLM status:** {_ctx_err}")
                                    else:
                                        st.warning("⚠️ No `[Context: ...]` prefix found — LLM enrichment failed.")
                                        _model_used = _context_llm_status.get("model") or _strip_ollama_provider_prefix(
                                            os.environ.get("RAG_CONTEXT_LLM_MODEL",
                                                get_agent_model("main"))
                                        )
                                        st.caption(f"Model attempted: `{_model_used}`")
                                        _ctx_err = _context_llm_status.get("error", "")
                                        if _ctx_err:
                                            st.error(f"**LLM error:** {_ctx_err}")
                                        else:
                                            st.caption(
                                                "No error captured — the model returned no final answer text. "
                                                "Thinking models are still allowed; if they only emit reasoning, "
                                                "RAG uses `[Topics: ...]` fallback tags instead."
                                            )
                                if _actual_method in ("late_chunking", "semantic_contextual_late"):
                                    if _has_late:
                                        st.success("✅ `[LateCtx: ...]` prefix found — late chunking context active")
                                    else:
                                        st.warning("⚠️ No `[LateCtx: ...]` prefix — late context window was empty (single-chunk document?)")
                                for _i, _chunk in enumerate(_preview_chunks[:3], 1):
                                    _label = _chunk["source"]
                                    _method = _chunk["chunking_method"] or "unknown"
                                    _marker_tags = _chunk_marker_tags(_chunk["content"])
                                    st.markdown(f"**Chunk {_i}**")
                                    st.caption(f"Source: `{_label}` · Method: `{_method}`")
                                    st.markdown(
                                        "Markers: "
                                        + " ".join(_marker_badge(tag) for tag in (_marker_tags or ["none"])),
                                        unsafe_allow_html=True,
                                    )
                                    st.code(_chunk["content"][:800], language="text")
                        elif st.session_state.rag_uploaded_files:
                            st.caption("Preview unavailable for the current upload batch.")
                    except Exception:
                        pass
                for fail in res.get("failures", []):
                    st.caption(f"⚠️ Failed: {fail}")

    try:
        rag_projects = _get_rag_projects() if _get_rag_projects else []
    except Exception:
        rag_projects = []
    filter_projects = ["All Projects", "Default"] + [p for p in rag_projects if p != "Default"]
    if st.session_state.rag_project_mode not in filter_projects:
        filter_projects.append(st.session_state.rag_project_mode)
    if st.session_state.get("rag_store_project_filter") and st.session_state["rag_store_project_filter"] not in filter_projects:
        filter_projects.append(st.session_state["rag_store_project_filter"])

    selected_project = (
        st.session_state.get("rag_store_project_filter")
        or st.session_state.rag_project_mode
        or "All Projects"
    )
    st.session_state.rag_project_mode = selected_project

    # Index summary — always visible
    try:
        if selected_project == "All Projects":
            # Merge index data from all projects
            all_files = {}
            total_chunks = 0
            for proj in filter_projects:
                if proj != "All Projects":
                    try:
                        proj_summary = (
                            _get_rag_idx(project=proj)
                            if _get_rag_idx
                            else {"total_chunks": 0, "files": {}}
                        )
                        for src, info in proj_summary.get("files", {}).items():
                            # Prefix file source with project name for clarity
                            prefixed_src = f"{proj}/{src}" if proj != "Default" else src
                            all_files[prefixed_src] = info
                        total_chunks += proj_summary.get("total_chunks", 0)
                    except Exception:
                        pass
            summary = {"total_chunks": total_chunks, "files": all_files}
        else:
            summary = (
                _get_rag_idx(project=selected_project)
                if _get_rag_idx
                else {"total_chunks": 0, "files": {}}
            )
    except Exception:
        summary = {"total_chunks": 0, "files": {}}

    try:
        themes = _get_rag_themes(project=selected_project) if _get_rag_themes else []
    except Exception:
        themes = []

    # Compute store_themes for Stored Documents section (handles "All Projects" case)
    try:
        if selected_project == "All Projects":
            store_themes_set = set()
            for proj in filter_projects:
                if proj != "All Projects":
                    try:
                        proj_themes = _get_rag_themes(project=proj) if _get_rag_themes else []
                        store_themes_set.update(proj_themes)
                    except Exception:
                        pass
            store_themes = sorted(list(store_themes_set))
        else:
            store_themes = themes
    except Exception:
        store_themes = []

    files = summary.get("files", {})
    _idx_error = summary.get("error", "")

    if not files and selected_project != "All Projects":
        fallback_project = None
        fallback_summary: dict[str, Any] | None = None
        for proj in project_options:
            if proj in {"All Projects", selected_project}:
                continue
            try:
                candidate_summary = get_rag_index_summary(project=proj)
            except Exception:
                continue
            candidate_files = candidate_summary.get("files", {})
            if candidate_files:
                fallback_project = proj
                fallback_summary = candidate_summary
                break
        if fallback_project and fallback_summary:
            st.session_state.rag_project_mode = fallback_project
            st.session_state.rag_project_filter_selector = fallback_project
            st.rerun()
        elif selected_project == "Default" and not files and "All Projects" in project_options:
            try:
                all_summary = {}
                all_files: dict[str, dict[str, Any]] = {}
                total_chunks = 0
                for proj in project_options:
                    if proj == "All Projects":
                        continue
                    candidate_summary = get_rag_index_summary(project=proj)
                    total_chunks += candidate_summary.get("total_chunks", 0)
                    for src, info in candidate_summary.get("files", {}).items():
                        prefixed_src = f"{proj}/{src}" if proj != "Default" else src
                        all_files[prefixed_src] = info
                if all_files:
                    st.session_state.rag_project_mode = "All Projects"
                    st.session_state.rag_project_filter_selector = ""
                    st.rerun()
            except Exception:
                pass

    with stored_docs_tab:
        _store_col1, _store_col2 = st.columns(2)
        with _store_col1:
            selected_project = st.selectbox(
                "Project filter",
                options=filter_projects,
                index=filter_projects.index(selected_project),
                key="rag_store_project_filter",
                help="Show stored documents for one project, or **All Projects** to view across all projects.",
            )
        with _store_col2:
            selected_theme_filter = st.selectbox(
                "Theme filter",
                ["All"] + store_themes,
                index=0,
                key="rag_store_theme_filter",
                help="Optionally limit the table to one theme.",
            )

        if st.session_state.rag_project_mode != selected_project:
            st.session_state.rag_project_mode = selected_project
            st.rerun()

        if _idx_error:
            st.error(
                f"**Index corrupted:** {_idx_error}\n\n"
                "Use **Delete Matching Chunks** below with no File or Theme filter, then re-ingest your documents."
            )
        elif files:
            import pandas as pd
            _STORE_CHUNK_LABELS = {
                "semantic": "Semantic",
                "recursive": "Recursive",
                "contextual": "Contextual — Semantic + LLM context (slow)",
                "late_chunking": "Late Chunking — Semantic + neighbour context window",
                "semantic_contextual_late": "Semantic + Contextual + Late — all three combined (slowest)",
            }
            rows = []
            for src, info in sorted(files.items()):
                # Extract project name from source if it's prefixed (e.g., "ProjectName/filename")
                if "/" in src and selected_project == "All Projects":
                    proj_name, file_name = src.split("/", 1)
                elif selected_project == "All Projects":
                    proj_name = "Default"
                    file_name = src
                else:
                    proj_name = selected_project
                    file_name = src
                
                rows.append({
                    "Project": proj_name,
                    "File": file_name,
                    "Theme": ", ".join(sorted(info["themes"])),
                    "Chunks": info["chunks"],
                    "Modalities": ", ".join(
                        f"{name}:{count}" for name, count in sorted(info.get("modalities", {}).items())
                    ) or "text",
                    "Table Extraction": ", ".join(
                        f"{name}:{count}" for name, count in sorted(info.get("table_extraction_methods", {}).items())
                    ) or "—",
                    "Vision Captioned": info.get("vision_captioned_chunks", 0),
                    "Chunking Method": ", ".join(
                        _STORE_CHUNK_LABELS.get(method, method)
                        for method in sorted(info.get("chunking_methods", set()))
                    ) or "—",
                    "Added": info.get("date_added", "—"),
                })
            
            if selected_theme_filter != "All":
                rows = [
                    row for row in rows
                    if selected_theme_filter in [part.strip() for part in row["Theme"].split(",")]
                ]
            st.dataframe(pd.DataFrame(rows), width="stretch", hide_index=True)
        else:
            st.info("No documents indexed yet. Use Ingest to add files.")

        filenames = sorted(files.keys())
        with st.expander("Delete chunks"):
            _delete_entries = []
            for src, info in sorted(files.items()):
                if "/" in src and selected_project == "All Projects":
                    _delete_proj, _delete_file = src.split("/", 1)
                elif selected_project == "All Projects":
                    _delete_proj, _delete_file = "Default", src
                else:
                    _delete_proj, _delete_file = selected_project, src
                _delete_entries.append(
                    {
                        "project": _delete_proj,
                        "file": _delete_file,
                        "themes": sorted(info.get("themes", set())),
                    }
                )

            _project_scope_options = ["Project filter"] + sorted(
                {entry["project"] for entry in _delete_entries if entry["project"]}
            )
            _project_scope_value = st.session_state.get("rag_del_project_scope_mode", "Project filter")
            if _project_scope_value not in _project_scope_options:
                _project_scope_value = "Project filter"

            _file_scope_state = st.session_state.get("rag_del_file_scope_mode", "File filter")
            _theme_scope_state = st.session_state.get("rag_del_theme_scope_mode", "Theme filter")

            del_c1, del_c2, del_c3, del_c4 = st.columns(4)
            with del_c1:
                sel_project_delete = st.selectbox(
                    "Project Scope",
                    _project_scope_options,
                    key="rag_del_project_scope_mode",
                    disabled=not _delete_entries,
                    index=_project_scope_options.index(_project_scope_value),
                    help="Choose the project to delete from. Filtered deletion no longer runs across all projects.",
                )
                if sel_project_delete == "Project filter":
                    sel_project_delete = ""
            with del_c2:
                _theme_entries = [
                    entry
                    for entry in _delete_entries
                    if (not sel_project_delete or entry["project"] == sel_project_delete)
                    and (
                        _file_scope_state == "File filter"
                        or entry["file"] == _file_scope_state
                    )
                ]
                _theme_scope_options = ["Theme filter"] + sorted(
                    {
                        theme
                        for entry in _theme_entries
                        for theme in entry["themes"]
                        if theme
                    }
                )
                _theme_scope_value = _theme_scope_state if _theme_scope_state in _theme_scope_options else "Theme filter"
                sel_theme = st.selectbox(
                    "Theme Scope",
                    _theme_scope_options,
                    key="rag_del_theme_scope_mode",
                    disabled=len(_theme_scope_options) == 1,
                    index=_theme_scope_options.index(_theme_scope_value),
                    help="Choose a theme to narrow deletion, or leave this unset.",
                )
                if sel_theme == "Theme filter":
                    sel_theme = ""
            with del_c3:
                _file_entries = [
                    entry
                    for entry in _delete_entries
                    if (not sel_project_delete or entry["project"] == sel_project_delete)
                    and (
                        not sel_theme
                        or sel_theme in entry["themes"]
                    )
                ]
                _file_scope_options = ["File filter"] + sorted({entry["file"] for entry in _file_entries if entry["file"]})
                _file_scope_value = _file_scope_state if _file_scope_state in _file_scope_options else "File filter"
                sel_file = st.selectbox(
                    "File Scope",
                    _file_scope_options,
                    key="rag_del_file_scope_mode",
                    disabled=len(_file_scope_options) == 1,
                    index=_file_scope_options.index(_file_scope_value),
                    help="Choose a file to narrow deletion, or leave this unset.",
                )
                if sel_file == "File filter":
                    sel_file = ""
                if sel_project_delete and st.button("Delete project", key="rag_del_project_btn", width="stretch"):
                    deleted = (
                        delete_rag_documents(
                            project=sel_project_delete,
                            allow_delete_all=True,
                        )
                        if delete_rag_documents
                        else 0
                    )
                    st.success(f"Deleted {deleted} chunks from project '{sel_project_delete}'")
                    st.rerun()
            with del_c4:
                st.write("")
                if selected_project == "All Projects":
                    st.caption("Select a single project in the table above to enable filtered deletion.")
                elif st.button(
                    "🧹 Delete Matching Chunks", key="rag_delete_filtered",
                    width="stretch", disabled=not files,
                    help=(
                        "Delete chunks from the selected project scope, optionally narrowed "
                        "to a specific theme and/or file."
                    ),
                ):
                    if not sel_project_delete:
                        st.warning("Choose a project scope before deleting matching chunks.")
                    else:
                        allow_delete_all = not sel_file and not sel_theme
                        total_deleted = (
                            delete_rag_documents(
                                project=sel_project_delete,
                                source=sel_file,
                                theme=sel_theme,
                                allow_delete_all=allow_delete_all,
                            )
                            if delete_rag_documents
                            else 0
                        )
                        st.success(
                            f"Deleted {total_deleted} chunks matching the selected filters in project '{sel_project_delete}'"
                        )
                        st.rerun()
            st.caption(
                "Select a project scope first. Theme and file options update dynamically from the available chunks in that project."
            )

        if get_last_rag_query_diagnostics:
            latest_diag = get_last_rag_query_diagnostics()
        else:
            latest_diag = {}
        if latest_diag and latest_diag.get("project", selected_project) in (
            "",
            selected_project,
        ):
            with st.expander("📈 Latest Retrieval Diagnostics", expanded=False):
                _diag_top = {
                    "status": latest_diag.get("status", ""),
                    "project": latest_diag.get("project", ""),
                    "themes": latest_diag.get("themes", []),
                    "mode": latest_diag.get("mode", ""),
                    "top_k": latest_diag.get("top_k", ""),
                    "fetch_k": latest_diag.get("fetch_k", ""),
                    "candidate_count": latest_diag.get("candidate_count", ""),
                    "final_chunk_count": latest_diag.get("final_chunk_count", ""),
                    "selected_files": latest_diag.get("selected_files", []),
                    "reranker": latest_diag.get("reranker", ""),
                }
                st.json(_diag_top, expanded=False)
                _diag_modalities = latest_diag.get("selected_modalities", {})
                if _diag_modalities:
                    st.caption(
                        "Selected modalities: "
                        + ", ".join(f"{name}:{count}" for name, count in sorted(_diag_modalities.items()))
                    )
                _diag_chunking = latest_diag.get("selected_chunking_methods", {})
                if _diag_chunking:
                    st.caption(
                        "Selected chunking methods: "
                        + ", ".join(f"{name}:{count}" for name, count in sorted(_diag_chunking.items()))
                    )
                _diag_table = latest_diag.get("selected_table_extraction_methods", {})
                if _diag_table:
                    st.caption(
                        "Selected table extraction: "
                        + ", ".join(f"{name}:{count}" for name, count in sorted(_diag_table.items()))
                    )
                _remaining_diag = {
                    k: v for k, v in latest_diag.items()
                    if k not in _diag_top
                    and k not in {"selected_modalities", "selected_chunking_methods", "selected_table_extraction_methods"}
                }
                if _remaining_diag:
                    st.markdown("**Raw diagnostics**")
                    st.json(_remaining_diag, expanded=False)

    # ── Retrieval Settings (project, theme, presets, manual controls) ──
    retrieval_selected_project = st.session_state.get("rag_project_filter_selector", "")
    rag_projects = _get_rag_projects() if _get_rag_projects else []
    active_project_chunk_count = None
    if retrieval_selected_project and retrieval_selected_project != "All Projects" and _get_rag_idx:
        try:
            active_project_chunk_count = int(_get_rag_idx(project=retrieval_selected_project).get("total_chunks", 0))
        except Exception:
            active_project_chunk_count = None

    try:
        # When "All Projects" is selected, retrieve themes from all projects
        if retrieval_selected_project == "All Projects":
            all_themes_set = set()
            for proj in filter_projects:
                if proj != "All Projects":
                    try:
                        proj_themes = _get_rag_themes(project=proj) if _get_rag_themes else []
                        all_themes_set.update(proj_themes)
                    except Exception:
                        pass
            global_themes = sorted(list(all_themes_set))
        elif retrieval_selected_project:
            global_themes = (
                _get_rag_themes(project=retrieval_selected_project)
                if _get_rag_themes
                else []
            )
        else:
            global_themes = []
    except Exception:
        global_themes = []

    with retrieval_settings_tab:
        if not rag_projects:
            st.error(
                "RAG index not found or empty. The `rag-chroma` store currently has no chunks, "
                "so retrieval cannot return anything until you ingest documents again."
            )
            if st.session_state.rag_uploaded_files:
                st.button(
                    "Re-ingest files",
                    key="rag_reingest_recommended_btn",
                    use_container_width=True,
                    on_click=_queue_rag_reingest,
                    help="Re-run ingestion for the files already staged in the current session.",
                )
                st.caption("This will reuse the files already staged in the Ingest tab.")
            else:
                st.info("Upload files in Ingest first, then use this button to re-ingest them.")
        elif active_project_chunk_count == 0:
            st.warning(
                f"Project `{retrieval_selected_project}` has no indexed chunks in the current RAG store. "
                "If you deleted `rag-chroma` or changed `RAG_EMBED_MODEL`, re-ingest the documents with the same embedding setup."
            )
            if st.session_state.rag_uploaded_files:
                st.button(
                    "Re-ingest files",
                    key="rag_reingest_recommended_btn",
                    use_container_width=True,
                    on_click=_queue_rag_reingest,
                    help="Re-run ingestion for the files already staged in the current session.",
                )
                st.caption("This will reuse the files already staged in the Ingest tab.")
            else:
                st.info("Upload files in Ingest first, then use this button to re-ingest them.")

        retrieval_model_col1, retrieval_model_col2 = st.columns(2)
        with retrieval_model_col1:
            _render_model_selector(
                "RAG Model", session_key="model_rag", env_key=AGENT_MODEL_ENV_KEYS["ragsub"],
                widget_key="rag_model_selector", allow_inherit=True,
            )
        
        with retrieval_model_col2:
            # Show current embedding configuration
            _current_embed_config = os.getenv("RAG_EMBED_MODEL", "llama_server:").strip()
            _current_embed_provider, _current_embed_model = resolve_provider_and_model(_current_embed_config)
            st.caption("Embedding config")
            st.code(_current_embed_config, language="text")
            st.caption("Set it in Ingest → Provider.")

        # ...existing code...


        # ── Project & Theme scope ──
        _scope_col1, _scope_col2, _scope_col3 = st.columns(3)
        with _scope_col1:
            _project_options = [""] + filter_projects
            _current_project = st.session_state.get("rag_project_filter_selector", "")
            if _current_project not in _project_options:
                _current_project = ""
            st.selectbox(
                "Project",
                options=_project_options,
                key="rag_project_filter_selector",
                format_func=lambda value: "No project filter" if value == "" else value,
                help="Leave blank for no project filter. Select **All Projects** to search across all projects explicitly, or choose one project to scope retrieval.",
            )
        with _scope_col2:
            _stored_themes = st.session_state.get("rag_query_themes", [])
            if _stored_themes == "select_all":
                _stored_themes = []
            _default_themes = [t for t in _stored_themes if t in global_themes]
            _selected = st.multiselect(
                "Themes",
                options=global_themes,
                default=_default_themes,
                key="rag_themes_mode",
                disabled=not global_themes,
                help="Leave blank for no theme filter. Select a project or **All Projects** first to choose themes.",
            )
            st.session_state.rag_query_themes = _selected
            if not global_themes:
                st.caption("Select a project or **All Projects** to choose themes.")
        with _scope_col3:
            _modality_options = ["text", "table", "image"]
            _stored_modalities = [
                modality for modality in st.session_state.get("rag_query_modalities", [])
                if modality in _modality_options
            ]
            _selected_modalities = st.multiselect(
                "Modalities",
                options=_modality_options,
                default=_stored_modalities,
                key="rag_modalities_mode",
                help="Leave blank for all modalities. Select one or more to constrain retrieval to text, table, or image chunks.",
            )
            st.session_state.rag_query_modalities = _selected_modalities

        if st.session_state.get("rag_active_preset") == "📄 File summary":
            st.markdown("---")
            st.subheader("File Summary")
            st.caption("Pick a file and the app will show every stored chunk for that file. You can still ask questions below.")

            file_options = sorted(files.keys())
            if not file_options:
                st.info("No indexed files found for the current project scope.")
                return

            selected_file = st.selectbox(
                "File",
                options=file_options,
                index=file_options.index(st.session_state.rag_file_summary_selected_file)
                if st.session_state.rag_file_summary_selected_file in file_options
                else 0,
                key="rag_file_summary_file_selector",
                help="Choose one stored file to inspect all of its chunks.",
            )
            st.session_state.rag_file_summary_selected_file = selected_file

            if selected_project == "All Projects" and "/" in selected_file:
                file_project, source_file = selected_file.split("/", 1)
            else:
                file_project = selected_project if selected_project != "All Projects" else ""
                source_file = selected_file

            summary_chunks = _list_rag_file_chunks(
                _rag_tools,
                project=file_project,
                file_name=source_file,
            )

            st.caption(
                f"Project: `{file_project or 'Default'}` · File: `{Path(source_file).name}` · "
                f"Chunks found: {len(summary_chunks)}"
            )

            if not summary_chunks:
                st.session_state.rag_file_summary_context = ""
                st.info("No stored chunks were returned for this file.")
            else:
                st.session_state.rag_file_summary_context = "\n\n".join(
                    f"[File Chunk {idx}]\n{str(chunk.get('text') or '').strip()}"
                    for idx, chunk in enumerate(summary_chunks, 1)
                    if str(chunk.get("text") or "").strip()
                )
                for idx, chunk in enumerate(summary_chunks, 1):
                    chunk_text = str(chunk.get("text") or "").strip()
                    chunk_meta = chunk.get("metadata") or {}
                    chunk_source = str(chunk_meta.get("source") or chunk_meta.get("file") or source_file)
                    chunk_page = str(chunk_meta.get("page") or chunk_meta.get("page_number") or "").strip()
                    with st.expander(f"Chunk {idx}", expanded=idx == 1):
                        meta_bits = [f"Source: `{Path(chunk_source).name}`"]
                        if chunk_page:
                            meta_bits.append(f"Page: `{chunk_page}`")
                        st.caption(" · ".join(meta_bits))
                        st.code(chunk_text or "[empty chunk]", language="text")

        st.markdown("---")

        # ── Quick presets ──
        st.caption("**Quick profiles** — click to apply a tuned configuration. The built-in default profile is **💬 Q&A** unless you choose a different profile.")
        _preset_names = list(_RAG_PRESETS.keys())
        _active_p = st.session_state.get("rag_active_preset", "")
        _preset_row1 = st.columns(4)
        _preset_row2 = st.columns(3)
        for _col, _pname in zip(_preset_row1 + _preset_row2, _preset_names):
            _p = _RAG_PRESETS[_pname]
            _is_active = _pname == _active_p
            _btn_label = f"✅ {_pname}" if _is_active else _pname
            _col.button(
                _btn_label,
                key=f"rag_preset_{_pname}",
                on_click=_apply_rag_preset,
                args=((_pname, _p["mode"], _p["top_k"], _p["fetch_k"], _p["max_files"], _p["min_rerank"])),
                help=(
                    f"{_p['desc']}\n\n"
                    f"Mode: **{_p['mode']}** · top_k={_p['top_k']} · fetch_k={_p['fetch_k']} "
                    f"· max_files={_p['max_files']} · min_rerank={_p['min_rerank']}\n\n"
                    f"Tuned for: **{_p['agent_label']}**"
                ),
                width="stretch",
            )
        st.markdown("---")

        ret_col1, ret_col2 = st.columns(2)
        with ret_col1:
            _ret_modes = ["Top-K Globally", "Top-K Per File", "MMR", "Hybrid"]
            st.session_state.rag_retrieval_mode = st.selectbox(
                "Retrieval mode",
                options=_ret_modes,
                index=_ret_modes.index(st.session_state.rag_retrieval_mode),
                help=(
                    "**Top-K Globally** — returns the best chunks across all documents (best for focused Q&A).\n\n"
                    "**Top-K Per File** — returns the best chunks *per document*, then merges (best for synthesis/presentation).\n\n"
                    "**MMR** — Maximal Marginal Relevance: removes near-duplicate chunks to maximise diversity (best for multi-source synthesis).\n\n"
                    "**Hybrid** — dense vector search fused with BM25 keyword search via RRF, then CrossEncoder reranked (best recall for exact terms or rare keywords)."
                ),
            )
        with ret_col2:
            st.session_state.rag_top_k = st.slider(
                "Top-K (final chunks)",
                min_value=1,
                max_value=20,
                step=1,
                key="rag_top_k_slider_v2",
                help=(
                    "Number of chunks kept **after reranking** and passed to the LLM as context.\n\n"
                    "Lower = faster, more focused. Higher = more context, potentially more complete answers.\n\n"
                    "Typical values: Q&A → 4–6 · Academic/Synthesis → 8–12 · Presentation → 8–10."
                ),
            )

        ret_col3, ret_col4 = st.columns(2)
        with ret_col3:
            st.session_state.rag_fetch_k = st.slider(
                "Fetch-K (candidate pool)",
                min_value=10,
                max_value=300,
                step=10,
                key="rag_fetch_k_slider_v2",
                help=(
                    "Number of chunks fetched from the vector store **before reranking**.\n\n"
                    "A larger pool gives the reranker more candidates to choose from, improving recall at the cost of speed.\n\n"
                    "Typical values: 50–100 for quick Q&A · 150–300 for academic or synthesis tasks."
                ),
            )
        with ret_col4:
            _mf_disabled = st.session_state.rag_retrieval_mode == "Top-K Globally"
            st.session_state.rag_max_files = st.slider(
                "Max files",
                min_value=1,
                max_value=10,
                step=1,
                key="rag_max_files_slider_v2",
                disabled=_mf_disabled,
                help=(
                    "Maximum number of source files considered when mode is **Top-K Per File** or **MMR**.\n\n"
                    "Ignored for *Top-K Globally* and *Hybrid* modes.\n\n"
                    "Increase to draw from more documents; decrease to keep answers tightly scoped."
                ),
            )

        st.session_state.rag_min_rerank_score = st.slider(
            "Min rerank score",
            min_value=0.0,
            max_value=1.0,
            step=0.05,
            format="%.2f",
            key="rag_min_rerank_score_slider_v2",
            help=(
                "Minimum CrossEncoder relevance score a chunk must achieve to be included in the final context.\n\n"
                "**0.0** (default) — no filtering, all Top-K chunks are passed to the LLM.\n\n"
                "Raise this (e.g. 0.3–0.6) to discard weakly-relevant chunks and keep only high-confidence matches.\n\n"
                "If you get empty results, lower this value or rephrase the query."
            ),
        )

        # ── Chat history control for RAG sub-agent ──────────────────────
        st.session_state.rag_history_mode = st.selectbox(
            "RAG sub-agent chat history",
            options=["Full history", "Current question only", "No history"],
            index=st.session_state.get("rag_history_mode_index", 2),
            help=(
                "Controls how much of the prior conversation the RAG sub-agent sees.\n\n"
                "**Full history** — sends all prior messages (default in some agents).\n"
                "**Current question only** — only sends your latest question, no prior context.\n"
                "**No history** — runs the RAG query with an empty conversation, so the agent "
                "relies purely on retrieved documents."
            ),
        )

        st.markdown("---")
        _active_preset_name = st.session_state.get("rag_active_preset", "")

        _chat_col1, _chat_col2 = st.columns(2)
        with _chat_col1:
            st.radio(
                "Chat with",
                options=["ragsub", "main"],
                key="rag_chat_target",
                horizontal=True,
                format_func=lambda value: (
                    "RAG sub-agent" if value == "ragsub" else "Main agent"
                ),
                help=(
                    "**Preference hint** (agents may override):\n"
                    "• **RAG sub-agent** — specialized retrieval mode\n"
                    "• **Main agent** — full orchestration & tool access\n\n"
                    "With full DeepAgents autonomy enabled, the agent may choose the optimal path regardless of this hint."
                ),
            )
        with _chat_col2:
            st.markdown("**Performance**")
            st.session_state.disable_agent_cache = st.checkbox(
                "🔴 Force cache clear (NOT recommended)",
                value=st.session_state.disable_agent_cache,
                help="⚠️ Disables all DeepAgents & LangChain caching. Dramatically slower. Only for debugging specific cache-related issues.",
            )
            
        
    st.divider()
    _render_chat_history()
    user_input = st.chat_input("Ask me to analyze data, research topics, write content...")
    if not user_input:
        return

    # Clear cache if disabled
    if st.session_state.disable_agent_cache:
        get_cached_agent_bundle.clear()
        get_cached_data_scientist_bundle.clear()
        get_cached_ragsub_bundle.clear()
        get_cached_specialist_router_bundle.clear()
        get_cached_presenter_bundle.clear()

    def _resolve_active_rag_theme_filters() -> list[str]:
        raw_themes = st.session_state.get("rag_query_themes", [])
        if raw_themes == "select_all" or not raw_themes:
            return []
        if isinstance(raw_themes, str):
            return [raw_themes]
        return [str(theme).strip() for theme in raw_themes if str(theme).strip()]

    def _resolve_active_rag_project_filter() -> str:
        retrieval_project = st.session_state.get("rag_project_filter_selector", "")
        return "" if retrieval_project in ("", "All Projects") else retrieval_project

    def _render_active_rag_scope_label() -> str:
        return "across all projects" if not active_project_filter else f"from **{active_project_filter}**"

    def _render_active_rag_scope_caption() -> str:
        project_label = (
            st.session_state.get("rag_project_filter_selector", "")
            or "No project filter"
        )
        theme_label = ", ".join(active_theme_filters) if active_theme_filters else "No theme filter"
        return (
            f"Scope: **{project_label}** · "
            f"Themes: **{theme_label}** · "
            f"Mode: **{st.session_state.rag_retrieval_mode}** · "
            f"top_k: {st.session_state.rag_top_k} · "
            f"fetch_k: {st.session_state.rag_fetch_k} · "
            f"max_files: {st.session_state.rag_max_files} · "
            f"min_rerank: {st.session_state.rag_min_rerank_score:.2f}"
        )

    def _render_fallback_retrieval_notice() -> None:
        project_label = st.session_state.get("rag_project_filter_selector", "") or "No project filter"
        theme_label = ", ".join(active_theme_filters) if active_theme_filters else "No theme filter"
        st.caption(
            "Fallback retrieval scope: "
            f"project=`{project_label}` | "
            f"themes=`{theme_label}` | "
            f"mode=`{st.session_state.rag_retrieval_mode}` | "
            f"top_k=`{st.session_state.rag_top_k}` | "
            f"fetch_k=`{st.session_state.rag_fetch_k}` | "
            f"max_files=`{st.session_state.rag_max_files}` | "
            f"min_rerank=`{st.session_state.rag_min_rerank_score:.2f}`"
        )

    user_message = user_input.strip()
    active_project_filter = _resolve_active_rag_project_filter()
    active_theme_filters = _resolve_active_rag_theme_filters()

    if st.session_state.rag_uploaded_files:
        valid = [
            _safe_project_relative_path(Path(p))
            for p in st.session_state.rag_uploaded_files
            if Path(p).exists() and _safe_project_relative_path(Path(p))
        ]
        valid_strs = [v for v in valid if v is not None]
        if valid_strs:
            prefix_lines = []
            if active_project_filter:
                prefix_lines.append(f"[RAG PROJECT: {active_project_filter}]")
            prefix_lines.extend([
                f"[RAG UPLOADED FILES: {', '.join(valid_strs)}]",
                "[RAG WORKFLOW: Ingest files first, then retrieve and rerank.]",
                "",
                user_message,
            ])
            user_message = "\n".join(prefix_lines)

    if not st.session_state.rag_uploaded_files:
        if active_project_filter:
            user_message = f"[RAG PROJECT: {active_project_filter}]\n{user_message}"

    if active_theme_filters:
        user_message = f"[RAG THEMES: {', '.join(active_theme_filters)}]\n{user_message}"

    # Inject retrieval parameters (informational — used by RAG SubAgent when it handles retrieval itself)
    user_message = (
        f"[RAG TOP_K: {st.session_state.rag_top_k}]\n"
        f"[RAG FETCH_K: {st.session_state.rag_fetch_k}]\n"
        f"[RAG MODE: {st.session_state.rag_retrieval_mode}]\n"
        f"[RAG MAX_FILES: {st.session_state.rag_max_files}]\n"
        f"[RAG MIN_RERANK_SCORE: {st.session_state.rag_min_rerank_score:.2f}]\n"
        + (
            f"[RAG TASK PRESET: {_active_preset_name} — {_RAG_PRESETS[_active_preset_name]['desc']}]\n"
            if _active_preset_name and _active_preset_name in _RAG_PRESETS
            else ""
        )
        + f"{user_message}"
    )

    if st.session_state.rag_knowledge_only:
        user_message = (
            "[RAG KNOWLEDGE ONLY: true]\n"
            "[RAG RULE: Answer strictly from retrieved RAG context. If insufficient, say so.]\n"
            f"{user_message}"
        )

    # ── Stage 1: Pre-retrieve for shared main-chat grounding ────────────────
    # RAG sub-agent turns do their own retrieval inside the agent, so the
    # preview is only useful when grounding main-chat augmentation.
    _chat_target = st.session_state.get("rag_chat_target", "ragsub")
    _preset_forced_agent = None  # Allow agent full autonomy in routing
    _pre_context = ""
    if _chat_target == "main":
        try:
            _pre_context, _diag = _run_active_rag_retrieval(user_input.strip())
            if _pre_context and not _pre_context.startswith("[WARN]"):
                _reranker_mode = _diag.get("reranker", _diag.get("mode", ""))
                _final_count = _diag.get("final_chunk_count", "?")
                _rerank_count = _diag.get("rerank_count", "?")
                _filtered_count = _diag.get("filtered_count", "?")
                _sel_files = _diag.get("selected_files", [])
                _is_fallback = _reranker_mode == "semantic-fallback"
                _retrieval_fallback = str(_diag.get("retrieval_fallback", "") or "").strip()

                _retrieval_status = st.status(
                    f"✅ Retrieved {_final_count} chunks {_render_active_rag_scope_label()}",
                    expanded=False,
                )
                _summary_parts = [f"**{_final_count}** chunks returned"]
                if _rerank_count != "?":
                    _summary_parts.append(f"{_rerank_count} scored by reranker")
                if _filtered_count != "?":
                    _summary_parts.append(f"{_filtered_count} passed threshold")
                if _sel_files:
                    _summary_parts.append(f"from {len(_sel_files)} files")
                if _retrieval_fallback == "bm25":
                    _summary_parts.append("⚠️ keyword fallback used")
                if _is_fallback:
                    _summary_parts.append("⚠️ reranker filtered all → semantic fallback")
                with _retrieval_status:
                    st.caption(_render_active_rag_scope_caption())
                    st.caption(" · ".join(_summary_parts))

                    if _sel_files:
                        st.caption(f"Files: {', '.join(str(f) for f in _sel_files[:6])}")
                    if _diag.get("dense_candidate_count"):
                        st.caption(
                            f"Hybrid: {_diag.get('dense_candidate_count', 0)} dense + "
                            f"{_diag.get('bm25_result_count', 0)} BM25 → "
                            f"{_diag.get('fused_candidate_count', 0)} fused"
                        )
                _retrieval_status.update(
                    label=(
                        f"✅ Retrieved {_diag.get('final_chunk_count', '?')} chunks "
                        f"{_render_active_rag_scope_label()}"
                    ),
                    state="complete",
                    expanded=False,
                )
                user_message = (
                    f"[RAG CONTEXT — retrieved {'across all projects' if not active_project_filter else f'from project {active_project_filter!r}'} "
                    f"using {st.session_state.rag_retrieval_mode}, top_k={st.session_state.rag_top_k}"
                    f"{', keyword fallback' if str(_diag.get('retrieval_fallback', '') or '').strip() == 'bm25' else ''}]\n"
                    f"{_pre_context}\n"
                    f"[END RAG CONTEXT]\n\n"
                    f"{user_message}"
                )
        except Exception as _exc:
            pass

    # ── Stage 2: Route to the selected chat target ──────────────────────────
    # With full DeepAgents autonomy, allow agents to choose optimal routing
    forced_agent_type = AgentType.RAG_SUB.value if _chat_target == "ragsub" else None
    force_main_agent = _chat_target == "main"
    _chat_execution_block(
        user_message,
        user_input.strip(),
        force_agent_type=forced_agent_type,
        force_main_agent=force_main_agent,
        selected_rag_context=_pre_context,
    )


def _mode_literature() -> None:
    """Literature Review mode — 5-step human-in-the-loop academic pipeline."""
    st.markdown("### 🔬 Literature Review")

    col_m, _ = st.columns([2, 2])
    with col_m:
        _render_model_selector(
            "Synthesis Model", session_key="model_literature", env_key="LITERATURE_MODEL",
            widget_key="lit_model_selector", allow_inherit=True,
        )

    # ── Helper: reset to a given step ────────────────────────────────────────
    def _lr_reset_to(step: str) -> None:
        """Clear session state from *step* onward and rerun."""
        _keys_by_step = {
            "search": [
                "lr_found_papers", "lr_search_query",
                "lr_selected_download", "lr_download_result",
                "lr_files_for_ingest", "lr_extra_files", "lr_ingest_result",
                "lr_last_result", "lr_last_docx",
            ],
            "download": [
                "lr_download_result",
                "lr_files_for_ingest", "lr_extra_files", "lr_ingest_result",
                "lr_last_result", "lr_last_docx",
            ],
            "ingest": [
                "lr_ingest_result", "lr_last_result", "lr_last_docx",
            ],
            "report": [
                "lr_last_docx",
            ],
        }
        for k in _keys_by_step.get(step, []):
            st.session_state.pop(k, None)
        st.rerun()

    # ── Draft save / load helpers ────────────────────────────────────────────
    _DRAFT_DIR = PROJECT_DIR / "literature_drafts"
    _DRAFT_DIR.mkdir(parents=True, exist_ok=True)

    # Keys persisted in a draft (order matters for step detection)
    _LR_PERSIST_KEYS = [
        "lr_found_papers", "lr_search_query",
        "lr_sel_download", "lr_download_result", "lr_sel_ingest",
        "lr_last_result", "lr_last_docx",
    ]

    def _lr_current_step() -> str:
        if st.session_state.get("lr_last_result"):
            return "Step 4 — Report"
        if st.session_state.get("lr_download_result"):
            return "Step 3 — Ingest"
        if st.session_state.get("lr_found_papers"):
            return "Step 2 — Download"
        return "Step 1 — Search"

    def _lr_save_draft(name: str) -> Path:
        """Persist current workflow state to a JSON file."""
        data: dict[str, Any] = {}
        for k in _LR_PERSIST_KEYS:
            v = st.session_state.get(k)
            if v is not None:
                data[k] = v
        data["_saved_at"] = datetime.datetime.now().isoformat()
        data["_step"] = _lr_current_step()
        safe = re.sub(r'[^\w\- ]', '', name).strip().replace(' ', '_')[:60] or "draft"
        path = _DRAFT_DIR / f"{safe}.json"
        path.write_text(json.dumps(data, default=str, ensure_ascii=False), encoding="utf-8")
        return path

    def _lr_load_draft(path: Path) -> None:
        """Restore workflow state from a saved draft."""
        data = json.loads(path.read_text(encoding="utf-8"))
        # Clear all LR keys first
        for k in _LR_PERSIST_KEYS:
            st.session_state.pop(k, None)
        # Restore saved keys
        for k in _LR_PERSIST_KEYS:
            if k in data:
                st.session_state[k] = data[k]
        # Migrate old drafts: lr_search_project → rag_project_mode
        if "lr_search_project" in data:
            st.session_state.rag_project_mode = data["lr_search_project"]

    # ── Restart & Draft toolbar ──────────────────────────────────────────────
    _has_progress = any(
        st.session_state.get(k)
        for k in ("lr_found_papers", "lr_download_result", "lr_last_result")
    )

    with st.expander("💾 Save / Load Draft", expanded=False):
        # ── Save ──
        _save_col, _save_btn_col = st.columns([3, 1])
        with _save_col:
            _draft_name = st.text_input(
                "Draft name",
                value=st.session_state.rag_project_mode,
                key="lr_draft_name",
                placeholder="My research draft",
            )
        with _save_btn_col:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("💾 Save", key="lr_save_draft_btn", width="stretch",
                         disabled=not _has_progress):
                _saved_path = _lr_save_draft(_draft_name or "draft")
                st.success(f"Draft saved: {_saved_path.name}")

        # ── Load ──
        _existing_drafts = sorted(_DRAFT_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True)
        if _existing_drafts:
            st.markdown("**Saved drafts:**")
            for _dp in _existing_drafts:
                try:
                    _meta = json.loads(_dp.read_text(encoding="utf-8"))
                except Exception:
                    _meta = {}
                _step_label = _meta.get("_step", "?")
                _saved_at = _meta.get("_saved_at", "")
                if _saved_at:
                    try:
                        _dt = datetime.datetime.fromisoformat(_saved_at)
                        _saved_at = _dt.strftime("%b %d, %Y %H:%M")
                    except Exception:
                        pass
                _load_col, _del_col = st.columns([5, 1])
                with _load_col:
                    if st.button(
                        f"📂 {_dp.stem}  ·  {_step_label}  ·  {_saved_at}",
                        key=f"lr_load_{_dp.stem}",
                        width="stretch",
                    ):
                        _lr_load_draft(_dp)
                        st.success(f"Draft loaded: {_dp.stem}")
                        st.rerun()
                with _del_col:
                    if st.button("🗑", key=f"lr_del_{_dp.stem}"):
                        _dp.unlink(missing_ok=True)
                        st.rerun()
        else:
            st.caption("No saved drafts yet.")

    if _has_progress:
        if st.button("🔄 Restart from scratch", key="lr_restart_all"):
            _lr_reset_to("search")

    # ══════════════════════════════════════════════════════════════════════════
    # STEP 1 — Search
    # ══════════════════════════════════════════════════════════════════════════
    with st.container(border=True):
        st.markdown("**🔍 Step 1 — Search Papers**")
        lr_query = st.text_area(
            "Research query",
            placeholder="e.g. deep learning for drug discovery",
            key="lr_query_mode",
            height=80,
        )
        col_proj, col_max = st.columns([2, 1])
        with col_proj:
            lr_project = st.text_input(
                "Project name",
                value=st.session_state.rag_project_mode,
                key="lr_project_mode",
                help="Shared with RAG Knowledge Base — documents ingested here are visible on both pages.",
            )
        with col_max:
            lr_max = st.slider("Max papers", 5, 50, 20, 5, key="lr_max_mode")

        _search_col, _stop_search_col = st.columns([3, 1])
        with _search_col:
            _do_search = st.button("🔍 Search Papers", key="lr_search_btn", width="stretch")
        with _stop_search_col:
            if st.button("⏹ Stop", key="lr_stop_search_btn", width="stretch"):
                st.session_state["lr_stop_search"] = True

        if _do_search:
            if not lr_query.strip():
                st.warning("Enter a research query first.")
            else:
                import threading
                stop_event = threading.Event()
                st.session_state["lr_stop_search"] = False

                try:
                    from literature_review.pipeline import search_papers_only
                    import asyncio as _asyncio

                    _bar = st.progress(0, text="Searching academic databases…")
                    _txt = st.empty()

                    def _on_prog(step: str, idx: int, total: int, msg: str) -> None:
                        if st.session_state.get("lr_stop_search"):
                            stop_event.set()
                        if total > 0:
                            _bar.progress(min(idx / max(total, 1), 1.0), text=msg)
                        _txt.caption(msg)

                    try:
                        _loop = _asyncio.get_event_loop()
                        if _loop.is_running():
                            import concurrent.futures as _cf
                            with _cf.ThreadPoolExecutor(max_workers=1) as _pool:
                                _search_res = _pool.submit(
                                    _asyncio.run,
                                    search_papers_only(
                                        query=lr_query.strip(),
                                        max_papers=lr_max,
                                        on_progress=_on_prog,
                                        stop_event=stop_event,
                                    ),
                                ).result(timeout=300)
                        else:
                            _search_res = _loop.run_until_complete(
                                search_papers_only(
                                    query=lr_query.strip(),
                                    max_papers=lr_max,
                                    on_progress=_on_prog,
                                    stop_event=stop_event,
                                )
                            )
                    except Exception as _e:
                        _search_res = None
                        st.error(f"Search error: {_e}")

                    _bar.empty()
                    _txt.empty()

                    if _search_res:
                        if _search_res.get("stopped"):
                            st.warning("🛑 Search stopped by user.")
                        else:
                            st.success(f"✅ Found **{_search_res['papers_found']}** papers")
                        st.session_state["lr_found_papers"] = _search_res.get("paper_records", [])
                        st.session_state["lr_search_query"] = lr_query.strip()
                        # Sync project name to the shared RAG project state
                        st.session_state.rag_project_mode = lr_project.strip() or "Default"
                        # Clear downstream state for a fresh run
                        for _k in ("lr_download_result", "lr_files_for_ingest",
                                    "lr_extra_files", "lr_ingest_result",
                                    "lr_last_result", "lr_last_docx"):
                            st.session_state.pop(_k, None)
                        st.rerun()
                except ImportError:
                    st.error("literature_review module not found. Check installation.")

    # ══════════════════════════════════════════════════════════════════════════
    # STEP 2 — Select & Download
    # ══════════════════════════════════════════════════════════════════════════
    _found = st.session_state.get("lr_found_papers")
    _already_downloaded = st.session_state.get("lr_download_result") is not None

    if _found and _already_downloaded:
        # ── Collapsed summary after download ──
        _dl_res_summary = st.session_state["lr_download_result"]
        with st.container(border=True):
            _n_dl = _dl_res_summary.get('papers_downloaded', 0)
            _n_ab = _dl_res_summary.get('papers_abstract_only', 0)
            _n_sk = _dl_res_summary.get('papers_skipped', 0)
            _icon = "✅" if _n_dl > 0 or _n_ab > 0 else "⚠️"
            st.markdown(
                f"**{_icon} Step 2 — Downloaded** — "
                f"{_n_dl} PDFs · "
                f"{_n_ab} abstract-only · "
                f"{_n_sk} skipped"
            )
            _dl_failures = _dl_res_summary.get("failures") or []
            if _dl_failures:
                with st.expander(f"⚠️ {len(_dl_failures)} download issues", expanded=True):
                    for _f in _dl_failures:
                        st.caption(_f)
            if st.button("↩ Re-select papers", key="lr_redo_download"):
                st.session_state.pop("lr_download_result", None)
                st.session_state.pop("lr_sel_ingest", None)
                st.rerun()

    elif _found:
        # ── Full selection UI (not yet downloaded) ──
        with st.container(border=True):
            st.markdown(f"**📥 Step 2 — Select & Download** — {len(_found)} papers found")

            import pandas as pd
            _rows = []
            for idx, rec in enumerate(_found):
                _rows.append({
                    "#": idx,
                    "Title": rec["title"][:80],
                    "Authors": (rec.get("authors") or "")[:40],
                    "Year": rec.get("year", ""),
                    "Journal": (rec.get("journal") or "")[:30],
                    "Citations": rec.get("citation_count", 0),
                    "OA": "✅" if rec.get("open_access") else "",
                })
            _df = pd.DataFrame(_rows)

            with st.expander(f"📄 {len(_found)} papers from search", expanded=True):
                st.dataframe(_df, width="stretch", hide_index=True)

            # Selection checkboxes
            st.caption("Select papers to download (uncheck to skip):")
            _sel_key = "lr_sel_download"
            if _sel_key not in st.session_state:
                st.session_state[_sel_key] = list(range(len(_found)))

            _sel_all_col, _sel_none_col, _ = st.columns([1, 1, 4])
            with _sel_all_col:
                if st.button("Select All", key="lr_sel_all_dl"):
                    st.session_state[_sel_key] = list(range(len(_found)))
                    st.rerun()
            with _sel_none_col:
                if st.button("Select None", key="lr_sel_none_dl"):
                    st.session_state[_sel_key] = []
                    st.rerun()

            _selected_indices: list[int] = []
            _cols_per_row = 2
            for row_start in range(0, len(_found), _cols_per_row):
                _cols = st.columns(_cols_per_row)
                for col_offset, _c in enumerate(_cols):
                    _i = row_start + col_offset
                    if _i >= len(_found):
                        break
                    with _c:
                        _checked = st.checkbox(
                            f"{_found[_i]['title'][:65]}",
                            value=_i in st.session_state[_sel_key],
                            key=f"lr_dl_cb_{_i}",
                        )
                        if _checked:
                            _selected_indices.append(_i)
            st.session_state[_sel_key] = _selected_indices

            _n_sel = len(_selected_indices)
            _dl_col, _back_col = st.columns([3, 1])
            with _dl_col:
                _do_dl = st.button(
                    f"📥 Download {_n_sel} papers",
                    key="lr_download_btn",
                    type="primary",
                    width="stretch",
                    disabled=_n_sel == 0,
                )
            with _back_col:
                if st.button("🔍 New Search", key="lr_back_to_search", width="stretch"):
                    _lr_reset_to("search")

            if _do_dl:
                _sel_papers = [_found[i] for i in _selected_indices]

                import threading
                stop_event = threading.Event()
                st.session_state["lr_stop_download"] = False

                try:
                    from literature_review.pipeline import download_selected_papers
                    import asyncio as _asyncio

                    _bar = st.progress(0, text="Starting downloads…")
                    _txt = st.empty()
                    _stop_dl_btn = st.empty()
                    if _stop_dl_btn.button("⏹ Stop Download", key="lr_stop_dl_btn"):
                        st.session_state["lr_stop_download"] = True

                    def _on_dl_prog(step: str, idx: int, total: int, msg: str) -> None:
                        if st.session_state.get("lr_stop_download"):
                            stop_event.set()
                        if total > 0:
                            _bar.progress(min(idx / max(total, 1), 1.0), text=msg)
                        _txt.caption(msg)

                    try:
                        _loop = _asyncio.get_event_loop()
                        if _loop.is_running():
                            import concurrent.futures as _cf
                            with _cf.ThreadPoolExecutor(max_workers=1) as _pool:
                                _dl_res = _pool.submit(
                                    _asyncio.run,
                                    download_selected_papers(
                                        paper_records=_sel_papers,
                                        on_progress=_on_dl_prog,
                                        stop_event=stop_event,
                                    ),
                                ).result(timeout=600)
                        else:
                            _dl_res = _loop.run_until_complete(
                                download_selected_papers(
                                    paper_records=_sel_papers,
                                    on_progress=_on_dl_prog,
                                    stop_event=stop_event,
                                )
                            )
                    except Exception as _e:
                        _dl_res = None
                        st.error(f"Download error: {_e}")

                    _bar.empty()
                    _txt.empty()
                    _stop_dl_btn.empty()

                    if _dl_res:
                        if _dl_res.get("stopped"):
                            st.warning("🛑 Download stopped by user.")
                        else:
                            st.success(
                                f"✅ **{_dl_res['papers_downloaded']}** PDFs · "
                                f"**{_dl_res['papers_abstract_only']}** abstract-only · "
                                f"**{_dl_res['papers_skipped']}** skipped"
                            )
                        st.session_state["lr_download_result"] = _dl_res
                        # Debug: show file count in toast
                        _n_paths = len(_dl_res.get("paths_to_ingest") or [])
                        st.toast(f"📁 {_n_paths} files stored for ingestion")
                        if _dl_res.get("failures"):
                            with st.expander(f"⚠️ {len(_dl_res['failures'])} warnings", expanded=True):
                                for fail in _dl_res["failures"]:
                                    st.caption(fail)
                        st.rerun()
                except ImportError:
                    st.error("literature_review module not found.")

    # ══════════════════════════════════════════════════════════════════════════
    # STEP 3 — Select & Ingest
    # ══════════════════════════════════════════════════════════════════════════
    _dl_result = st.session_state.get("lr_download_result")
    if _dl_result is not None:
        with st.container(border=True):
            _dl_paths = _dl_result.get("paths_to_ingest") or []
            _dl_metas = _dl_result.get("extra_metas") or []
            _dl_records = _dl_result.get("paper_records", [])

            # Fallback: rebuild paths from paper_records if paths_to_ingest is empty
            if not _dl_paths and _dl_records:
                for _rec in _dl_records:
                    _lp = _rec.get("local_path", "")
                    if _lp and Path(_lp).exists():
                        _dl_paths.append(_lp)
                        _dl_metas.append({
                            "doi": _rec.get("doi", ""),
                            "journal": _rec.get("journal", ""),
                            "pub_date": str(_rec.get("year", "")),
                            "citation_count": str(_rec.get("citation_count", 0)),
                            "pdf_method": _rec.get("method", ""),
                            "paper_abstract": (_rec.get("abstract") or "")[:500],
                        })

            _n_dl_files = len(_dl_paths)
            _n_dl_stats = (
                _dl_result.get("papers_downloaded", 0)
                + _dl_result.get("papers_abstract_only", 0)
            )

            st.markdown(
                f"**📂 Step 3 — Select & Ingest** — "
                f"{max(_n_dl_files, _n_dl_stats)} files ready"
            )
            st.caption(
                f"Project: **{st.session_state.rag_project_mode}** · "
                "Select which files to ingest into the RAG knowledge base."
            )

            # --- File selection checkboxes ---
            _ingest_sel_key = "lr_sel_ingest"
            if _ingest_sel_key not in st.session_state:
                st.session_state[_ingest_sel_key] = list(range(len(_dl_paths)))

            _sel_all_i_col, _sel_none_i_col, _ = st.columns([1, 1, 4])
            with _sel_all_i_col:
                if st.button("Select All", key="lr_sel_all_ing"):
                    st.session_state[_ingest_sel_key] = list(range(len(_dl_paths)))
                    st.rerun()
            with _sel_none_i_col:
                if st.button("Select None", key="lr_sel_none_ing"):
                    st.session_state[_ingest_sel_key] = []
                    st.rerun()

            _ing_selected: list[int] = []
            for fi, fpath in enumerate(_dl_paths):
                _fname = Path(fpath).name
                _checked = st.checkbox(
                    _fname,
                    value=fi in st.session_state[_ingest_sel_key],
                    key=f"lr_ing_cb_{fi}",
                )
                if _checked:
                    _ing_selected.append(fi)
            st.session_state[_ingest_sel_key] = _ing_selected

            # --- Upload extra files ---
            st.markdown("**➕ Add extra files** (optional)")
            _extra_files = st.file_uploader(
                "Upload additional documents to include in ingestion",
                type=["pdf", "docx", "doc", "odt", "txt", "md", "csv", "json"],
                accept_multiple_files=True,
                key="lr_extra_upload",
            )
            _extra_saved: list[str] = []
            if _extra_files:
                _upload_dir = PROJECT_DIR / "temp_litreview"
                _upload_dir.mkdir(parents=True, exist_ok=True)
                for uf in _extra_files:
                    _dest = _upload_dir / uf.name
                    _dest.write_bytes(uf.getvalue())
                    _extra_saved.append(str(_dest))
                st.caption(f"{len(_extra_saved)} extra file(s) staged for ingestion")

            # --- Chunking & Embedding Settings ---
            with st.expander("⚙️ Chunking & Embedding Settings", expanded=False):
                _LR_CHUNK_LABELS = {
                    "semantic": "Semantic (recommended)",
                    "recursive": "Recursive",
                    "contextual": "Contextual — Semantic + LLM context (slow)",
                    "late_chunking": "Late Chunking — Semantic + neighbour context window",
                    "semantic_contextual_late": "Semantic + Contextual + Late — all three combined (slowest)",
                }
                _lr_chunk_method = st.selectbox(
                    "Chunking",
                    list(_LR_CHUNK_LABELS.keys()),
                    index=0,
                    format_func=lambda x: _LR_CHUNK_LABELS.get(x, x),
                    key="lr_chunk_method",
                )
                _lr_needs_semantic = _lr_chunk_method in (
                    "semantic", "contextual", "late_chunking", "semantic_contextual_late"
                )
                _lr_needs_llm = _lr_chunk_method in ("contextual", "semantic_contextual_late")
                if _lr_needs_semantic:
                    _lr_bp = st.slider(
                        "Breakpoint threshold", 50, 99, 95, 1, key="lr_bp",
                    )
                    _lr_cs, _lr_co = 1500, 300
                    if _lr_needs_llm:
                        _lr_ctx_model = _render_model_selector(
                            "Context LLM",
                            session_key="model_lr_context",
                            env_key="RAG_CONTEXT_LLM_MODEL",
                            widget_key="lr_context_model_selector",
                            allow_inherit=True,
                            provider_override="all",
                        )
                        _lr_ctx_label = _lr_ctx_model or os.environ.get(
                            AGENT_MODEL_ENV_KEYS["main"], get_agent_model("main"),
                        )
                        st.caption(
                            f"⚠️ Each chunk calls the LLM (`{_lr_ctx_label}`) — slow."
                        )
                    if _lr_chunk_method in ("late_chunking", "semantic_contextual_late"):
                        st.caption(
                            "ℹ️ **Late chunking**: each chunk is enriched with sentences "
                            "from neighbouring chunks (±2 window) before embedding."
                        )
                else:
                    _lr_bp = 95
                    _lr_cs_col, _lr_co_col = st.columns(2)
                    with _lr_cs_col:
                        _lr_cs = st.number_input(
                            "Chunk size", 100, 4000, 1500, 100, key="lr_cs",
                        )
                    with _lr_co_col:
                        _lr_co = st.number_input(
                            "Overlap", 0, 1000, 300, 50, key="lr_co",
                        )
                st.caption(
                    f"**{_lr_chunk_method}** · size={_lr_cs} · overlap={_lr_co} · "
                    f"breakpoint={_lr_bp}"
                )

            # --- Action buttons ---
            _total_to_ingest = len(_ing_selected) + len(_extra_saved)
            _ing_btn_col, _stop_ing_col, _skip_ing_col, _back_dl_col = st.columns([2, 1, 1, 1])
            with _ing_btn_col:
                _do_ingest = st.button(
                    f"⚡ Ingest {_total_to_ingest} files",
                    key="lr_ingest_btn",
                    type="primary",
                    width="stretch",
                    disabled=_total_to_ingest == 0,
                )
            with _stop_ing_col:
                if st.button("⏹ Stop", key="lr_stop_ingest_btn", width="stretch"):
                    st.session_state["lr_stop_ingest"] = True
            with _skip_ing_col:
                _do_skip = st.button("⏭ Skip", key="lr_skip_ingest_btn", width="stretch")
            with _back_dl_col:
                if st.button("↩ Back", key="lr_back_to_dl", width="stretch"):
                    _lr_reset_to("download")

            if _do_skip:
                _all_records = _dl_records
                st.session_state["lr_last_result"] = {
                    "papers_found": len(st.session_state.get("lr_found_papers", [])),
                    "paper_records": _all_records,
                    "ingested_chunks": 0,
                    "project": st.session_state.rag_project_mode,
                    "theme": "Literature Review",
                    "query": st.session_state.get("lr_search_query", ""),
                }
                st.info("⏭ Ingestion skipped — papers available for report but not indexed.")
                st.rerun()

            if _do_ingest:
                import threading
                stop_event = threading.Event()
                st.session_state["lr_stop_ingest"] = False

                # Build final lists: selected downloaded + extra uploads
                _final_paths = [_dl_paths[i] for i in _ing_selected] + _extra_saved
                _final_metas = [_dl_metas[i] for i in _ing_selected if i < len(_dl_metas)]
                # For extra files, add minimal metadata
                for _ep in _extra_saved:
                    _final_metas.append({"doi": "", "journal": "", "pub_date": "",
                                         "citation_count": "0", "pdf_method": "manual_upload",
                                         "paper_abstract": ""})

                try:
                    from literature_review.pipeline import ingest_papers_pipeline

                    _bar = st.progress(0, text="Starting ingestion…")
                    _txt = st.empty()

                    def _on_ing_prog(step: str, idx: int, total: int, msg: str) -> None:
                        if st.session_state.get("lr_stop_ingest"):
                            stop_event.set()
                        if total > 0:
                            _bar.progress(min(idx / max(total, 1), 1.0), text=msg)
                        _txt.caption(msg)

                    _ing_res = ingest_papers_pipeline(
                        paths_to_ingest=_final_paths,
                        extra_metas=_final_metas,
                        project=st.session_state.rag_project_mode,
                        theme="Literature Review",
                        on_progress=_on_ing_prog,
                        stop_event=stop_event,
                        chunking_method=_lr_chunk_method,
                        chunk_size=int(_lr_cs),
                        chunk_overlap=int(_lr_co),
                        breakpoint_threshold=float(_lr_bp),
                    )

                    _bar.empty()
                    _txt.empty()

                    if _ing_res.get("stopped"):
                        st.warning(
                            f"🛑 Ingestion stopped — "
                            f"{_ing_res['ingested_chunks']} chunks indexed before stop."
                        )
                    else:
                        st.success(
                            f"✅ Ingested **{_ing_res['ingested_chunks']}** chunks from "
                            f"**{_total_to_ingest}** files"
                        )

                    if _ing_res.get("failures"):
                        with st.expander(f"⚠️ {len(_ing_res['failures'])} warnings"):
                            for fail in _ing_res["failures"]:
                                st.caption(fail)

                    st.session_state["lr_last_result"] = {
                        "papers_found": len(st.session_state.get("lr_found_papers", [])),
                        "paper_records": _dl_records,
                        "ingested_chunks": _ing_res["ingested_chunks"],
                        "project": st.session_state.rag_project_mode,
                        "theme": "Literature Review",
                        "query": st.session_state.get("lr_search_query", ""),
                        "failures": _dl_result.get("failures", []) + _ing_res.get("failures", []),
                    }
                    st.rerun()
                except ImportError:
                    st.error("literature_review module not found.")

    # ══════════════════════════════════════════════════════════════════════════
    # STEP 4 — Report
    # ══════════════════════════════════════════════════════════════════════════
    lr_last = st.session_state.get("lr_last_result")
    if lr_last:
        with st.container(border=True):
            _lr_proj = lr_last.get('project', 'Literature Review')
            _lr_model = st.session_state.get("model_literature", "") or "(no model set)"
            st.markdown(
                f"**📄 Step 4 — Report** — *{lr_last.get('papers_found', '?')} papers · "
                f"{lr_last.get('ingested_chunks', 0)} chunks · "
                f"project: **{_lr_proj}***"
            )
            st.caption(f"RAG project: `{_lr_proj}` · Synthesis model: `{_lr_model}`")

            # Live RAG status check
            try:
                _rag_summary = _get_rag_tools().get_rag_index_summary(project=_lr_proj)
                _rag_chunks = _rag_summary.get("total_chunks", 0)
                _rag_files = len(_rag_summary.get("files", {}))
                if _rag_chunks == 0:
                    st.warning(
                        f"⚠️ RAG database has **0 chunks** for project `{_lr_proj}`. "
                        "The report will use metadata fallback only. "
                        "Go back to **Step 3** and run ingestion to enable RAG-based synthesis.",
                        icon="⚠️",
                    )
                else:
                    st.info(
                        f"✅ RAG database: **{_rag_chunks} chunks** from **{_rag_files} files** "
                        f"indexed for project `{_lr_proj}`.",
                        icon="✅",
                    )
            except Exception:
                pass

            lr_synthesis = st.text_area(
                "Synthesis text (optional — paste from chat below)",
                placeholder="## Findings\n...\n## Consensus\n...",
                key="lr_synthesis_mode",
                height=120,
            )
            col_gen, col_dl, col_redo = st.columns([2, 2, 1])
            with col_gen:
                if st.button("📄 Generate .docx", key="lr_gen_mode", width="stretch"):
                    try:
                        from literature_review.report_generator import generate_docx_report
                        _out_dir = PROJECT_DIR / "literature_reports"
                        _out_dir.mkdir(exist_ok=True)
                        _out_path = _out_dir / f"{lr_last.get('project', 'LR').replace(' ', '_')[:40]}_review.docx"
                        _synth_bar = st.progress(0, text="🔬 Querying RAG database…")
                        _synth_txt = st.empty()
                        _section_labels = ["findings", "consensus", "debates", "gaps", "conclusion"]
                        _total_sections = len(_section_labels)

                        def _on_synth_prog(step: str, idx: int, total: int, msg: str) -> None:
                            _synth_bar.progress(
                                min(idx / max(total, 1), 1.0),
                                text=f"🔬 {msg}",
                            )
                            _synth_txt.caption(msg)

                        _docx_path = generate_docx_report(
                            project=lr_last.get("project", "Literature Review"),
                            query=lr_last.get("query", ""),
                            synthesis_text=lr_synthesis or "",
                            paper_records=lr_last.get("paper_records", []),
                            output_path=_out_path,
                            pipeline_stats=lr_last,
                            model_name=st.session_state.get("model_literature", ""),
                            on_progress=_on_synth_prog,
                        )
                        _synth_bar.empty()
                        _synth_txt.empty()
                        st.session_state["lr_last_docx"] = str(_docx_path)
                        st.success(f"Report generated: {_docx_path.name}")
                    except Exception as _e:
                        st.error(f"Report generation failed: {_e}")
            with col_dl:
                last_docx = st.session_state.get("lr_last_docx")
                if last_docx and Path(last_docx).exists():
                    with open(last_docx, "rb") as _f:
                        st.download_button(
                            label="⬇️ Download .docx",
                            data=_f,
                            file_name=Path(last_docx).name,
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key="lr_dl_mode",
                            width="stretch",
                        )
            with col_redo:
                if st.button("🔄 Redo", key="lr_redo_btn", width="stretch"):
                    _lr_reset_to("search")

    st.divider()
    st.caption("💬 Chat below to synthesise findings or ask follow-up questions about indexed papers.")
    _render_chat_history()

    user_input = st.chat_input("Ask about indexed literature or request a synthesis...")
    if not user_input:
        return

    user_message = user_input.strip()
    proj = st.session_state.rag_project_mode
    if proj:
        user_message = (
            f"[RAG PROJECT: {proj}]\n"
            f"[MODE: Literature synthesis — use RAG retrieval for this project.]\n"
            f"{user_message}"
        )

    if st.session_state.get("rag_active_preset") == "📄 File summary":
        file_context = str(st.session_state.get("rag_file_summary_context") or "").strip()
        selected_file = str(st.session_state.get("rag_file_summary_selected_file") or "").strip()
        if file_context and selected_file:
            user_message = (
                f"[FILE SUMMARY MODE]\n"
                f"[FILE: {selected_file}]\n"
                f"[GROUNDING: Answer only from the selected file chunks below. If the answer is not in the chunks, say so clearly.]\n"
                f"{file_context}\n\n{user_message}"
            )

    # Full DeepAgents autonomy: no forced routing
    _chat_execution_block(user_message, user_input.strip(), force_agent_type=None)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    """Main Streamlit app — multi-mode dashboard."""

    active_mode = render_sidebar(
        st=st,
        session_state=st.session_state,
        project_dir=PROJECT_DIR,
        load_chat_sessions=_load_chat_sessions,
        save_current_session=_save_current_session,
        new_chat_session=_new_chat_session,
        switch_to_session=_switch_to_session,
        delete_session=_delete_session,
    )

    # ── MAIN CONTENT — mode-specific ──────────────────────────────────────────
    if active_mode == "main":
        _mode_main()
    elif active_mode == "scan_to_text":
        _mode_scan_to_text()
    elif active_mode == "rag_databases":
        _mode_rag_databases()
    elif active_mode == "rag":
        _mode_rag()
    elif active_mode == "literature":
        _mode_literature()
    elif active_mode == "data_analysis":
        _mode_data_analysis()


if __name__ == "__main__":
    main()
