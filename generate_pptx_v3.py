import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(output_path):
    prs = Presentation()

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Understanding AI Fundamentals"
    subtitle.text = "Gain a foundational understanding of what Artificial Intelligence is"

    # 2. Table of Contents Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"
    content = slide.placeholders[1]
    content.text = (
        "• Introduction\n"
        "• Evolution and History\n"
        "• Levels and Types of AI\n"
        "• Core Components\n"
        "• Summary"
    )

    # 3. Introduction Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction to AI"
    content = slide.placeholders[1]
    content.text = (
        "Artificial Intelligence (AI) refers to the field of research and development focused on creating machines and systems capable of performing tasks that typically require human intelligence, such as reasoning, learning, and problem-solving.\n\n"
        "Three Foundational Elements:\n"
        "• High-Quality Data: The fuel that allows models to learn patterns\n"
        "• Algorithm Selection: The mathematical instructions that process the data\n"
        "• Structured Output: The organized result or decision produced by the system"
    )

    # 4. Evolution and History Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Evolution and History of AI"
    content = slide.placeholders[1]
    content.text = (
        "• 1950: The Theoretical Spark (Alan Turing)\n"
        "• 1956: The Big Bang (Dartmouth Workshop)\n"
        "• Deep Learning Milestone: Development of large neural networks\n"
        "• 2017–Present: The Generative Revolution (Transformer architectures & LLMs)"
    )

    # 5. Levels and Types of AI Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
ly_text = (
        "Levels of Intelligence:\n"
        "• Artificial Narrow Intelligence (ANI)\n"
        "• Artificial General Intelligence (AGI)\n"
        "• Artificial Superintelligence (ASI)\n\n"
        "Functional Types:\n"
        "• Predictive AI: Forecasting future outcomes\n"
        "• Generative AI (GenAI): Creating new content\n"
        "• Agentic AI: Completing complex workflows"
    )
    title.text = "Levels and Types of AI"
    content = slide.placeholders[1]
    content.text = (
        "Levels of Intelligence:\n"
        "• Artificial Narrow Intelligence (ANI)\n"
        "• Artificial General Intelligence (AGI)\n"
        "• Artificial Superintelligence (ASI)\n\n"
        "Functional Types:\n"
        "• Predictive AI: Forecasting future outcomes\n"
        "• Generative AI (GenAI): Creating new content\n"
        "• Agentic AI: Completing complex workflows"
    )

    # 6. Core Components Slide 1 (ML & DL)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Components: ML & DL"
    content = slide.placeholders[1]
    content.text = (
        "Machine Learning (ML):\n"
        "• Supervised Learning (Labeled data)\n"
        "• Unsupervised Learning (Unlabeled data)\n"
        "• Reinforcement Learning (Trial and error)\n"
        "• Examples: Netflix recommendations, Fraud detection\n\n"
        "Deep Learning (DL):\n"
        "• Uses multi-layered Neural Networks\n"
        "• Examples: Advanced image recognition, NLP"
    )

    # 7. Core Components Slide 2 (NLP & CV)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Components: NLP & CV"
    content = slide.placeholders[1]
    content.text = (
        "Natural Language Processing (NLP):\n"
        "• Understanding, interpreting, and generating human language\n"
        "• Examples: Chatbots, Sentiment analysis, Translation\n\n"
        "Computer Vision (CV):\n"
        "• Interpreting visual information from the world\n"
        "• Examples: Facial recognition, Automated quality control"
    )

    # 8. Core Components Slide 3 (GenAI)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Components: Generative AI"
    content = slide.placeholders[1]
    content.text = (
        "Generative AI (GenAI):\n"
        "• A subset of AI focused on the creation of new, original content\n"
        "• Examples:\n"
        "  - Automated content creation\n"
        "  - Graphic design\n"
        "  - Personalized marketing messages"
    )

    # 9. Summary Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Summary"
    content = slide.placeholders[1]
    content.text = (
        "• AI has evolved from early theories into a transformative force.\n"
        "• Key drivers: High-quality data, algorithms, and structured output.\n"
        "• Core technologies: ML, DL, NLP, CV, and GenAI.\n"
        "• Understanding these fundamentals is essential for innovation."
    )

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    output_file = "AI_Fundamentals_Presentation.pptx"
    create_presentation(output_file)
