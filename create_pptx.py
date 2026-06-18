import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Add a slide with a title and content
    slide_layout = prs.slide_layers[1] # 1 is Title and Content
    # Wait, slide_layouts is an indexable object. 
    # Let's use a safer way to get the layout.
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Introduction to KoboToolbox"
    
    content = slide.placeholders[1]
    content.text = "A platform for creating structured surveys, questionnaires, and quizzes"
    
    p = content.text_frame.add_paragraph()
    p.text = "Utilizes the XLSForm standard"
    p.level = 0
    
    p = content.text_frame.add_paragraph()
    p.text = "Key Features:"
    p.level = 0
    
    p = content.text_frame.add_paragraph()
    p.text = "Structured Organization (begin_group)"
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "Complex Logic and Skip Logic (select_multiple, relevant)"
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "Automated Calculations and Scoring (calculate, if-statements)"
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "Ideal for quantitative research and statistics"
    p.level = 0

    prs.save('KoboToolbox_Presentation.pptx')
    print("Presentation created successfully.")

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"An error occurred: {e}")
