"""Presentation generation tools for creating PowerPoint and PDF slide decks."""

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT


def create_pptx_presentation(
    slides_data: list[dict[str, Any]], output_path: str
) -> str:
    """Create a PowerPoint presentation from slide data.
    
    Args:
        slides_data: List of slide dictionaries with keys:
            - title: Slide title
            - content: List of bullet points or main text
            - slide_type: "title", "content", or "closing"
        output_path: Path to save the PPTX file
        
    Returns:
        Path to the created file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        title_shape.text_frame.paragraphs[0].font.size = Pt(54)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        content = slide_data.get("content", [])
        if isinstance(content, list):
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = str(point)
                p.level = 0
                p.font.size = Pt(32)
        else:
            body_shape = slide.placeholders[1]
            body_shape.text = str(content)
    
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    
    return str(output_file)


def create_pdf_presentation(
    slides_data: list[dict[str, Any]], output_path: str
) -> str:
    """Create a PDF presentation from slide data.
    
    Args:
        slides_data: List of slide dictionaries with keys:
            - title: Slide title
            - content: List of bullet points or main text
        output_path: Path to save the PDF file
        
    Returns:
        Path to the created file
    """
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    
    doc = SimpleDocTemplate(str(output_file), pagesize=letter)
    story = []
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=RGBColor(0, 51, 102),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    content_style = ParagraphStyle(
        'CustomContent',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=12,
        alignment=TA_LEFT,
        leading=18
    )
    
    for i, slide_data in enumerate(slides_data):
        # Add title
        title = slide_data.get("title", "")
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Add content
        content = slide_data.get("content", [])
        if isinstance(content, list):
            for point in content:
                bullet_text = f\"- {point}\"
                story.append(Paragraph(bullet_text, content_style))
        else:
            story.append(Paragraph(str(content), content_style))
        
        # Add page break between slides (except last one)
        if i < len(slides_data) - 1:
            story.append(PageBreak())
    
    doc.build(story)
    return str(output_file)


def create_presentation_from_json(
    json_data: str, output_dir: str = "./presentations"
) -> dict[str, str]:
    """Create both PPTX and PDF presentations from JSON slide data.
    
    Args:
        json_data: JSON string containing slides array
        output_dir: Directory to save files
        
    Returns:
        Dictionary with keys 'pptx' and 'pdf' containing file paths
    """
    try:
        slides_data = json.loads(json_data)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON data: {e}")
    
    if not isinstance(slides_data, list):
        raise ValueError("JSON data must be an array of slide objects")
    
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    pptx_path = output_dir / "presentation.pptx"
    pdf_path = output_dir / "presentation.pdf"
    
    pptx_file = create_pptx_presentation(slides_data, str(pptx_path))
    pdf_file = create_pdf_presentation(slides_data, str(pdf_path))
    
    return {
        "pptx": pptx_file,
        "pdf": pdf_file,
        "slides_count": len(slides_data)
    }
