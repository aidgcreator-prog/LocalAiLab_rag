import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(output_path):
    prs = Presentation()

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Understanding AI Fundamentals"
    subtitle.text = "Gain a foundational understanding of what Artificial Intelligence is"

    # 2. Table of Contents Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"
    content = slide.placeholders[1]
    content.text = (
        "• Introduction\n"
        "• Evolution and History\n"
        "• Levels and Types of AI\n"
        "• Core Components\n"
        "• Summary"
    )

    # 3. Introduction Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction to AI"
    content = slide.placeholders[1]
    content.text = (
        "Artificial Intelligence (AI) refers to the field of research and development focused on creating machines and systems capable of performing tasks that typically require human intelligence, such as reasoning, learning, and problem-solving.\n\n"
        "Three Foundational Elements:\n"
        "• High-Quality Data: The fuel that allows models to learn patterns\n"
        "• Algorithm Selection: The mathematical instructions that process the data\n"
        "• Structured Output: The organized result or decision produced by the system"
    )

    # 4. Evolution and History Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Evolution and History of AI"
    content = slide.placeholders[1]
    content.text = (
        "• 1950: The Theoretical Spark (Alan Turing)\n"
        "• 1956: The Big Bang (Dartmouth Workshop)\n"
        "• Deep Learning Milestone: Development of large neural networks\n"
        "• 2017–Present: The Generative Revolution (Transformer architectures & LLMs)"
    )

    # 5. Levels and Types of AI Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Levels and Types of AI"
    content = slide.placeholders[1]
    content.text = (
        "Levels of Intelligence:\n"
        "• Artificial Narrow Intelligence (ANI)\n"
        "• Artificial General Intelligence (AGI)\n"
        "• Artificial Superintelligence (ASI)\n\n"
        "Functional Types:\n"
        "• Predictive AI: Forecasting future outcomes\n"
        "• Generative AI (GenAI): Creating new content\n"
        "• Agentic AI: Completing complex workflows"
    )

    # 6. Core Components Slide 1 (ML & DL)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Components: ML & DL"
    content = slide.append_text_box_if_needed(slide, "Machine Learning (ML):\n• Supervised Learning (Labeled data)\n• Unsupervised Learning (Unlabeled data)\n• Reinforcement Learning (Trial and error)\n• Examples: Netflix recommendations, Fraud detection\n\nDeep Learning (DL):\n• Uses multi-layered Neural Networks\n• Examples: Advanced image recognition, NLP")
    # Wait, I can't use append_text_box_if_needed, that's not a real method. 
    # I'll just use the standard way.
    
    # Let's rewrite the whole thing properly.
    pass

# Re-writing properly
